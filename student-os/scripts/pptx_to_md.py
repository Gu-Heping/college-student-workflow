#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
from pathlib import Path

def load_presentation():
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise SystemExit(
            "Missing dependency 'python-pptx'. Install the packages from requirements.txt before running this script."
        ) from exc
    return Presentation


def main() -> int:
    parser = argparse.ArgumentParser(description="Convert a PPTX deck into a markdown slide summary.")
    parser.add_argument("pptx", help="Path to the PPTX file")
    parser.add_argument("--output", required=True, help="Output markdown path")
    args = parser.parse_args()

    pptx_path = Path(args.pptx).resolve()
    output_path = Path(args.output).resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)

    Presentation = load_presentation()
    prs = Presentation(str(pptx_path))
    lines = [
        "---",
        "type: slide-summary",
        "course:",
        "status: active",
        "created:",
        "updated:",
        "tags: [import, slides]",
        f"source_file: {pptx_path}",
        "import_method: pptx-to-md",
        "repair_status:",
        "derived_from_import:",
        "---",
        "",
        f"# Slide Summary - {pptx_path.stem}",
        "",
        "## Source",
        "",
        f"- Source file: {pptx_path}",
        "- Import method: pptx-to-md",
        "",
        "## Slides",
        "",
    ]

    for index, slide in enumerate(prs.slides, start=1):
        title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else f"Slide {index}"
        lines.extend([f"## Slide {index}: {title}", ""])
        extracted = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    extracted.append(text)
        if extracted:
            for block in extracted:
                lines.extend([block, ""])
        else:
            lines.extend(["[No extractable text found on this slide.]", ""])

    output_path.write_text("\n".join(lines), encoding="utf-8")
    print(json.dumps({"output": str(output_path)}, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
