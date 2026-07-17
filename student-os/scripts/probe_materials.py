#!/usr/bin/env python3
from __future__ import annotations

import os
import re
import shutil
from pathlib import Path
from typing import Any


PDF_SUFFIXES = {".pdf"}
DOCX_SUFFIXES = {".docx"}
PPTX_SUFFIXES = {".pptx"}
XLSX_SUFFIXES = {".xlsx"}
IMAGE_SUFFIXES = {".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp"}
LEGACY_OFFICE_SUFFIXES = {".doc", ".ppt", ".xls"}
BINARY_INDEX_SUFFIXES = {".bit", ".ms14", ".vhd"}
TEXT_SKIP_SUFFIXES = {
    ".c",
    ".cc",
    ".cpp",
    ".csv",
    ".json",
    ".md",
    ".py",
    ".rs",
    ".sh",
    ".tex",
    ".toml",
    ".ts",
    ".tsx",
    ".txt",
    ".yaml",
    ".yml",
}
API_SUPPORTED_SUFFIXES = (
    PDF_SUFFIXES
    | DOCX_SUFFIXES
    | PPTX_SUFFIXES
    | XLSX_SUFFIXES
    | IMAGE_SUFFIXES
    | LEGACY_OFFICE_SUFFIXES
)

FORCE_STRATEGIES = (
    "ocr",
    "mineru-api",
    "pymupdf",
    "pandoc",
    "python-pptx",
    "xlsx-to-md",
    "image-index",
    "legacy-office-index",
    "docx-to-md",
    "pdf-to-md",
    "binary-index",
)

PDF_SAMPLE_PAGES = 5
PDF_SCANNED_CHARS_PER_PAGE = 40
DOCX_IMAGE_HEAVY_TEXT = 50
PPTX_CHARS_PER_SLIDE = 40


def _env_int(name: str, default: int) -> int:
    raw = os.environ.get(name)
    if not raw:
        return default
    try:
        value = int(raw)
    except ValueError:
        return default
    return value if value > 0 else default


def _env_float(name: str, default: float) -> float:
    raw = os.environ.get(name)
    if not raw:
        return default
    try:
        value = float(raw)
    except ValueError:
        return default
    return value if value >= 0 else default


# Page count is only a supporting signal for manuals; density gates the pymupdf path.
PDF_MANUAL_MIN_PAGES = _env_int("STUDENT_OS_PDF_MANUAL_MIN_PAGES", 80)
PDF_MANUAL_CHARS_PER_PAGE = _env_int("STUDENT_OS_PDF_MANUAL_CHARS_PER_PAGE", 800)
PDF_FORMULA_RATIO = _env_float("STUDENT_OS_PDF_FORMULA_RATIO", 0.045)

FORMULA_CHAR_RE = re.compile(
    r"[\u2200-\u22FF\u27C0-\u27EF\u2980-\u29FF\u2A00-\u2AFF∫∑∏√∞≈≠≤≥±∂∇α-ωΑ-Ω]"
)


def pandoc_available() -> bool:
    return shutil.which("pandoc") is not None


def _base_result(
    *,
    strategy: str,
    tool: str,
    reason: str,
    needs_ocr: bool = False,
    needs_api: bool = False,
    **metrics: Any,
) -> dict[str, Any]:
    payload: dict[str, Any] = {
        "strategy": strategy,
        "tool": tool,
        "reason": reason,
        "needs_ocr": needs_ocr,
        "needs_api": needs_api,
    }
    payload.update(metrics)
    return payload


def formula_density(text: str) -> float:
    if not text:
        return 0.0
    hits = len(FORMULA_CHAR_RE.findall(text))
    return hits / max(len(text), 1)


def local_tool_for_suffix(suffix: str) -> str | None:
    if suffix in PDF_SUFFIXES:
        return "pdf-to-md"
    if suffix in DOCX_SUFFIXES:
        return "pandoc" if pandoc_available() else "docx-to-md"
    if suffix in PPTX_SUFFIXES:
        return "python-pptx"
    if suffix in XLSX_SUFFIXES:
        return "xlsx-to-md"
    if suffix in IMAGE_SUFFIXES:
        return "image-index"
    if suffix in LEGACY_OFFICE_SUFFIXES:
        return "legacy-office-index"
    if suffix in BINARY_INDEX_SUFFIXES:
        return "binary-index"
    if suffix in TEXT_SKIP_SUFFIXES:
        return "skip"
    return "binary-index"


def probe_pdf(path: Path) -> dict[str, Any]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise SystemExit(
            "Missing dependency 'pypdf'. Install the packages from requirements.txt before probing PDFs."
        ) from exc

    try:
        reader = PdfReader(str(path))
        page_count = len(reader.pages)
        sample_n = min(PDF_SAMPLE_PAGES, page_count) if page_count else 0
        texts: list[str] = []
        for index in range(sample_n):
            try:
                texts.append(reader.pages[index].extract_text() or "")
            except Exception:
                texts.append("")
    except Exception as exc:
        raise RuntimeError(f"Failed to read PDF for page-count probe: {exc}") from exc

    joined = "\n".join(texts)
    text_len = len(joined.strip())
    chars_per_page = text_len / sample_n if sample_n else 0.0
    has_text_layer = chars_per_page >= PDF_SCANNED_CHARS_PER_PAGE
    density = formula_density(joined)

    metrics = {
        "page_count": page_count,
        "text_len": text_len,
        "chars_per_page": round(chars_per_page, 1),
        "has_text_layer": has_text_layer,
        "formula_density": round(density, 4),
        "file_size_bytes": path.stat().st_size,
    }

    if page_count == 0 or not has_text_layer:
        return _base_result(
            strategy="scanned",
            tool="mineru-api",
            reason=f"no/low text layer ({chars_per_page:.0f} chars/page over {sample_n} sampled pages)",
            needs_ocr=True,
            needs_api=True,
            **metrics,
        )

    # Density gates the local PyMuPDF path; page count is only mentioned as supporting context.
    if chars_per_page >= PDF_MANUAL_CHARS_PER_PAGE and density < PDF_FORMULA_RATIO:
        page_note = (
            f", {page_count} pages (>= {PDF_MANUAL_MIN_PAGES} suggests a manual)"
            if page_count >= PDF_MANUAL_MIN_PAGES
            else f", {page_count} pages"
        )
        return _base_result(
            strategy="text-manual",
            tool="pymupdf",
            reason=(
                f"text-heavy PDF ({chars_per_page:.0f} chars/page, formula density {density:.3f}"
                f"{page_note})"
            ),
            needs_ocr=False,
            needs_api=False,
            **metrics,
        )

    return _base_result(
        strategy="academic",
        tool="mineru-api",
        reason=f"mixed academic PDF ({page_count} pages, formula density {density:.3f})",
        needs_ocr=False,
        needs_api=True,
        **metrics,
    )


def _docx_text_len(document: Any) -> int:
    text_len = sum(len(paragraph.text.strip()) for paragraph in document.paragraphs)
    for table in getattr(document, "tables", []) or []:
        for row in table.rows:
            for cell in row.cells:
                text_len += len(cell.text.strip())
    return text_len


def _docx_image_count(document: Any) -> int:
    try:
        rels = getattr(getattr(document, "part", None), "rels", {}) or {}
        count = 0
        for rel in rels.values():
            reltype = str(getattr(rel, "reltype", "") or "")
            if "image" in reltype.lower():
                count += 1
        if count:
            return count
    except Exception:
        pass
    return len(getattr(document, "inline_shapes", []) or [])


def probe_docx(path: Path) -> dict[str, Any]:
    try:
        from docx import Document
    except ImportError as exc:
        raise SystemExit(
            "Missing dependency 'python-docx'. Install the packages from requirements.txt before probing DOCX."
        ) from exc

    try:
        document = Document(str(path))
        text_len = _docx_text_len(document)
        image_count = _docx_image_count(document)
    except Exception as exc:
        raise RuntimeError(f"Failed to probe DOCX {path}: {exc}") from exc

    metrics = {"text_len": text_len, "image_count": image_count}

    if text_len < DOCX_IMAGE_HEAVY_TEXT and image_count > 0:
        return _base_result(
            strategy="image-heavy-docx",
            tool="mineru-api",
            reason=f"image-heavy DOCX ({image_count} images, {text_len} chars)",
            needs_ocr=True,
            needs_api=True,
            **metrics,
        )

    tool = "pandoc" if pandoc_available() else "docx-to-md"
    return _base_result(
        strategy="text-docx",
        tool=tool,
        reason=f"text DOCX ({text_len} chars); using {tool}",
        needs_ocr=False,
        needs_api=False,
        **metrics,
    )


def _shape_is_picture(shape: Any) -> bool:
    shape_type = getattr(shape, "shape_type", None)
    if shape_type is None:
        return False
    name = str(getattr(shape_type, "name", shape_type)).upper()
    return "PICTURE" in name or str(shape_type) in {"13", "MSO_SHAPE_TYPE.PICTURE"}


def probe_pptx(path: Path) -> dict[str, Any]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise SystemExit(
            "Missing dependency 'python-pptx'. Install the packages from requirements.txt before probing PPTX."
        ) from exc

    try:
        presentation = Presentation(str(path))
        slide_count = len(presentation.slides)
        per_slide_chars: list[int] = []
        picture_count = 0
        for slide in presentation.slides:
            slide_text = 0
            for shape in slide.shapes:
                text = getattr(shape, "text", None)
                if text:
                    slide_text += len(str(text).strip())
                if _shape_is_picture(shape):
                    picture_count += 1
            per_slide_chars.append(slide_text)
    except Exception as exc:
        raise RuntimeError(f"Failed to probe PPTX {path}: {exc}") from exc

    text_len = sum(per_slide_chars)
    chars_per_slide = text_len / slide_count if slide_count else 0.0
    metrics = {
        "text_len": text_len,
        "slide_count": slide_count,
        "chars_per_slide": round(chars_per_slide, 1),
        "picture_count": picture_count,
    }

    # Short text-only decks stay local. Image-heavy / sparse-text decks go to MinerU.
    if picture_count == 0:
        return _base_result(
            strategy="text-pptx",
            tool="python-pptx",
            reason=f"text-only PPTX ({text_len} chars across {slide_count} slides)",
            needs_ocr=False,
            needs_api=False,
            **metrics,
        )

    if chars_per_slide >= PPTX_CHARS_PER_SLIDE and text_len >= max(picture_count * 40, PPTX_CHARS_PER_SLIDE):
        return _base_result(
            strategy="text-pptx",
            tool="python-pptx",
            reason=(
                f"text-forward PPTX ({chars_per_slide:.0f} chars/slide, {picture_count} pictures)"
            ),
            needs_ocr=False,
            needs_api=False,
            **metrics,
        )

    return _base_result(
        strategy="image-heavy-pptx",
        tool="mineru-api",
        reason=(
            f"image-heavy/sparse-text PPTX ({chars_per_slide:.0f} chars/slide, {picture_count} pictures)"
        ),
        needs_ocr=True,
        needs_api=True,
        **metrics,
    )


def probe_file(path: Path) -> dict[str, Any]:
    suffix = path.suffix.lower()
    result: dict[str, Any]
    if suffix in PDF_SUFFIXES:
        result = probe_pdf(path)
    elif suffix in DOCX_SUFFIXES:
        result = probe_docx(path)
    elif suffix in PPTX_SUFFIXES:
        result = probe_pptx(path)
    elif suffix in XLSX_SUFFIXES:
        result = _base_result(
            strategy="spreadsheet",
            tool="xlsx-to-md",
            reason="tabular spreadsheet",
            needs_ocr=False,
            needs_api=False,
        )
    elif suffix in IMAGE_SUFFIXES:
        result = _base_result(
            strategy="image",
            tool="mineru-api",
            reason="image file without a text layer",
            needs_ocr=True,
            needs_api=True,
        )
    elif suffix in LEGACY_OFFICE_SUFFIXES:
        result = _base_result(
            strategy="legacy-office",
            tool="mineru-api",
            reason="legacy Office binary; local converters unsupported",
            needs_ocr=False,
            needs_api=True,
        )
    elif suffix in BINARY_INDEX_SUFFIXES:
        result = _base_result(
            strategy="binary",
            tool="binary-index",
            reason="tool-specific binary; keep searchable index sidecar",
            needs_ocr=False,
            needs_api=False,
        )
    elif suffix in TEXT_SKIP_SUFFIXES:
        result = _base_result(
            strategy="already-text",
            tool="skip",
            reason="already text-friendly; skip conversion",
            needs_ocr=False,
            needs_api=False,
        )
    else:
        result = _base_result(
            strategy="binary",
            tool="binary-index",
            reason="unknown binary; keep searchable index sidecar",
            needs_ocr=False,
            needs_api=False,
        )

    result["source"] = str(path)
    result["suffix"] = suffix
    return result


def force_strategy_result(strategy: str, path: Path) -> dict[str, Any]:
    mapping = {
        "ocr": ("ocr", "mineru-api", True, True, "forced OCR via MinerU API"),
        "mineru-api": ("forced-api", "mineru-api", False, True, "forced MinerU API"),
        "pymupdf": ("forced-pymupdf", "pymupdf", False, False, "forced PyMuPDF local PDF extract"),
        "pandoc": ("forced-pandoc", "pandoc", False, False, "forced pandoc DOCX convert"),
        "docx-to-md": ("forced-docx", "docx-to-md", False, False, "forced local docx_to_md"),
        "python-pptx": ("forced-pptx", "python-pptx", False, False, "forced local pptx_to_md"),
        "xlsx-to-md": ("forced-xlsx", "xlsx-to-md", False, False, "forced local xlsx_to_md"),
        "pdf-to-md": ("forced-pdf", "pdf-to-md", False, False, "forced local pdf_to_markdown"),
        "image-index": ("forced-image-index", "image-index", False, False, "forced image index sidecar"),
        "legacy-office-index": (
            "forced-legacy-index",
            "legacy-office-index",
            False,
            False,
            "forced legacy Office index sidecar",
        ),
        "binary-index": ("forced-binary-index", "binary-index", False, False, "forced binary index sidecar"),
    }
    if strategy not in mapping:
        raise ValueError(f"Unsupported force strategy: {strategy}")
    name, tool, needs_ocr, needs_api, reason = mapping[strategy]
    result = _base_result(
        strategy=name,
        tool=tool,
        reason=reason,
        needs_ocr=needs_ocr,
        needs_api=needs_api,
    )
    result["source"] = str(path)
    result["suffix"] = path.suffix.lower()
    result["forced"] = True
    return result


def _optional_ocr_hint(path: Path) -> bool:
    suffix = path.suffix.lower()
    if suffix in IMAGE_SUFFIXES:
        return True
    if suffix not in PDF_SUFFIXES:
        return False
    try:
        return bool(probe_pdf(path).get("needs_ocr"))
    except Exception:
        return False


def resolve_probe(
    path: Path,
    *,
    method: str,
    force_strategy: str | None,
    has_api_token: bool,
) -> dict[str, Any]:
    suffix = path.suffix.lower()
    if force_strategy:
        probe = force_strategy_result(force_strategy, path)
    elif method == "auto":
        probe = probe_file(path)
    elif method == "api":
        # Do not let content probing block an explicit API request.
        if suffix in API_SUPPORTED_SUFFIXES:
            probe = _base_result(
                strategy="forced-method-api",
                tool="mineru-api",
                reason="--method api forces MinerU when the suffix is API-supported",
                needs_ocr=_optional_ocr_hint(path),
                needs_api=True,
            )
        else:
            tool = local_tool_for_suffix(suffix) or "binary-index"
            probe = _base_result(
                strategy="forced-method-api-unsupported",
                tool=tool,
                reason=f"--method api unsupported suffix; keep {tool}",
                needs_ocr=False,
                needs_api=False,
            )
        probe["source"] = str(path)
        probe["suffix"] = suffix
        # Explicit --method api must not silently degrade when a token is missing.
        probe["forced"] = True
    else:  # local
        tool = local_tool_for_suffix(suffix) or "binary-index"
        probe = _base_result(
            strategy="forced-method-local",
            tool=tool,
            reason=f"--method local forces {tool}",
            needs_ocr=False,
            needs_api=False,
        )
        probe["source"] = str(path)
        probe["suffix"] = suffix

    # Without a token, API-needed strategies degrade unless the user forced an API strategy.
    if probe.get("needs_api") and not has_api_token and not probe.get("forced"):
        if suffix in PDF_SUFFIXES:
            fallback_tool = "pdf-to-md"
            fallback_strategy = "api-unavailable-pdf"
        elif suffix in DOCX_SUFFIXES:
            fallback_tool = "pandoc" if pandoc_available() else "docx-to-md"
            fallback_strategy = "api-unavailable-docx"
        elif suffix in PPTX_SUFFIXES:
            fallback_tool = "python-pptx"
            fallback_strategy = "api-unavailable-pptx"
        elif suffix in IMAGE_SUFFIXES:
            fallback_tool = "image-index"
            fallback_strategy = "api-unavailable-image"
        elif suffix in LEGACY_OFFICE_SUFFIXES:
            fallback_tool = "legacy-office-index"
            fallback_strategy = "api-unavailable-legacy"
        else:
            fallback_tool = "binary-index"
            fallback_strategy = "api-unavailable-binary"
        probe = {
            **probe,
            "strategy": fallback_strategy,
            "tool": fallback_tool,
            "needs_api": False,
            "needs_ocr": False,
            "degraded": True,
            "reason": (
                f"{probe.get('reason')}; MinerU token unavailable, degraded to {fallback_tool}"
            ),
        }
    return probe
