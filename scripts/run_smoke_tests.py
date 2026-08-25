#!/usr/bin/env python3
from __future__ import annotations

import argparse
import http.server
import importlib.util
import json
import os
import re
import shutil
import socket
import subprocess
import sys
import tempfile
import threading
import types
from datetime import date, timedelta
from pathlib import Path, PurePosixPath, PureWindowsPath
from urllib.parse import urlparse


ROOT = Path(__file__).resolve().parents[1]
ROOT_SCRIPTS = ROOT / "scripts"
STUDENT_OS_SCRIPTS = ROOT / "student-os" / "scripts"
EXAMPLES_ROOT = ROOT / "examples"


def load_group_git_changes_module():
    script_path = STUDENT_OS_SCRIPTS / "group_git_changes.py"
    spec = importlib.util.spec_from_file_location("student_os_group_git_changes", script_path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"Unable to load module spec for {script_path}")
    module = importlib.util.module_from_spec(spec)
    original_flag = sys.dont_write_bytecode
    sys.dont_write_bytecode = True
    try:
        spec.loader.exec_module(module)
    finally:
        sys.dont_write_bytecode = original_flag
    return module


def run_script(name: str, *args: str, cwd: Path = ROOT, env: dict[str, str] | None = None) -> str:
    script_path = STUDENT_OS_SCRIPTS / name
    result = subprocess.run(
        [sys.executable, "-B", str(script_path), *args],
        check=True,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=cwd,
        env={
            **os.environ,
            "PYTHONDONTWRITEBYTECODE": "1",
            "PYTHONIOENCODING": "utf-8",
            **(env or {}),
        },
    )
    return result.stdout.strip()


def run_script_with_stdin(
    name: str,
    stdin_text: str,
    *args: str,
    cwd: Path = ROOT,
    check: bool = True,
) -> subprocess.CompletedProcess[str]:
    script_path = STUDENT_OS_SCRIPTS / name
    return subprocess.run(
        [sys.executable, "-B", str(script_path), *args],
        check=check,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=cwd,
        input=stdin_text,
        env={
            **os.environ,
            "PYTHONDONTWRITEBYTECODE": "1",
            "PYTHONIOENCODING": "utf-8",
        },
    )


def run_root_script(name: str, *args: str, cwd: Path = ROOT, env: dict[str, str] | None = None) -> str:
    script_path = ROOT_SCRIPTS / name
    result = subprocess.run(
        [sys.executable, "-B", str(script_path), *args],
        check=True,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=cwd,
        env={
            **os.environ,
            "PYTHONDONTWRITEBYTECODE": "1",
            "PYTHONIOENCODING": "utf-8",
            **(env or {}),
        },
    )
    return result.stdout.strip()


def run_path_script(script_path: Path, *args: str, cwd: Path, env: dict[str, str] | None = None) -> str:
    result = subprocess.run(
        [sys.executable, "-B", str(script_path), *args],
        check=True,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=cwd,
        env={
            **os.environ,
            "PYTHONDONTWRITEBYTECODE": "1",
            "PYTHONIOENCODING": "utf-8",
            **(env or {}),
        },
    )
    return result.stdout.strip()


def run_script_failure(name: str, *args: str, cwd: Path = ROOT) -> str:
    script_path = STUDENT_OS_SCRIPTS / name
    result = subprocess.run(
        [sys.executable, "-B", str(script_path), *args],
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=cwd,
        env={**os.environ, "PYTHONDONTWRITEBYTECODE": "1", "PYTHONIOENCODING": "utf-8"},
    )
    if result.returncode == 0:
        raise AssertionError(f"Expected {name} to fail for args {args!r}")
    return (result.stderr or result.stdout).strip()


def run_root_script_failure(name: str, *args: str, cwd: Path = ROOT, env: dict[str, str] | None = None) -> str:
    script_path = ROOT_SCRIPTS / name
    result = subprocess.run(
        [sys.executable, "-B", str(script_path), *args],
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=cwd,
        env={
            **os.environ,
            "PYTHONDONTWRITEBYTECODE": "1",
            "PYTHONIOENCODING": "utf-8",
            **(env or {}),
        },
    )
    if result.returncode == 0:
        raise AssertionError(f"Expected {name} to fail for args {args!r}")
    return (result.stderr or result.stdout).strip()


def run_path_script_failure(script_path: Path, *args: str, cwd: Path, env: dict[str, str] | None = None) -> str:
    result = subprocess.run(
        [sys.executable, "-B", str(script_path), *args],
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=cwd,
        env={
            **os.environ,
            "PYTHONDONTWRITEBYTECODE": "1",
            "PYTHONIOENCODING": "utf-8",
            **(env or {}),
        },
    )
    if result.returncode == 0:
        raise AssertionError(f"Expected {script_path.name} to fail for args {args!r}")
    return (result.stderr or result.stdout).strip()


def ensure_contains(path: Path, needle: str) -> None:
    text = path.read_text(encoding="utf-8")
    if needle not in text:
        raise AssertionError(f"{path} does not contain expected text: {needle}")


def ensure_not_contains(path: Path, needle: str) -> None:
    text = path.read_text(encoding="utf-8")
    if needle in text:
        raise AssertionError(f"{path} contains unexpected text: {needle}")


def ensure_exists(path: Path) -> None:
    if not path.exists():
        raise AssertionError(f"Expected path to exist: {path}")


def _run_cross_validate(repo: Path, course: str = "linear-algebra", exam_scope: str = "期中") -> tuple[int, dict]:
    result = subprocess.run(
        [
            sys.executable,
            "-B",
            str(STUDENT_OS_SCRIPTS / "cross_validate_exam_census.py"),
            str(repo),
            "--course",
            course,
            "--exam-scope",
            exam_scope,
        ],
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=ROOT,
        env={**os.environ, "PYTHONDONTWRITEBYTECODE": "1", "PYTHONIOENCODING": "utf-8"},
    )
    payload: dict = {}
    stdout = (result.stdout or "").strip()
    if stdout:
        try:
            payload = json.loads(stdout)
        except json.JSONDecodeError as exc:
            raise AssertionError(f"cross_validate stdout not JSON: {stdout[:500]}") from exc
    return result.returncode, payload


def write_minimal_prep_pack(reviews_root: Path, skeleton_names: list[str]) -> None:
    """Write L1/L3/L4 files that satisfy Phase E prep-pack mechanical checks."""
    reviews_root.mkdir(parents=True, exist_ok=True)
    type_rows = []
    for name in skeleton_names:
        type_id = name.split("-", 1)[-1].removesuffix(".md") if "-" in name else name
        type_rows.append(
            f"| P0 | {type_id} | high | 1h | [题型解析/{name}](题型解析/{name}) |"
        )
    if not type_rows:
        type_rows.append("| P0 | sample | — | — | [题型解析/](题型解析/) |")
    type_table = "\n".join(type_rows)
    first_skel = skeleton_names[0] if skeleton_names else ""
    first_link = f"[题型解析/{first_skel}](题型解析/{first_skel})" if first_skel else "[题型解析/](题型解析/)"

    (reviews_root / "备考指南.md").write_text(
        "\n".join(
            [
                "---",
                "type: exam-prep-guide",
                'course: "linear-algebra"',
                'exam_scope: "期中"',
                "status: active",
                "review_scope: exam-census",
                "---",
                "",
                "# linear-algebra · 期中 · 备考指南",
                "",
                "## 怎么使用这套资料",
                "",
                "| 层级 | 文件 | 什么时候用 | 目标 |",
                "| --- | --- | --- | --- |",
                "| L1 | [备考指南.md](备考指南.md) | 开始前 | 规划 |",
                "| L2 | [题型解析/](题型解析/) | 学题型 | 深入 |",
                "| L3 | [公式总卡.md](公式总卡.md) / [答题模板速查.md](答题模板速查.md) | 速查 | 背诵 |",
                "| L4 | [考前1小时清单.md](考前1小时清单.md) | 考前1h | 冲刺 |",
                "",
                "## 题型优先级",
                "",
                "| 优先级 | 题型 | 出现率 | 建议投入时间 | 入口 |",
                "| --- | --- | ---: | --- | --- |",
                type_table,
                "",
                "## 复习时间分配",
                "",
                "| 时间总量 | 先做什么 | 目标 |",
                "| --- | --- | --- |",
                "| 1 小时 | P0 + 公式总卡 | 定向 |",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    (reviews_root / "公式总卡.md").write_text(
        "\n".join(
            [
                "---",
                "type: formula-cheat-sheet",
                'course: "linear-algebra"',
                'exam_scope: "期中"',
                "status: active",
                "review_scope: exam-census",
                "---",
                "",
                "# linear-algebra · 期中 · 公式总卡",
                "",
                "## 高频公式速查",
                "",
                "| 题型 | 看到什么 | 公式 / 结论 | 先算什么 | 最容易错 | 来源 |",
                "| --- | --- | --- | --- | --- | --- |",
                f"| matrix-rank | 秩 | $\\lvert A\\rvert$ | 化阶梯 | 漏零行 | {first_link} |",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    (reviews_root / "答题模板速查.md").write_text(
        "\n".join(
            [
                "---",
                "type: answer-template-quickref",
                'course: "linear-algebra"',
                'exam_scope: "期中"',
                "status: active",
                "review_scope: exam-census",
                "---",
                "",
                "# linear-algebra · 期中 · 答题模板速查",
                "",
                "## 标准答题模板",
                "",
                "| 题型 | 看到什么 | 第一句写什么 | 填空式模板 | 来源 |",
                "| --- | --- | --- | --- | --- |",
                f"| matrix-rank | 求秩 | 先化阶梯 | 由[条件]，得[表达式]=[答案]。 | {first_link} |",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    (reviews_root / "考前1小时清单.md").write_text(
        "\n".join(
            [
                "---",
                "type: pre-exam-one-hour-checklist",
                'course: "linear-algebra"',
                'exam_scope: "期中"',
                "status: active",
                "review_scope: exam-census",
                "---",
                "",
                "# linear-algebra · 期中 · 考前1小时清单",
                "",
                "## 最后 60 分钟怎么用",
                "",
                "| 时间 | 做什么 | 文件 | 目标 |",
                "| --- | --- | --- | --- |",
                "| 60-45 分钟 | P0 | [备考指南.md](备考指南.md) / [题型解析/](题型解析/) | 方法 |",
                "| 45-30 分钟 | 公式 | [公式总卡.md](公式总卡.md) | 背诵 |",
                "| 30-15 分钟 | 模板 | [答题模板速查.md](答题模板速查.md) | 步骤分 |",
                f"| 15-5 分钟 | 易错 | {first_link} | 避坑 |",
                "| 5-0 分钟 | checklist | 本文件 | 稳住 |",
                "",
                "## 最后检查",
                "",
                "- [ ] 高频题型入口看过",
                "- [ ] 公式总卡看过",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )


def assert_prep_pack_templates() -> None:
    templates = ROOT / "student-os" / "templates"
    expected = {
        "exam-prep-guide.md": "exam-prep-guide",
        "formula-cheat-sheet.md": "formula-cheat-sheet",
        "answer-template-quickref.md": "answer-template-quickref",
        "pre-exam-one-hour-checklist.md": "pre-exam-one-hour-checklist",
    }
    for filename, type_name in expected.items():
        path = templates / filename
        ensure_exists(path)
        text = path.read_text(encoding="utf-8")
        if f"type: {type_name}" not in text and f'type: "{type_name}"' not in text:
            raise AssertionError(f"{filename} missing frontmatter type: {type_name}")
        if filename == "exam-prep-guide.md":
            for needle in ("怎么使用这套资料", "题型优先级", "复习时间分配", "L1", "L2", "L3", "L4"):
                if needle not in text:
                    raise AssertionError(f"{filename} missing {needle!r}")
        if filename == "pre-exam-one-hour-checklist.md":
            for slot in ("60-45", "45-30", "30-15", "15-5", "5-0"):
                if slot not in text:
                    raise AssertionError(f"{filename} missing time slot {slot}")


def copy_repo(src: Path, dest: Path) -> None:
    if dest.exists():
        shutil.rmtree(dest)
    shutil.copytree(src, dest)


def load_root_script_module(name: str, module_name: str) -> object:
    script_path = ROOT_SCRIPTS / name
    spec = importlib.util.spec_from_file_location(module_name, script_path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"Unable to load module spec for {script_path}")
    module = importlib.util.module_from_spec(spec)
    original_flag = sys.dont_write_bytecode
    sys.dont_write_bytecode = True
    previous_module = sys.modules.get(module_name)
    sys.modules[module_name] = module
    try:
        spec.loader.exec_module(module)
    finally:
        if previous_module is None:
            sys.modules.pop(module_name, None)
        else:
            sys.modules[module_name] = previous_module
        sys.dont_write_bytecode = original_flag
    return module


def load_student_os_script_module(name: str, module_name: str) -> object:
    script_path = STUDENT_OS_SCRIPTS / name
    spec = importlib.util.spec_from_file_location(module_name, script_path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"Unable to load module spec for {script_path}")
    module = importlib.util.module_from_spec(spec)
    original_flag = sys.dont_write_bytecode
    sys.dont_write_bytecode = True
    previous_module = sys.modules.get(module_name)
    sys.modules[module_name] = module
    try:
        spec.loader.exec_module(module)
    finally:
        if previous_module is None:
            sys.modules.pop(module_name, None)
        else:
            sys.modules[module_name] = previous_module
        sys.dont_write_bytecode = original_flag
    return module


def load_import_dependencies() -> tuple[object, object, object, object]:
    try:
        from docx import Document
        from openpyxl import Workbook
        from pypdf import PdfWriter
        from pptx import Presentation
    except ImportError as exc:
        raise SystemExit(
            "Missing one or more import-workflow dependencies. Install requirements.txt before running smoke tests."
        ) from exc
    return Document, Workbook, Presentation, PdfWriter


def normalize_text_files(root: Path) -> None:
    for path in root.rglob("*"):
        if not path.is_file():
            continue
        if path.name == ".gitkeep":
            continue
        try:
            text = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            continue
        normalized = text.replace("\r\n", "\n")
        path.write_text(normalized, encoding="utf-8", newline="\n")


def scrub_example_paths(root: Path, source_root: Path) -> None:
    abs_root = str(source_root.resolve())
    escaped_root = abs_root.replace("\\", "\\\\")
    for path in root.rglob("*"):
        if not path.is_file() or path.name == ".gitkeep":
            continue
        try:
            text = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            continue
        scrubbed = text
        scrubbed = scrubbed.replace(f"{escaped_root}\\\\", "")
        scrubbed = scrubbed.replace(f"{escaped_root}/", "")
        scrubbed = scrubbed.replace(f"{abs_root}\\", "")
        scrubbed = scrubbed.replace(f"{abs_root}/", "")
        if scrubbed != text:
            path.write_text(scrubbed, encoding="utf-8", newline="\n")


def materialize_empty_dirs(root: Path) -> None:
    for path in sorted([candidate for candidate in root.rglob("*") if candidate.is_dir()], key=lambda item: len(item.parts), reverse=True):
        if any(path.iterdir()):
            continue
        (path / ".gitkeep").write_text("", encoding="utf-8")


def rewrite_legacy_task_link(task_path: Path) -> None:
    text = task_path.read_text(encoding="utf-8")
    old = "- Course: ../../courses/legacy-course/index.md"
    new = "- Course: legacy course folder without generated course home"
    task_path.write_text(text.replace(old, new), encoding="utf-8", newline="\n")


def write_task_fixture(
    path: Path,
    *,
    title: str,
    due: str = "",
    area: str = "",
    priority: str = "",
    course: str = "",
    tags: str = "[task]",
    course_link: str = "",
) -> None:
    lines = [
        "---",
        "type: task",
        f"course: {course}",
        "status: active",
        f"created: {date.today().isoformat()}",
        f"updated: {date.today().isoformat()}",
        f"tags: {tags}",
        "---",
        "",
        f"# {title}",
        "",
        "## Details",
        "",
        f"- Due: {due}",
        f"- Area: {area}",
        f"- Priority: {priority}",
        "",
        "## Checklist",
        "",
        "- [ ] Clarify scope",
        "- [ ] Start work",
        "- [ ] Finish",
        "",
        "## Links",
        "",
        f"- Course: {course_link}",
        "- Project:",
        "- Source file:",
    ]
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text("\n".join(lines) + "\n", encoding="utf-8", newline="\n")


def seed_planning_inputs(repo: Path, today: date) -> None:
    write_task_fixture(
        repo / "tasks" / "inbox" / "capture-linear-algebra-questions.md",
        title="Capture linear algebra questions",
        area="inbox",
        priority="medium",
        course="Linear Algebra",
        tags="[task, inbox]",
    )
    write_task_fixture(
        repo / "tasks" / "deadlines" / "linear-algebra-midterm-checkpoint.md",
        title="Linear Algebra Midterm",
        due=(today + timedelta(days=10)).isoformat(),
        area="exam",
        priority="high",
        course="Linear Algebra",
        tags="[task, exam]",
    )
    write_task_fixture(
        repo / "tasks" / "deadlines" / "linear-algebra-overdue-reading.md",
        title="Linear Algebra Overdue Reading",
        due=(today - timedelta(days=2)).isoformat(),
        area="reading",
        priority="medium",
        course="Linear Algebra",
    )
    archived_path = repo / "tasks" / "deadlines" / "linear-algebra-archived-quiz.md"
    write_task_fixture(
        archived_path,
        title="Linear Algebra Archived Quiz",
        due=(today - timedelta(days=1)).isoformat(),
        area="exam",
        priority="low",
        course="Linear Algebra",
        tags="[task, exam]",
    )
    archived_text = archived_path.read_text(encoding="utf-8").replace("status: active", "status: archived", 1)
    archived_path.write_text(archived_text, encoding="utf-8", newline="\n")
    review_path = repo / "courses" / "linear-algebra" / "reviews" / "chapter-1-review.md"
    review_path.parent.mkdir(parents=True, exist_ok=True)
    review_path.write_text(
        "\n".join(
            [
                "---",
                "type: chapter-review",
                "course: Linear Algebra",
                "status: active",
                f"created: {today.isoformat()}",
                f"updated: {today.isoformat()}",
                "tags: [review]",
                "---",
                "",
                "# Chapter 1 Review",
                "",
                "## Concepts",
                "",
                "- Basis changes",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    dashboard_path = repo / "courses" / "linear-algebra" / "dashboard.md"
    dashboard_text = dashboard_path.read_text(encoding="utf-8")
    dashboard_text = dashboard_text.replace("- Next exam:", f"- Next exam: {(today + timedelta(days=9)).isoformat()}", 1)
    dashboard_path.write_text(dashboard_text, encoding="utf-8", newline="\n")


def write_docx_fixture(path: Path) -> None:
    Document, _, _, _ = load_import_dependencies()
    document = Document()
    document.add_heading("Linear Algebra Import Outline", level=1)
    document.add_paragraph("Focus on eigenvalues, diagonalization, and orthogonality.")
    table = document.add_table(rows=3, cols=2)
    table.rows[0].cells[0].text = "Week"
    table.rows[0].cells[1].text = "Topic"
    table.rows[1].cells[0].text = "3"
    table.rows[1].cells[1].text = "Eigenvectors"
    table.rows[2].cells[0].text = "4"
    table.rows[2].cells[1].text = "Orthogonal bases"
    document.save(str(path))


def write_xlsx_fixture(path: Path) -> None:
    _, Workbook, _, _ = load_import_dependencies()
    wb = Workbook()
    ws = wb.active
    ws.title = "Progress"
    ws.append(["Task", "Score", "Weight", "Weighted"])
    ws.append(["Quiz 1", 92, 0.2, "=B2*C2"])
    ws.append(["Homework 1", 88, 0.3, "=B3*C3"])
    ws.append(["Midterm Prep", None, 0.5, "=B4*C4"])
    second = wb.create_sheet("Deadlines")
    second.append(["Item", "Date"])
    second.append(["Worksheet A", "2026-07-11"])
    second.append(["Review Session", "2026-07-12"])
    wb.save(str(path))


def write_pptx_fixture(path: Path) -> None:
    _, _, Presentation, _ = load_import_dependencies()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Linear Algebra Week 2"
    slide.placeholders[1].text = "Diagonalization and basis changes"
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Key reminders"
    slide2.placeholders[1].text = "Check eigenvalue multiplicities\nLink notes to review sheets"
    prs.save(str(path))


def write_pdf_fixture(path: Path) -> None:
    _, _, _, PdfWriter = load_import_dependencies()
    from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

    writer = PdfWriter()
    page = writer.add_blank_page(width=612, height=792)
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    page[NameObject("/Resources")] = DictionaryObject(
        {NameObject("/Font"): DictionaryObject({NameObject("/F1"): font})}
    )
    stream = DecodedStreamObject()
    stream.set_data(
        b"BT\n/F1 18 Tf\n72 720 Td\n(#Broken linear algebra handout) Tj\n0 -24 Td\n(Page 1) Tj\n0 -24 Td\n(##Next section ........ 4) Tj\n0 -24 Td\n(-  Orthogonal projection summary) Tj\nET"
    )
    page[NameObject("/Contents")] = writer._add_object(stream)
    writer.add_metadata({"/Title": "Linear Algebra Import Handout"})
    with path.open("wb") as handle:
        writer.write(handle)


def write_multipage_pdf_fixture(path: Path, page_count: int) -> None:
    _, _, _, PdfWriter = load_import_dependencies()
    writer = PdfWriter()
    for _ in range(page_count):
        writer.add_blank_page(width=612, height=792)
    with path.open("wb") as handle:
        writer.write(handle)


def write_text_heavy_pdf_fixture(path: Path, page_count: int) -> None:
    _, _, _, PdfWriter = load_import_dependencies()
    from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

    writer = PdfWriter()
    filler = "Linear algebra handbook text sample. " * 40
    for page_number in range(1, page_count + 1):
        page = writer.add_blank_page(width=612, height=792)
        font = DictionaryObject(
            {
                NameObject("/Type"): NameObject("/Font"),
                NameObject("/Subtype"): NameObject("/Type1"),
                NameObject("/BaseFont"): NameObject("/Helvetica"),
            }
        )
        page[NameObject("/Resources")] = DictionaryObject(
            {NameObject("/Font"): DictionaryObject({NameObject("/F1"): font})}
        )
        stream = DecodedStreamObject()
        # Keep content ASCII-safe for the PDF content stream.
        payload = f"BT\n/F1 10 Tf\n72 720 Td\n(Page {page_number} {filler}) Tj\nET".encode("latin-1", errors="ignore")
        stream.set_data(payload)
        page[NameObject("/Contents")] = writer._add_object(stream)
    with path.open("wb") as handle:
        writer.write(handle)


def write_png_fixture(path: Path) -> None:
    path.write_bytes(
        bytes.fromhex(
            "89504E470D0A1A0A0000000D4948445200000001000000010802000000907753DE"
            "0000000C49444154789C6360600000000400010D0A2DB40000000049454E44AE426082"
        )
    )


def write_fake_mineru_sdk(root: Path) -> Path:
    module_path = root / "mineru.py"
    module_path.write_text(
        "\n".join(
            [
                "from dataclasses import dataclass, field",
                "from pathlib import Path",
                "",
                "@dataclass",
                "class Image:",
                "    name: str",
                "    data: bytes",
                "",
                "@dataclass",
                "class ExtractResult:",
                "    markdown: str",
                "    images: list[Image] = field(default_factory=list)",
                "    error: str | None = None",
                "",
                "class MinerU:",
                "    def __init__(self, token=None, base_url='https://mineru.net/api/v4', flash_base_url=None):",
                "        self.token = token",
                "",
                "    def extract(self, source, *, model=None, ocr=None, formula=None, table=None, language=None, pages=None, extra_formats=None, file_params=None, timeout=300):",
                "        path = Path(source)",
                "        body = [",
                "            f'# API Parsed - {path.name}',",
                "            '',",
                "            f'- model: {model}',",
                "            f'- language: {language}',",
                "            f'- pages: {pages}',",
                "        ]",
                "        images = []",
                "        if '.part' in path.name:",
                "            body.append('')",
                "            body.append('![figure](images/image1.png)')",
                "            images.append(Image(name='image1.png', data=b'\\x89PNG'))",
                "        return ExtractResult(markdown='\\n'.join(body), images=images)",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    return root


class FakeMineruAgentServer:
    def __init__(self) -> None:
        self.tasks: dict[str, dict[str, str | bytes]] = {}
        self.counter = 0
        self.server: http.server.ThreadingHTTPServer | None = None
        self.thread: threading.Thread | None = None

    @property
    def base_url(self) -> str:
        if self.server is None:
            raise RuntimeError("Fake MinerU Agent server is not running")
        host, port = self.server.server_address
        return f"http://{host}:{port}/api/v1/agent"

    def __enter__(self) -> "FakeMineruAgentServer":
        owner = self

        class Handler(http.server.BaseHTTPRequestHandler):
            def log_message(self, _format: str, *_args: object) -> None:
                return

            def send_json(self, payload: dict[str, object]) -> None:
                data = json.dumps(payload).encode("utf-8")
                self.send_response(200)
                self.send_header("Content-Type", "application/json")
                self.send_header("Content-Length", str(len(data)))
                self.end_headers()
                self.wfile.write(data)

            def do_POST(self) -> None:
                if self.path != "/api/v1/agent/parse/file":
                    self.send_error(404)
                    return
                length = int(self.headers.get("Content-Length", "0"))
                body = self.rfile.read(length).decode("utf-8") if length else "{}"
                payload = json.loads(body)
                owner.counter += 1
                task_id = f"task-{owner.counter}"
                filename = payload.get("file_name") or payload.get("filename") or ""
                owner.tasks[task_id] = {"filename": str(filename), "content": b"", "headers": ""}
                self.send_json(
                    {
                        "code": 0,
                        "data": {
                            "task_id": task_id,
                            "file_url": f"http://{self.server.server_address[0]}:{self.server.server_address[1]}/upload/{task_id}",
                        },
                    }
                )

            def do_PUT(self) -> None:
                parsed = urlparse(self.path)
                if not parsed.path.startswith("/upload/"):
                    self.send_error(404)
                    return
                task_id = parsed.path.rsplit("/", 1)[-1]
                length = int(self.headers.get("Content-Length", "0"))
                if task_id not in owner.tasks:
                    self.send_error(404)
                    return
                owner.tasks[task_id]["content"] = self.rfile.read(length)
                owner.tasks[task_id]["headers"] = json.dumps(dict(self.headers), sort_keys=True)
                self.send_response(200)
                self.end_headers()

            def do_GET(self) -> None:
                parsed = urlparse(self.path)
                if parsed.path.startswith("/api/v1/agent/parse/"):
                    task_id = parsed.path.rsplit("/", 1)[-1]
                    if task_id not in owner.tasks:
                        self.send_error(404)
                        return
                    self.send_json(
                        {
                            "code": 0,
                            "data": {
                                "task_id": task_id,
                                "state": "done",
                                "markdown_url": f"http://{self.server.server_address[0]}:{self.server.server_address[1]}/markdown/{task_id}",
                            },
                        }
                    )
                    return
                if parsed.path.startswith("/markdown/"):
                    task_id = parsed.path.rsplit("/", 1)[-1]
                    task = owner.tasks.get(task_id)
                    if task is None:
                        self.send_error(404)
                        return
                    filename = str(task.get("filename", "unknown"))
                    uploaded = len(task.get("content") or b"")
                    headers = str(task.get("headers", "{}"))
                    data = (
                        f"# Agent Parsed - {filename}\n\n"
                        f"- task: {task_id}\n"
                        f"- uploaded_bytes: {uploaded}\n"
                        f"- upload_headers: {headers}\n"
                    ).encode("utf-8")
                    self.send_response(200)
                    self.send_header("Content-Type", "text/markdown; charset=utf-8")
                    self.send_header("Content-Length", str(len(data)))
                    self.end_headers()
                    self.wfile.write(data)
                    return
                self.send_error(404)

        self.server = http.server.ThreadingHTTPServer(("127.0.0.1", 0), Handler)
        self.thread = threading.Thread(target=self.server.serve_forever, daemon=True)
        self.thread.start()
        return self

    def __exit__(self, exc_type: object, exc: object, traceback: object) -> None:
        if self.server is not None:
            self.server.shutdown()
            self.server.server_close()
        if self.thread is not None:
            self.thread.join(timeout=5)


def verify_material_type_constants() -> None:
    previous_sys_path = list(sys.path)
    sys.path.insert(0, str(STUDENT_OS_SCRIPTS))
    try:
        material_types = load_student_os_script_module("material_types.py", "student_os_material_types_smoke")
        materials_convert = load_student_os_script_module(
            "materials_convert.py", "student_os_materials_convert_types_smoke"
        )
        probe_materials = load_student_os_script_module("probe_materials.py", "student_os_probe_materials_types_smoke")
    finally:
        sys.path = previous_sys_path

    shared_names = [
        "PDF_SUFFIXES",
        "DOCX_SUFFIXES",
        "PPTX_SUFFIXES",
        "XLSX_SUFFIXES",
        "IMAGE_SUFFIXES",
        "LEGACY_OFFICE_SUFFIXES",
        "BINARY_INDEX_SUFFIXES",
        "TEXT_SKIP_SUFFIXES",
        "API_SUPPORTED_SUFFIXES",
    ]
    for name in shared_names:
        expected = getattr(material_types, name)
        if getattr(materials_convert, name) != expected:
            raise AssertionError(f"materials_convert.py drifted from material_types.{name}")
        if getattr(probe_materials, name) != expected:
            raise AssertionError(f"probe_materials.py drifted from material_types.{name}")


def verify_mineru_agent_helper_guards() -> None:
    previous_sys_path = list(sys.path)
    previous_high = os.environ.get("STUDENT_OS_MINERU_AGENT_MAX_FILE_BYTES")
    sys.path.insert(0, str(STUDENT_OS_SCRIPTS))
    try:
        os.environ["STUDENT_OS_MINERU_AGENT_MAX_FILE_BYTES"] = str(20 * 1024 * 1024)
        high_limit_types = load_student_os_script_module("material_types.py", "student_os_material_types_high_limit_smoke")
        if high_limit_types.MINERU_AGENT_MAX_FILE_BYTES != 10 * 1024 * 1024:
            raise AssertionError("MinerU Agent limit override must not exceed the 10 MiB hard limit")

        os.environ["STUDENT_OS_MINERU_AGENT_MAX_FILE_BYTES"] = "1234"
        low_limit_types = load_student_os_script_module("material_types.py", "student_os_material_types_low_limit_smoke")
        if low_limit_types.MINERU_AGENT_MAX_FILE_BYTES != 1234:
            raise AssertionError("MinerU Agent limit override should allow lower test limits")
    finally:
        if previous_high is None:
            os.environ.pop("STUDENT_OS_MINERU_AGENT_MAX_FILE_BYTES", None)
        else:
            os.environ["STUDENT_OS_MINERU_AGENT_MAX_FILE_BYTES"] = previous_high
        sys.path = previous_sys_path

    previous_sys_path = list(sys.path)
    sys.path.insert(0, str(STUDENT_OS_SCRIPTS))
    try:
        materials_convert = load_student_os_script_module(
            "materials_convert.py", "student_os_materials_convert_agent_guard_smoke"
        )
        probe_materials = load_student_os_script_module(
            "probe_materials.py", "student_os_probe_materials_agent_guard_smoke"
        )
    finally:
        sys.path = previous_sys_path

    if not materials_convert.mineru_agent_url_allowed("https://mineru.oss-cn-shanghai.aliyuncs.com/demo/result.md"):
        raise AssertionError("Expected MinerU OSS result host to be allowed by default")
    if not probe_materials.mojibake_metrics("\u0080" * 200)["mojibake_suspect"]:
        raise AssertionError("Expected repeated C1-control mojibake sample to be detected")

    # Latin-1 Supplement / punctuation mojibake (e.g. no-ToUnicode SimSun subset PDFs)
    # must be detected even when C1-control ratio alone is below the threshold.
    latin_mojibake_sample = "\u00C6\u00B5\u201AS\u00D8\u2021\u2030K" * 100
    latin_mojibake_metrics = probe_materials.mojibake_metrics(latin_mojibake_sample)
    if not latin_mojibake_metrics["mojibake_suspect"]:
        raise AssertionError(
            f"Expected Latin-1/punctuation mojibake sample to be detected, got: {latin_mojibake_metrics}"
        )

    # Normal Latin-script text with accents must not be flagged as mojibake.
    for normal_sample in (
        "tiếng Việt có dấu tiếng Việt có dấu " * 10,
        "Le français utilise des accents comme é è à ç " * 10,
        "Polski używa znaków takich jak ą ę ś ć ń " * 10,
    ):
        normal_metrics = probe_materials.mojibake_metrics(normal_sample)
        if normal_metrics["mojibake_suspect"]:
            raise AssertionError(
                f"Expected normal text not to be mojibake, got: {normal_metrics}"
            )

    class MojibakePage:
        def extract_text(self) -> str:
            return "\u0080" * 200

    class MojibakeReader:
        def __init__(self, _path: str) -> None:
            self.pages = [MojibakePage()]

    previous_pypdf = sys.modules.get("pypdf")
    sys.modules["pypdf"] = types.SimpleNamespace(PdfReader=MojibakeReader)
    try:
        mojibake_fd, mojibake_name = tempfile.mkstemp(suffix=".pdf")
        os.close(mojibake_fd)
        mojibake_pdf = Path(mojibake_name)
        mojibake_pdf.write_bytes(b"%PDF-1.4\n")
        try:
            mojibake_probe = probe_materials.probe_pdf(mojibake_pdf)
        finally:
            mojibake_pdf.unlink(missing_ok=True)
    finally:
        if previous_pypdf is None:
            sys.modules.pop("pypdf", None)
        else:
            sys.modules["pypdf"] = previous_pypdf
    if mojibake_probe["tool"] != "mineru-api" or mojibake_probe["strategy"] != "mojibake-text-layer":
        raise AssertionError(f"Expected mojibake PDF text layer to route to MinerU API, got: {mojibake_probe}")

    secret_fd, secret_name = tempfile.mkstemp()
    os.close(secret_fd)
    secret = Path(secret_name)
    secret.write_text("LOCAL_SECRET", encoding="utf-8")

    class RedirectToFile(http.server.BaseHTTPRequestHandler):
        def log_message(self, _format: str, *_args: object) -> None:
            return

        def do_GET(self) -> None:
            self.send_response(302)
            self.send_header("Location", secret.as_uri())
            self.send_header("Content-Length", "0")
            self.end_headers()

    redirect_server = http.server.ThreadingHTTPServer(("127.0.0.1", 0), RedirectToFile)
    redirect_thread = threading.Thread(target=redirect_server.serve_forever, daemon=True)
    redirect_thread.start()
    try:
        redirect_url = f"http://127.0.0.1:{redirect_server.server_port}/markdown"
        previous_base_url = os.environ.get("STUDENT_OS_MINERU_AGENT_BASE_URL")
        proxy_keys = ("HTTP_PROXY", "HTTPS_PROXY", "http_proxy", "https_proxy", "NO_PROXY", "no_proxy")
        previous_proxy_env = {key: os.environ.get(key) for key in proxy_keys}
        os.environ["STUDENT_OS_MINERU_AGENT_BASE_URL"] = f"http://127.0.0.1:{redirect_server.server_port}/api/v1/agent"
        os.environ["HTTP_PROXY"] = "http://127.0.0.1:9"
        os.environ["HTTPS_PROXY"] = "http://127.0.0.1:9"
        os.environ["NO_PROXY"] = ""
        os.environ["no_proxy"] = ""
        try:
            try:
                materials_convert.http_text(redirect_url, timeout=5)
            except RuntimeError as exc:
                if "unsupported URL scheme" not in str(exc):
                    raise AssertionError(f"Expected redirect URL scheme rejection, got: {exc}") from exc
            else:
                raise AssertionError("Expected MinerU Agent markdown redirect to file:// to be rejected")
        finally:
            if previous_base_url is None:
                os.environ.pop("STUDENT_OS_MINERU_AGENT_BASE_URL", None)
            else:
                os.environ["STUDENT_OS_MINERU_AGENT_BASE_URL"] = previous_base_url
            for key, value in previous_proxy_env.items():
                if value is None:
                    os.environ.pop(key, None)
                else:
                    os.environ[key] = value
    finally:
        redirect_server.shutdown()
        redirect_server.server_close()
        redirect_thread.join(timeout=5)
        secret.unlink(missing_ok=True)

    with FakeMineruAgentServer() as upload_server:
        upload_fd, upload_name = tempfile.mkstemp()
        os.close(upload_fd)
        upload_source = Path(upload_name)
        upload_source.write_bytes(b"student-os-upload")
        previous_base_url = os.environ.get("STUDENT_OS_MINERU_AGENT_BASE_URL")
        proxy_keys = ("NO_PROXY", "no_proxy")
        previous_proxy_env = {key: os.environ.get(key) for key in proxy_keys}
        try:
            os.environ["STUDENT_OS_MINERU_AGENT_BASE_URL"] = upload_server.base_url
            os.environ["NO_PROXY"] = "127.0.0.1,localhost,::1"
            os.environ["no_proxy"] = "127.0.0.1,localhost,::1"
            task_id = "manual-put"
            upload_server.tasks[task_id] = {"filename": "manual.pdf", "content": b"", "headers": ""}
            materials_convert.http_put_file(f"{upload_server.base_url.rsplit('/api/v1/agent', 1)[0]}/upload/{task_id}", upload_source)
            upload_headers = json.loads(str(upload_server.tasks[task_id]["headers"]))
            if any(key.lower() == "content-type" for key in upload_headers):
                raise AssertionError(f"MinerU Agent signed upload must not send Content-Type, got: {upload_headers}")
        finally:
            if previous_base_url is None:
                os.environ.pop("STUDENT_OS_MINERU_AGENT_BASE_URL", None)
            else:
                os.environ["STUDENT_OS_MINERU_AGENT_BASE_URL"] = previous_base_url
            for key, value in previous_proxy_env.items():
                if value is None:
                    os.environ.pop(key, None)
                else:
                    os.environ[key] = value
            upload_source.unlink(missing_ok=True)

    class FakeForwardProxy:
        def __init__(self) -> None:
            self.request_line: str | None = None
            self.server: http.server.ThreadingHTTPServer | None = None
            self.thread: threading.Thread | None = None

        @property
        def url(self) -> str:
            if self.server is None:
                raise RuntimeError("Fake proxy server is not running")
            host, port = self.server.server_address
            return f"http://{host}:{port}"

        def __enter__(self) -> "FakeForwardProxy":
            owner = self

            class Handler(http.server.BaseHTTPRequestHandler):
                def log_message(self, _format: str, *_args: object) -> None:
                    return

                def do_PUT(self) -> None:
                    owner.request_line = self.requestline
                    length = int(self.headers.get("Content-Length", "0"))
                    if length:
                        self.rfile.read(length)
                    self.send_response(200)
                    self.end_headers()

            self.server = http.server.ThreadingHTTPServer(("127.0.0.1", 0), Handler)
            self.thread = threading.Thread(target=self.server.serve_forever, daemon=True)
            self.thread.start()
            return self

        def __exit__(self, exc_type: object, exc: object, traceback: object) -> None:
            if self.server is not None:
                self.server.shutdown()
                self.server.server_close()
            if self.thread is not None:
                self.thread.join(timeout=5)

    with FakeForwardProxy() as proxy:
        proxy_fd, proxy_name = tempfile.mkstemp()
        os.close(proxy_fd)
        proxy_source = Path(proxy_name)
        proxy_source.write_bytes(b"proxy-upload")
        previous_base_url = os.environ.get("STUDENT_OS_MINERU_AGENT_BASE_URL")
        proxy_keys = ("HTTP_PROXY", "HTTPS_PROXY", "http_proxy", "https_proxy", "NO_PROXY", "no_proxy")
        previous_proxy_env = {key: os.environ.get(key) for key in proxy_keys}
        original_proxy_bypass = materials_convert.urllib.request.proxy_bypass
        materials_convert.urllib.request.proxy_bypass = lambda _host: False
        try:
            os.environ["STUDENT_OS_MINERU_AGENT_BASE_URL"] = "http://fake-target.example.com/api/v1/agent"
            os.environ["HTTP_PROXY"] = proxy.url
            os.environ["HTTPS_PROXY"] = proxy.url
            os.environ["http_proxy"] = proxy.url
            os.environ["https_proxy"] = proxy.url
            os.environ["NO_PROXY"] = ""
            os.environ["no_proxy"] = ""
            materials_convert.http_put_file(
                "http://fake-target.example.com/upload/proxy-task", proxy_source, timeout=5
            )
            expected_line = "PUT http://fake-target.example.com/upload/proxy-task HTTP/1.1"
            if proxy.request_line != expected_line:
                raise AssertionError(
                    f"Expected HTTP proxy request line {expected_line!r}, got {proxy.request_line!r}"
                )
        finally:
            materials_convert.urllib.request.proxy_bypass = original_proxy_bypass
            if previous_base_url is None:
                os.environ.pop("STUDENT_OS_MINERU_AGENT_BASE_URL", None)
            else:
                os.environ["STUDENT_OS_MINERU_AGENT_BASE_URL"] = previous_base_url
            for key, value in previous_proxy_env.items():
                if value is None:
                    os.environ.pop(key, None)
                else:
                    os.environ[key] = value
            proxy_source.unlink(missing_ok=True)

    class FakeProxyForHttpsTunnel:
        def __init__(self) -> None:
            self.connect_request: str | None = None
            self.tls_handshake_byte: bytes | None = None
            self._sock: socket.socket | None = None
            self._thread: threading.Thread | None = None
            self._stop = threading.Event()

        @property
        def url(self) -> str:
            if self._sock is None:
                raise RuntimeError("Fake proxy server is not running")
            host, port = self._sock.getsockname()
            return f"http://{host}:{port}"

        def _serve(self) -> None:
            self._sock.settimeout(1.0)
            while not self._stop.is_set():
                try:
                    conn, _addr = self._sock.accept()
                except (socket.timeout, OSError):
                    continue
                try:
                    conn.settimeout(2.0)
                    data = b""
                    while b"\r\n\r\n" not in data:
                        chunk = conn.recv(1024)
                        if not chunk:
                            break
                        data += chunk
                    if not data:
                        conn.close()
                        continue
                    first_line = data.split(b"\r\n", 1)[0].decode("latin-1")
                    self.connect_request = first_line
                    conn.sendall(b"HTTP/1.0 200 Connection established\r\n\r\n")
                    try:
                        self.tls_handshake_byte = conn.recv(1)
                    except OSError:
                        pass
                finally:
                    conn.close()

        def __enter__(self) -> "FakeProxyForHttpsTunnel":
            self._sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
            self._sock.bind(("127.0.0.1", 0))
            self._sock.listen(1)
            self._thread = threading.Thread(target=self._serve, daemon=True)
            self._thread.start()
            return self

        def __exit__(self, exc_type: object, exc: object, traceback: object) -> None:
            self._stop.set()
            if self._sock is not None:
                self._sock.close()
            if self._thread is not None:
                self._thread.join(timeout=5)

    with FakeProxyForHttpsTunnel() as proxy:
        proxy_fd, proxy_name = tempfile.mkstemp()
        os.close(proxy_fd)
        proxy_source = Path(proxy_name)
        proxy_source.write_bytes(b"proxy-upload")
        previous_base_url = os.environ.get("STUDENT_OS_MINERU_AGENT_BASE_URL")
        proxy_keys = ("HTTP_PROXY", "HTTPS_PROXY", "http_proxy", "https_proxy", "NO_PROXY", "no_proxy")
        previous_proxy_env = {key: os.environ.get(key) for key in proxy_keys}
        original_proxy_bypass = materials_convert.urllib.request.proxy_bypass
        materials_convert.urllib.request.proxy_bypass = lambda _host: False
        expected_prefix = "CONNECT fake-target.example.com:443 HTTP/1."
        try:
            os.environ["STUDENT_OS_MINERU_AGENT_BASE_URL"] = "https://fake-target.example.com/api/v1/agent"
            os.environ["HTTP_PROXY"] = proxy.url
            os.environ["HTTPS_PROXY"] = proxy.url
            os.environ["http_proxy"] = proxy.url
            os.environ["https_proxy"] = proxy.url
            os.environ["NO_PROXY"] = ""
            os.environ["no_proxy"] = ""
            try:
                materials_convert.http_put_file(
                    "https://fake-target.example.com/upload/proxy-task", proxy_source, timeout=5
                )
            except RuntimeError:
                # TLS handshake with the fake proxy will fail; we only care about CONNECT + TLS start.
                pass
        finally:
            materials_convert.urllib.request.proxy_bypass = original_proxy_bypass
            if previous_base_url is None:
                os.environ.pop("STUDENT_OS_MINERU_AGENT_BASE_URL", None)
            else:
                os.environ["STUDENT_OS_MINERU_AGENT_BASE_URL"] = previous_base_url
            for key, value in previous_proxy_env.items():
                if value is None:
                    os.environ.pop(key, None)
                else:
                    os.environ[key] = value
            proxy_source.unlink(missing_ok=True)
        if not (proxy.connect_request and proxy.connect_request.startswith(expected_prefix)):
            raise AssertionError(
                f"Expected HTTPS proxy CONNECT starting with {expected_prefix!r}, got {proxy.connect_request!r}"
            )
        if proxy.tls_handshake_byte != b"\x16":
            raise AssertionError(
                f"Expected TLS handshake first byte 0x16, got {proxy.tls_handshake_byte!r}"
            )

    with tempfile.TemporaryDirectory(prefix="student-os-agent-chunks-") as tmp:
        work_dir = Path(tmp)
        original_split_pdf = materials_convert.split_pdf
        original_limit = materials_convert.MINERU_AGENT_MAX_FILE_BYTES

        def fake_split_pdf(_source_file: Path, chunk_size: int, split_dir: Path) -> list[dict[str, object]]:
            for child in split_dir.glob("*"):
                child.unlink()
            sizes = [120, 50] if chunk_size > 10 else [80, 80, 40]
            chunks: list[dict[str, object]] = []
            start_page = 1
            for index, size in enumerate(sizes, start=1):
                path = split_dir / f"part{index}.pdf"
                path.write_bytes(b"x" * size)
                chunks.append(
                    {
                        "path": path,
                        "part_index": index,
                        "start_page": start_page,
                        "end_page": start_page + chunk_size - 1,
                    }
                )
                start_page += chunk_size
            return chunks

        try:
            materials_convert.split_pdf = fake_split_pdf
            materials_convert.MINERU_AGENT_MAX_FILE_BYTES = 100
            chunks, chunk_pages = materials_convert.split_pdf_for_mineru_agent(Path("source.pdf"), 20, work_dir)
            if chunk_pages != 10 or any(Path(chunk["path"]).stat().st_size > 100 for chunk in chunks):
                raise AssertionError("Expected MinerU Agent chunk splitter to shrink oversized chunks")
        finally:
            materials_convert.split_pdf = original_split_pdf
            materials_convert.MINERU_AGENT_MAX_FILE_BYTES = original_limit


def verify_mineru_v1_repair_rules() -> None:
    previous_sys_path = list(sys.path)
    sys.path.insert(0, str(STUDENT_OS_SCRIPTS))
    try:
        repair_module = load_student_os_script_module("repair_markdown_import.py", "student_os_repair_smoke")
    finally:
        sys.path = previous_sys_path

    noisy = (
        "---\n"
        "import_method: mineru-agent-v1\n"
        "---\n"
        "\n"
        "## Imported Content\n"
        "\n"
        "{ \\dot { 1 } } + \\dot { 0 } = \\dot { \\lambda }\n"
        "\n"
        "\\mathrm { ~ } and \\mathrm { ~ . ~ } here.\n"
        "\n"
        "\\overset { \\cdot } { A } \\stackrel { r } {  } = \\overset {  } { B }\n"
        "\n"
        "{ \\texttt { \\textbf { # } : } }\n"
        "\n"
        "\\sharp A\n"
        "\n"
        "$$\\beta_2$$\n"
        "\n"
        "$$\\{x\\}$$\n"
        "\n"
        "$$\\gamma$$\n"
        "\n"
        "Keep real math: $\\dot{x}$ and $\\stackrel{r}{=}$.\n"
        "\n"
        "Keep \\binom fragment for manual review.\n"
        "## 一. Promoted question heading\n"
        "$unterminated\n"
    )
    repaired, summary = repair_module.repair_text(noisy)

    if "{ 1 } + 0 = \\lambda" not in repaired:
        raise AssertionError("Expected \\dot{} wrappers around digits/Greek to be removed")
    if "\\mathrm" in repaired:
        raise AssertionError("Expected \\mathrm{} noise wrappers to be removed")
    if "\\overset { \\cdot } { A }" in repaired or "\\stackrel { r } {  }" in repaired:
        raise AssertionError("Expected empty \\overset/\\stackrel wrappers to be simplified")
    if "\\texttt" in repaired:
        raise AssertionError("Expected garbled \\texttt placeholder to be replaced")
    if "\\sharp A" in repaired:
        raise AssertionError("Expected text-mode \\sharp to be replaced")
    if "$$\\beta_2 \\{x\\} \\gamma$$" not in repaired:
        raise AssertionError("Expected isolated single-line $$ fragments to be collapsed")
    if "$\\dot{x}$" not in repaired or "$\\stackrel{r}{=}$" not in repaired:
        raise AssertionError("Expected legitimate math accents to be preserved")
    if "\\binom" not in repaired:
        raise AssertionError("Expected \\binom fragments to remain for manual review")
    if not any("remaining noise signatures" in line for line in summary):
        raise AssertionError("Expected repair summary to report remaining noise signatures")

    risks = repair_module.diagnose_import_risks(repaired)
    risk_codes = {str(item["code"]) for item in risks}
    for expected in {
        "latex-binom-fragment",
        "lossy-ocr-placeholder",
        "question-heading-promoted",
        "math-dollar-unbalanced",
    }:
        if expected not in risk_codes:
            raise AssertionError(f"Expected repair diagnostics to report {expected}, got: {risks}")
    valid_binom = "---\nimport_method: mineru-agent-v1\n---\n\nA valid formula: $\\binom{n}{k}$.\n"
    if any(item["code"] == "latex-binom-fragment" for item in repair_module.diagnose_import_risks(valid_binom)):
        raise AssertionError("Valid \\binom{n}{k} formulas must not be flagged as malformed fragments")
    malformed_binom = "---\nimport_method: mineru-agent-v1\n---\n\nBroken formula: $\\binom fragment$.\n"
    if not any(item["code"] == "latex-binom-fragment" for item in repair_module.diagnose_import_risks(malformed_binom)):
        raise AssertionError("Malformed \\binom fragments should be flagged for review")

    english_import = "---\nimport_method: mineru-agent-v1\n---\n\n" + ("plain English text " * 160)
    if any(item["code"] == "low-cjk-density" for item in repair_module.diagnose_import_risks(english_import)):
        raise AssertionError("English imports must not be flagged as low-CJK-density without an expected CJK language marker")
    cjk_expected_import = "---\nimport_method: mineru-agent-v1\nlanguage: ch\n---\n\n" + ("plain English text " * 160)
    if not any(item["code"] == "low-cjk-density" for item in repair_module.diagnose_import_risks(cjk_expected_import)):
        raise AssertionError("Expected CJK-language imports with no CJK text to be flagged")
    verified_with_comment = "---\nverify_status: verified # checked against source\n---\n\nBody.\n"
    if not repair_module.is_verified(verified_with_comment):
        raise AssertionError("verify_status with a YAML inline comment should still be treated as verified")

    missing_repair_status = (
        "---\n"
        "type: imported-reference\n"
        "verify_status: legacy-value\n"
        "---\n"
        "\n"
        "Quoted body metadata must remain intact:\n"
        "verify_status: verified\n"
    )
    marked = repair_module.mark_auto_repaired(missing_repair_status, needs_review=True)
    frontmatter = marked.split("---\n", 2)[1]
    body = marked.split("---\n", 2)[2]
    if "repair_status: auto-repaired" not in frontmatter:
        raise AssertionError(f"Missing repair_status should be inserted into frontmatter:\n{marked}")
    if "verify_status: unverified" not in frontmatter:
        raise AssertionError(f"verify_status should be replaced inside frontmatter:\n{marked}")
    if "repair_risk: needs-human-review" not in frontmatter:
        raise AssertionError(f"repair_risk should be inserted into frontmatter:\n{marked}")
    if "verify_status: verified" not in body:
        raise AssertionError(f"Body metadata-looking text must not be rewritten:\n{marked}")


def verify_local_pdf_risk_forwarding(tmp_root: Path) -> None:
    previous_sys_path = list(sys.path)
    sys.path.insert(0, str(STUDENT_OS_SCRIPTS))
    try:
        materials_convert = load_student_os_script_module(
            "materials_convert.py",
            "student_os_materials_convert_pdf_risk_smoke",
        )
    finally:
        sys.path = previous_sys_path

    wrapped = materials_convert.wrap_mineru_markdown(
        source_file=tmp_root / "chinese-source.pdf",
        markdown_body="plain English text " * 160,
        import_method="mineru-api:pipeline",
        course=None,
        language="ch",
    )
    if "language: ch" not in wrapped:
        raise AssertionError(f"MinerU sidecars should persist the requested language:\n{wrapped}")
    if not any(item["code"] == "low-cjk-density" for item in materials_convert.diagnose_import_risks(wrapped)):
        raise AssertionError("MinerU sidecars with requested CJK language should surface low-CJK repair risks")

    output = tmp_root / "paper.pdf.md"
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text("course:\n", encoding="utf-8", newline="\n")
    original_run_script = materials_convert.run_script
    try:
        materials_convert.run_script = lambda *_args: {
            "output": str(output),
            "raw_output": str(tmp_root / "paper.pdf.raw.md"),
            "repair_summary": str(tmp_root / "paper.pdf-repair-summary.md"),
            "repairs": ["Manual review risk items:"],
            "risk_items": [{"code": "latex-binom-fragment", "count": 1}],
        }
        ctx = materials_convert.ConversionContext(
            method="local",
            course=None,
            overwrite=False,
            repair=False,
            repair_only=False,
            api_token=None,
            api_model="vlm",
            language="ch",
            ocr=None,
            formula=None,
            table=None,
            pages=None,
            timeout=300,
        )
        payload = materials_convert.convert_with_local_plan(
            tmp_root / "paper.pdf",
            output,
            materials_convert.ConversionPlan(kind="pdf", import_method="pdf-to-md"),
            ctx,
        )
    finally:
        materials_convert.run_script = original_run_script
    if not payload.get("risk_items"):
        raise AssertionError(f"Local PDF conversion should forward repair risk_items: {payload}")


def exercise_import_workflows(repo: Path) -> None:
    verify_material_type_constants()
    verify_mineru_agent_helper_guards()
    verify_mineru_v1_repair_rules()
    verify_local_pdf_risk_forwarding(repo / "references" / "imports" / "pdf-risk-forwarding")
    fixture_root = repo / "references" / "imports" / "source"
    fixture_root.mkdir(parents=True, exist_ok=True)
    docx_path = fixture_root / "linear-algebra-outline.docx"
    xlsx_path = fixture_root / "linear-algebra-progress.xlsx"
    pptx_path = fixture_root / "linear-algebra-week-2.pptx"
    pdf_path = fixture_root / "linear-algebra-handout.pdf"
    png_path = fixture_root / "homework-photo.png"
    binary_path = fixture_root / "fpga-lab.bit"
    write_docx_fixture(docx_path)
    write_xlsx_fixture(xlsx_path)
    write_pptx_fixture(pptx_path)
    write_pdf_fixture(pdf_path)
    write_png_fixture(png_path)
    binary_path.write_bytes(b"\x00\x01\x02fpga")

    docx_output = repo / "courses" / "linear-algebra" / "references" / "outline-import.md"
    course_repair_summary = repo / "courses" / "linear-algebra" / "references" / "outline-import-repair-summary.md"
    xlsx_output = repo / "dashboards" / "linear-algebra-progress-import.md"
    pptx_output = repo / "references" / "slides" / "linear-algebra-week-2.md"
    pdf_generic_output = repo / "courses" / "linear-algebra" / "references" / "handout-generic-import.md"
    pdf_mineru_output = repo / "references" / "imports" / "raw" / "linear-algebra-handout.md"
    repair_input = repo / "references" / "imports" / "raw" / "manual-repair-sample.md"
    repair_output = repo / "references" / "imports" / "repaired" / "manual-repair-sample.md"
    repair_summary = repo / "references" / "imports" / "repaired" / "manual-repair-sample-repair-summary.md"

    docx_payload = json.loads(run_script("docx_to_md.py", str(docx_path), "--output", str(docx_output)))
    xlsx_payload = json.loads(run_script("xlsx_to_md.py", str(xlsx_path), "--output", str(xlsx_output), "--max-rows", "3"))
    pptx_payload = json.loads(run_script("pptx_to_md.py", str(pptx_path), "--output", str(pptx_output)))
    pdf_probe_payload = json.loads(run_script("pdf_probe.py", str(pdf_path)))
    pdf_generic_payload = json.loads(
        run_script("pdf_to_markdown.py", str(pdf_path), "--output", str(pdf_generic_output), "--mode", "generic")
    )
    pdf_mineru_payload = json.loads(
        run_script("pdf_to_markdown.py", str(pdf_path), "--output", str(pdf_mineru_output), "--mode", "mineru-style")
    )
    with FakeMineruAgentServer() as default_agent_server:
        materials_payload = json.loads(
            run_script(
                "materials_convert.py",
                str(fixture_root),
                "--course",
                "Linear Algebra",
                env={
                    "MINERU_TOKEN": "",
                    "MINERU_API_TOKEN": "",
                    "STUDENT_OS_MINERU_AGENT_BASE_URL": default_agent_server.base_url,
                },
            )
        )
    materials_repair_payload = json.loads(
        run_script(
            "materials_convert.py",
            str(docx_path),
            "--course",
            "Linear Algebra",
            "--output-root",
            str(repo / "references" / "imports" / "repair-output"),
            "--repair",
        )
    )

    repair_input.write_text(
        "\n".join(
            [
                "---",
                "type: pdf-import-note",
                "course:",
                "status: active",
                "created:",
                "updated:",
                "tags: [import, pdf]",
                f"source_file: {json.dumps(str(pdf_path))}",
                "import_method: manual-test",
                "repair_status: raw",
                'derived_from_import: ""',
                "---",
                "",
                "#Broken heading",
                "",
                "Page 1",
                "",
                "-  Bullet with extra spacing",
                "",
                "## 一. OCR promoted question heading",
                "",
                "Keep \\binom fragment for manual review.",
                "",
                "## Next section ........ 4",
                "",
            ]
        ),
        encoding="utf-8",
        newline="\n",
    )
    repair_payload = json.loads(
        run_script(
            "repair_markdown_import.py",
            str(repair_input),
            "--output",
            str(repair_output),
            "--summary-path",
            str(repair_summary),
            "--derived-from",
            str(repair_input),
        )
    )

    ensure_exists(Path(docx_payload["output"]))
    ensure_contains(docx_output, "Linear Algebra Import Outline")
    ensure_contains(docx_output, "## Table 1")
    course_repair_summary.write_text("# Repair Summary\n\n- Example only.\n", encoding="utf-8", newline="\n")
    ensure_exists(Path(xlsx_payload["output"]))
    ensure_contains(xlsx_output, "| Task | Score | Weight | Weighted |")
    ensure_contains(xlsx_output, "=B2*C2")
    ensure_contains(xlsx_output, "Truncated after 3 rows.")
    ensure_exists(Path(pptx_payload["output"]))
    ensure_contains(pptx_output, "## Slide 1: Linear Algebra Week 2")
    ensure_contains(pptx_output, "Diagonalization and basis changes")
    if pdf_probe_payload["page_count"] != 1:
        raise AssertionError(f"Expected a one-page PDF fixture, got: {pdf_probe_payload}")
    ensure_exists(Path(pdf_generic_payload["output"]))
    ensure_contains(pdf_generic_output, "Import method: generic")
    ensure_contains(pdf_generic_output, "#Broken linear algebra handout")
    ensure_exists(Path(pdf_mineru_payload["output"]))
    ensure_contains(repo / "references" / "imports" / "raw" / "linear-algebra-handout.md", "repair_status: raw")
    ensure_contains(repo / "references" / "imports" / "repaired" / "linear-algebra-handout.md", "repair_status: auto-repaired")
    ensure_contains(repo / "references" / "imports" / "repaired" / "linear-algebra-handout.md", "verify_status: unverified")
    ensure_contains(repo / "references" / "imports" / "repaired" / "linear-algebra-handout-repair-summary.md", "# Repair Summary")
    ensure_contains(repo / "references" / "imports" / "repaired" / "linear-algebra-handout-repair-summary.md", "Removed isolated page labels.")
    ensure_contains(repo / "references" / "imports" / "repaired" / "linear-algebra-handout-repair-summary.md", "Normalized heading spacing.")
    ensure_exists(Path(repair_payload["output"]))
    if not repair_payload.get("risk_items"):
        raise AssertionError(f"Expected repair payload to include risk_items for promoted headings: {repair_payload}")
    ensure_contains(repair_output, "# Broken heading")
    ensure_contains(repair_output, "- Bullet with extra spacing")
    ensure_contains(repair_output, "repair_status: auto-repaired")
    ensure_contains(repair_output, "verify_status: unverified")
    ensure_contains(repair_output, "repair_risk: needs-human-review")
    ensure_contains(repair_summary, "Removed isolated page labels.")
    ensure_contains(repair_summary, "Normalized heading spacing.")
    ensure_contains(repair_summary, "Trimmed heading dot leaders or page-number residue.")
    ensure_contains(repair_summary, "Manual review risk items:")

    legacy_repaired_input = repo / "references" / "imports" / "repaired" / "legacy-repaired-input.md"
    legacy_repaired_output = repo / "references" / "imports" / "repaired" / "legacy-repaired-output.md"
    legacy_repaired_input.write_text(
        "\n".join(
            [
                "---",
                "type: imported-reference",
                "repair_status: repaired",
                "derived_from_import:",
                "---",
                "",
                "#Legacy heading",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    legacy_repaired_payload = json.loads(
        run_script(
            "repair_markdown_import.py",
            str(legacy_repaired_input),
            "--output",
            str(legacy_repaired_output),
        )
    )
    if legacy_repaired_payload.get("skipped"):
        raise AssertionError(f"Legacy repaired status must not be treated as verified: {legacy_repaired_payload}")
    ensure_contains(legacy_repaired_output, "repair_status: auto-repaired")
    ensure_contains(legacy_repaired_output, "verify_status: unverified")

    unicode_derived_root = repo / "references" / "imports" / "用户资料"
    unicode_derived_root.mkdir(parents=True, exist_ok=True)
    unicode_derived_source = unicode_derived_root / "raw-import.md"
    unicode_repair_input = repo / "references" / "imports" / "repaired" / "unicode-path-repair-input.md"
    unicode_repair_output = repo / "references" / "imports" / "repaired" / "unicode-path-repair-output.md"
    unicode_derived_source.write_text("# raw\n", encoding="utf-8", newline="\n")
    unicode_repair_input.write_text(
        "\n".join(
            [
                "---",
                "type: pdf-import-note",
                "course:",
                "status: active",
                "created:",
                "updated:",
                "tags: [import, pdf]",
                'source_file: "sample.pdf"',
                "import_method: manual-test",
                "repair_status:",
                "derived_from_import:",
                "---",
                "",
                "#Broken heading",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    unicode_repair_payload = json.loads(
        run_script(
            "repair_markdown_import.py",
            str(unicode_repair_input),
            "--output",
            str(unicode_repair_output),
            "--derived-from",
            str(unicode_derived_source),
        )
    )
    ensure_exists(Path(unicode_repair_payload["output"]))
    ensure_contains(unicode_repair_output, "derived_from_import:")
    ensure_contains(unicode_repair_output, "raw-import.md")
    unicode_output_text = unicode_repair_output.read_text(encoding="utf-8")
    if "用户资料" not in unicode_output_text or "\\u7528\\u6237" in unicode_output_text:
        raise AssertionError("Expected unicode path segment to remain literal in derived_from_import")
    ensure_contains(unicode_repair_output, "repair_status: auto-repaired")
    ensure_contains(unicode_repair_output, "verify_status: unverified")
    if 'derived_from_import:\n' in unicode_output_text or 'derived_from_import: ""' in unicode_output_text:
        raise AssertionError("derived_from_import should be filled for unicode Windows paths")

    if len(materials_payload["converted"]) != 6:
        raise AssertionError(f"Expected six converted material outputs, got: {materials_payload}")
    ensure_exists(fixture_root / "linear-algebra-outline.docx.md")
    ensure_contains(fixture_root / "linear-algebra-outline.docx.md", 'course: "Linear Algebra"')
    ensure_exists(fixture_root / "linear-algebra-progress.xlsx.md")
    ensure_contains(fixture_root / "linear-algebra-progress.xlsx.md", "| Task | Score | Weight | Weighted |")
    ensure_exists(fixture_root / "linear-algebra-week-2.pptx.md")
    ensure_contains(fixture_root / "linear-algebra-week-2.pptx.md", "## Slide 1: Linear Algebra Week 2")
    ensure_exists(fixture_root / "linear-algebra-handout.pdf.md")
    ensure_contains(fixture_root / "linear-algebra-handout.pdf.md", "Import method: mineru-agent-v1")
    ensure_contains(fixture_root / "linear-algebra-handout.pdf.md", "# Agent Parsed - linear-algebra-handout.pdf")
    ensure_not_contains(fixture_root / "linear-algebra-handout.pdf.md", "- uploaded_bytes: 0")
    ensure_exists(fixture_root / "homework-photo.png.md")
    ensure_contains(fixture_root / "homework-photo.png.md", "Import method: mineru-agent-v1")
    ensure_contains(fixture_root / "homework-photo.png.md", "# Agent Parsed - homework-photo.png")
    ensure_not_contains(fixture_root / "homework-photo.png.md", "- uploaded_bytes: 0")
    ensure_exists(fixture_root / "fpga-lab.bit.md")
    ensure_contains(fixture_root / "fpga-lab.bit.md", "Binary or tool-specific source detected.")
    if len(materials_repair_payload["converted"]) != 1:
        raise AssertionError(f"Expected one repaired material output, got: {materials_repair_payload}")
    repair_output_root = repo / "references" / "imports" / "repair-output"
    ensure_exists(repair_output_root / "linear-algebra-outline.docx.md")
    ensure_exists(repair_output_root / "linear-algebra-outline.docx.raw.md")
    ensure_exists(repair_output_root / "linear-algebra-outline.docx-repair-summary.md")
    ensure_contains(repair_output_root / "linear-algebra-outline.docx.md", "repair_status: auto-repaired")
    ensure_contains(repair_output_root / "linear-algebra-outline.docx.md", "verify_status: unverified")
    ensure_contains(repair_output_root / "linear-algebra-outline.docx-repair-summary.md", "# Repair Summary")
    ensure_contains(repair_output_root / "linear-algebra-outline.docx-repair-summary.md", "Output:")

    fake_sdk_root = repo / "references" / "imports" / "fake-sdk"
    fake_sdk_root.mkdir(parents=True, exist_ok=True)
    write_fake_mineru_sdk(fake_sdk_root)
    api_output_root = repo / "references" / "imports" / "api-output"
    api_payload = json.loads(
        run_script(
            "materials_convert.py",
            str(fixture_root),
            "--output-root",
            str(api_output_root),
            "--course",
            "Linear Algebra",
            "--method",
            "api",
            "--api-token",
            "test-token",
            "--language",
            "en",
            "--pages",
            "1-1",
            env={"PYTHONPATH": str(fake_sdk_root)},
        )
    )
    if api_payload["applied_method"] != "api":
        raise AssertionError(f"Expected API mode to be applied, got: {api_payload}")
    ensure_exists(api_output_root / "linear-algebra-outline.docx.md")
    ensure_contains(api_output_root / "linear-algebra-outline.docx.md", "Import method: mineru-api:vlm")
    ensure_contains(api_output_root / "linear-algebra-outline.docx.md", "# API Parsed - linear-algebra-outline.docx")
    ensure_exists(api_output_root / "homework-photo.png.md")
    ensure_contains(api_output_root / "homework-photo.png.md", "# API Parsed - homework-photo.png")
    ensure_exists(api_output_root / "linear-algebra-handout.pdf.md")
    ensure_contains(api_output_root / "linear-algebra-handout.pdf.md", "- pages: 1-1")
    ensure_exists(api_output_root / "fpga-lab.bit.md")
    ensure_contains(api_output_root / "fpga-lab.bit.md", "Binary or tool-specific source detected.")
    api_repair_payload = json.loads(
        run_script(
            "materials_convert.py",
            str(docx_path),
            "--output-root",
            str(repo / "references" / "imports" / "api-repair-output"),
            "--course",
            "Linear Algebra",
            "--method",
            "api",
            "--api-token",
            "test-token",
            "--repair",
            env={"PYTHONPATH": str(fake_sdk_root)},
        )
    )
    if len(api_repair_payload["converted"]) != 1:
        raise AssertionError(f"Expected one API repaired output, got: {api_repair_payload}")
    api_repair_root = repo / "references" / "imports" / "api-repair-output"
    ensure_exists(api_repair_root / "linear-algebra-outline.docx.md")
    ensure_exists(api_repair_root / "linear-algebra-outline.docx.raw.md")
    ensure_exists(api_repair_root / "linear-algebra-outline.docx-repair-summary.md")
    ensure_contains(api_repair_root / "linear-algebra-outline.docx.md", "repair_status: auto-repaired")
    ensure_contains(api_repair_root / "linear-algebra-outline.docx.md", "verify_status: unverified")

    large_pdf_path = fixture_root / "large-textbook.pdf"
    write_multipage_pdf_fixture(large_pdf_path, 5)
    split_output_root = repo / "references" / "imports" / "api-split-output"
    split_payload = json.loads(
        run_script(
            "materials_convert.py",
            str(large_pdf_path),
            "--output-root",
            str(split_output_root),
            "--course",
            "Linear Algebra",
            "--method",
            "api",
            "--api-token",
            "test-token",
            "--chunk-size",
            "2",
            env={"PYTHONPATH": str(fake_sdk_root)},
        )
    )
    if len(split_payload["converted"]) != 1:
        raise AssertionError(f"Expected one auto-split conversion, got: {split_payload}")
    split_info = split_payload["converted"][0].get("split")
    if not split_info or split_info.get("part_count") != 3 or not split_info.get("merged"):
        raise AssertionError(f"Expected 3 merged PDF chunks, got: {split_payload['converted'][0]}")
    split_output = split_output_root / "large-textbook.pdf.md"
    ensure_exists(split_output)
    ensure_contains(split_output, "<!-- MERGED from 3 parts: pages 1-2, pages 3-4, pages 5-5 -->")
    ensure_contains(split_output, "# API Parsed - large-textbook.pdf.part1.pdf")
    ensure_contains(split_output, "# API Parsed - large-textbook.pdf.part2.pdf")
    ensure_contains(split_output, "# API Parsed - large-textbook.pdf.part3.pdf")
    leftover_parts = list(split_output_root.glob("**/*.part*.pdf"))
    if leftover_parts:
        raise AssertionError(f"Auto-split should clean temporary PDF chunks, found: {leftover_parts}")
    # Each chunk emits an image named image1.png; they must not overwrite each other.
    split_image_dir = split_output_root / "images"
    for part_index in (1, 2, 3):
        ensure_exists(split_image_dir / f"part{part_index}-image1.png")
        ensure_contains(split_output, f"images/part{part_index}-image1.png")
    if (split_image_dir / "image1.png").exists():
        raise AssertionError("Auto-split should not leave unprefixed chunk images that collide across parts")

    no_split_output = run_path_script_failure(
        STUDENT_OS_SCRIPTS / "materials_convert.py",
        str(large_pdf_path),
        "--output-root",
        str(repo / "references" / "imports" / "api-no-split-output"),
        "--method",
        "api",
        "--api-token",
        "test-token",
        "--chunk-size",
        "2",
        "--no-auto-split",
        "--overwrite",
        cwd=ROOT,
        env={"PYTHONPATH": str(fake_sdk_root)},
    )
    no_split_payload = json.loads(no_split_output)
    if not no_split_payload.get("errors"):
        raise AssertionError(f"Expected --no-auto-split to record conversion errors, got: {no_split_payload}")
    if "exceeds --chunk-size" not in no_split_payload["errors"][0].get("error", ""):
        raise AssertionError(f"Expected chunk-size error message, got: {no_split_payload['errors']}")

    dotenv_cwd = repo / "references" / "imports" / "dotenv-cwd"
    dotenv_cwd.mkdir(parents=True, exist_ok=True)
    empty_skill_root = dotenv_cwd / "empty-skill-root"
    empty_skill_root.mkdir(parents=True, exist_ok=True)
    (dotenv_cwd / ".env").write_text('MINERU_TOKEN="dotenv-from-cwd"\n', encoding="utf-8", newline="\n")
    dotenv_output_root = repo / "references" / "imports" / "api-dotenv-output"
    dotenv_payload = json.loads(
        run_script(
            "materials_convert.py",
            str(pdf_path),
            "--output-root",
            str(dotenv_output_root),
            "--method",
            "api",
            "--overwrite",
            cwd=dotenv_cwd,
            env={
                "PYTHONPATH": str(fake_sdk_root),
                "MINERU_TOKEN": "",
                "MINERU_API_TOKEN": "",
                "STUDENT_OS_SKILL_ROOT": str(empty_skill_root),
            },
        )
    )
    if dotenv_payload["applied_method"] != "api":
        raise AssertionError(f"Expected cwd .env token to enable API mode, got: {dotenv_payload}")
    ensure_exists(dotenv_output_root / "linear-algebra-handout.pdf.md")
    ensure_contains(dotenv_output_root / "linear-algebra-handout.pdf.md", "Import method: mineru-api:vlm")

    # A --chunk-size above MinerU's 200-page limit must still split PDFs over 200 pages.
    over_limit_pdf = fixture_root / "over-limit.pdf"
    write_multipage_pdf_fixture(over_limit_pdf, 201)
    cap_payload = json.loads(
        run_script(
            "materials_convert.py",
            str(over_limit_pdf),
            "--output-root",
            str(repo / "references" / "imports" / "api-cap-output"),
            "--method",
            "api",
            "--api-token",
            "test-token",
            "--chunk-size",
            "500",
            env={"PYTHONPATH": str(fake_sdk_root)},
        )
    )
    cap_split = cap_payload["converted"][0].get("split")
    if not cap_split or cap_split.get("part_count") != 2 or cap_split.get("chunk_size") != 200:
        raise AssertionError(f"Expected --chunk-size to be capped at 200 pages, got: {cap_payload['converted'][0]}")
    over_limit_pdf.unlink()

    # A PDF within the page limit but over the byte cap must still be split by size.
    size_split_payload = json.loads(
        run_script(
            "materials_convert.py",
            str(large_pdf_path),
            "--output-root",
            str(repo / "references" / "imports" / "api-size-split-output"),
            "--method",
            "api",
            "--api-token",
            "test-token",
            "--chunk-size",
            "50",
            env={"PYTHONPATH": str(fake_sdk_root), "STUDENT_OS_MINERU_MAX_FILE_BYTES": "1"},
        )
    )
    size_split = size_split_payload["converted"][0].get("split")
    if not size_split or size_split.get("part_count") != 5:
        raise AssertionError(f"Expected size-driven split into single-page chunks, got: {size_split_payload['converted'][0]}")

    # A malformed PDF must be reported per-file (JSON stdout), not abort the batch.
    bad_pdf = fixture_root / "corrupt.pdf"
    bad_pdf.write_bytes(b"%PDF-1.4 not really a pdf")
    bad_result = subprocess.run(
        [
            sys.executable,
            "-B",
            str(STUDENT_OS_SCRIPTS / "materials_convert.py"),
            str(bad_pdf),
            "--output-root",
            str(repo / "references" / "imports" / "api-bad-pdf-output"),
            "--method",
            "api",
            "--api-token",
            "test-token",
        ],
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=ROOT,
        env={**os.environ, "PYTHONDONTWRITEBYTECODE": "1", "PYTHONIOENCODING": "utf-8", "PYTHONPATH": str(fake_sdk_root)},
    )
    if bad_result.returncode == 0:
        raise AssertionError("Expected malformed PDF to produce a non-zero exit code")
    bad_payload = json.loads(bad_result.stdout)
    if not bad_payload.get("errors") or "page-count probe" not in bad_payload["errors"][0].get("error", ""):
        raise AssertionError(f"Expected malformed PDF to be captured in errors, got: {bad_payload}")
    bad_pdf.unlink()

    # --no-merge must not clobber existing chunk sidecars without --overwrite, and
    # --repair must run on each chunk sidecar.
    no_merge_root = repo / "references" / "imports" / "api-no-merge-output"
    no_merge_payload = json.loads(
        run_script(
            "materials_convert.py",
            str(large_pdf_path),
            "--output-root",
            str(no_merge_root),
            "--method",
            "api",
            "--api-token",
            "test-token",
            "--chunk-size",
            "2",
            "--no-merge",
            "--repair",
            env={"PYTHONPATH": str(fake_sdk_root)},
        )
    )
    no_merge_converted = no_merge_payload["converted"][0]
    if not no_merge_converted.get("part_repairs") or len(no_merge_converted["part_repairs"]) != 3:
        raise AssertionError(f"Expected --no-merge --repair to repair each chunk, got: {no_merge_converted}")
    for part_index in (1, 2, 3):
        ensure_exists(no_merge_root / f"large-textbook.pdf.part{part_index}.md")
        ensure_exists(no_merge_root / f"large-textbook.pdf.part{part_index}-repair-summary.md")

    sentinel_part = no_merge_root / "large-textbook.pdf.part1.md"
    sentinel_part.write_text("SENTINEL EDIT\n", encoding="utf-8")
    rerun_payload = json.loads(
        run_script(
            "materials_convert.py",
            str(large_pdf_path),
            "--output-root",
            str(no_merge_root),
            "--method",
            "api",
            "--api-token",
            "test-token",
            "--chunk-size",
            "2",
            "--no-merge",
            env={"PYTHONPATH": str(fake_sdk_root)},
        )
    )
    if not rerun_payload["skipped"] or rerun_payload["skipped"][0].get("reason") != "output-exists":
        raise AssertionError(f"Expected --no-merge rerun to skip existing chunk sidecars, got: {rerun_payload}")
    if sentinel_part.read_text(encoding="utf-8") != "SENTINEL EDIT\n":
        raise AssertionError("--no-merge rerun without --overwrite must not clobber existing chunk sidecars")

    repair_only_input = repo / "references" / "imports" / "repair-only-sample.md"
    repair_only_input.write_text(
        "\n".join(
            [
                "---",
                "type: imported-reference",
                "course:",
                "status: active",
                "created:",
                "updated:",
                "tags: [import, reference]",
                'source_file: "sample.pdf"',
                "import_method: manual-test",
                "repair_status:",
                "derived_from_import:",
                "---",
                "",
                "#Broken heading",
                "",
                "Page 1",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    repair_only_payload = json.loads(
        run_script(
            "materials_convert.py",
            str(repair_only_input),
            "--repair-only",
        )
    )
    if len(repair_only_payload["converted"]) != 1:
        raise AssertionError(f"Expected one repair-only output, got: {repair_only_payload}")
    ensure_contains(repair_only_input, "# Broken heading")
    ensure_contains(repair_only_input, "repair_status: auto-repaired")
    ensure_contains(repair_only_input, "verify_status: unverified")
    ensure_exists(repo / "references" / "imports" / "repair-only-sample-repair-summary.md")
    ensure_contains(repo / "references" / "imports" / "repair-only-sample-repair-summary.md", "Removed isolated page labels.")

    verified_repair_input = repo / "references" / "imports" / "verified-repair-sample.md"
    verified_repair_input.write_text(
        "\n".join(
            [
                "---",
                "type: imported-reference",
                "repair_status: auto-repaired",
                "verify_status: verified",
                "derived_from_import:",
                "---",
                "",
                "#Broken verified heading",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    verified_payload = json.loads(run_script("materials_convert.py", str(verified_repair_input), "--repair-only"))
    if verified_payload["converted"] or not verified_payload["skipped_verified"]:
        raise AssertionError(f"Verified repair-only file should be skipped by default: {verified_payload}")
    ensure_contains(verified_repair_input, "#Broken verified heading")
    include_verified_payload = json.loads(
        run_script("materials_convert.py", str(verified_repair_input), "--repair-only", "--include-verified")
    )
    if len(include_verified_payload["converted"]) != 1:
        raise AssertionError(f"--include-verified should process verified files: {include_verified_payload}")
    ensure_contains(verified_repair_input, "# Broken verified heading")
    ensure_contains(verified_repair_input, "verify_status: unverified")

    verified_distinct_input = repo / "references" / "imports" / "verified-distinct-input.md"
    verified_distinct_input.write_text(
        "\n".join(
            [
                "---",
                "type: imported-reference",
                "repair_status: auto-repaired",
                "verify_status: verified",
                "derived_from_import:",
                "---",
                "",
                "#Broken verified heading",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    verified_distinct_output = repo / "references" / "imports" / "verified-distinct-output.md"
    verified_skip_payload = json.loads(
        run_script(
            "repair_markdown_import.py",
            str(verified_distinct_input),
            "--output",
            str(verified_distinct_output),
        )
    )
    if not verified_skip_payload.get("skipped") or verified_skip_payload.get("output") is not None:
        raise AssertionError(f"Verified skip should not claim a generated distinct output: {verified_skip_payload}")
    if verified_distinct_output.exists():
        raise AssertionError("Verified skip must not create a distinct output file")

    non_utf8_repair_input = repo / "references" / "imports" / "non-utf8-repair.md"
    non_utf8_repair_input.write_bytes(b"\xff\xfe# bad encoding\n")
    non_utf8_payload = json.loads(
        run_path_script_failure(
            STUDENT_OS_SCRIPTS / "materials_convert.py",
            str(non_utf8_repair_input),
            "--repair-only",
            cwd=ROOT,
        )
    )
    if non_utf8_payload["converted"] or not non_utf8_payload["errors"]:
        raise AssertionError(f"Non-UTF-8 repair-only files should be reported per-file as errors: {non_utf8_payload}")

    # Content-aware probing / smart routing.
    no_token_env = {
        "MINERU_TOKEN": "",
        "MINERU_API_TOKEN": "",
        "STUDENT_OS_SKILL_ROOT": str(empty_skill_root),
    }
    probe_only_payload = json.loads(
        run_script(
            "materials_convert.py",
            str(fixture_root),
            "--probe-only",
            cwd=empty_skill_root,
            env=no_token_env,
        )
    )
    if not probe_only_payload.get("probe_only") or not probe_only_payload.get("probes"):
        raise AssertionError(f"Expected --probe-only JSON report, got: {probe_only_payload}")
    probes_by_name = {Path(item["source"]).name: item for item in probe_only_payload["probes"]}
    if probes_by_name["linear-algebra-handout.pdf"]["tool"] != "mineru-api" or not probes_by_name[
        "linear-algebra-handout.pdf"
    ].get("agent_api"):
        raise AssertionError(
            f"Without a token, API-needed PDF probes should use MinerU v1 Agent, got: {probes_by_name['linear-algebra-handout.pdf']}"
        )
    if probes_by_name["homework-photo.png"]["tool"] != "mineru-api" or not probes_by_name["homework-photo.png"].get(
        "agent_api"
    ):
        raise AssertionError(f"Without a token, image probes should use MinerU v1 Agent, got: {probes_by_name['homework-photo.png']}")
    if probes_by_name["linear-algebra-outline.docx"]["tool"] not in {"pandoc", "docx-to-md"}:
        raise AssertionError(f"Unexpected DOCX tool: {probes_by_name['linear-algebra-outline.docx']}")

    scanned_pdf = fixture_root / "scanned-blank.pdf"
    write_multipage_pdf_fixture(scanned_pdf, 2)
    scanned_probe = json.loads(
        run_script(
            "materials_convert.py",
            str(scanned_pdf),
            "--probe-only",
            "--api-token",
            "test-token",
        )
    )["probes"][0]
    if (
        scanned_probe["strategy"] != "scanned"
        or scanned_probe.get("tool") != "mineru-api"
        or not scanned_probe.get("needs_ocr")
    ):
        raise AssertionError(f"Expected blank PDF to probe as scanned+mineru-api+OCR, got: {scanned_probe}")

    manual_pdf = fixture_root / "text-manual.pdf"
    write_text_heavy_pdf_fixture(manual_pdf, 3)
    manual_probe = json.loads(
        run_script(
            "materials_convert.py",
            str(manual_pdf),
            "--probe-only",
            cwd=empty_skill_root,
            env={
                **no_token_env,
                "STUDENT_OS_PDF_MANUAL_CHARS_PER_PAGE": "100",
            },
        )
    )["probes"][0]
    if manual_probe["tool"] != "pymupdf":
        raise AssertionError(f"Expected text-heavy PDF to prefer pymupdf, got: {manual_probe}")

    pymupdf_root = repo / "references" / "imports" / "pymupdf-output"
    pymupdf_payload = json.loads(
        run_script(
            "materials_convert.py",
            str(manual_pdf),
            "--output-root",
            str(pymupdf_root),
            "--force-strategy",
            "pymupdf",
            "--pages",
            "2",
            "--overwrite",
        )
    )
    if pymupdf_payload["converted"][0].get("import_method") != "pymupdf":
        raise AssertionError(f"Expected forced pymupdf conversion, got: {pymupdf_payload}")
    pymupdf_md = (pymupdf_root / "text-manual.pdf.md").read_text(encoding="utf-8")
    if "## Page 2" not in pymupdf_md or "## Page 1" in pymupdf_md or "## Page 3" in pymupdf_md:
        raise AssertionError(f"Expected pymupdf --pages 2 to extract only page 2, got:\n{pymupdf_md}")

    invalid_pages_result = subprocess.run(
        [
            sys.executable,
            "-B",
            str(STUDENT_OS_SCRIPTS / "materials_convert.py"),
            str(manual_pdf),
            "--output-root",
            str(repo / "references" / "imports" / "pymupdf-invalid-pages"),
            "--force-strategy",
            "pymupdf",
            "--pages",
            "99",
            "--overwrite",
        ],
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=ROOT,
        env={**os.environ, "PYTHONDONTWRITEBYTECODE": "1", "PYTHONIOENCODING": "utf-8"},
    )
    if invalid_pages_result.returncode == 0:
        raise AssertionError("Expected out-of-range --pages to exit nonzero")
    invalid_pages_payload = json.loads(invalid_pages_result.stdout)
    if not invalid_pages_payload.get("errors") or "does not select any pages" not in invalid_pages_payload["errors"][
        0
    ].get("error", ""):
        raise AssertionError(f"Expected invalid --pages error, got: {invalid_pages_payload}")

    method_api_no_token = subprocess.run(
        [
            sys.executable,
            "-B",
            str(STUDENT_OS_SCRIPTS / "materials_convert.py"),
            str(manual_pdf),
            "--output-root",
            str(repo / "references" / "imports" / "method-api-no-token"),
            "--method",
            "api",
            "--overwrite",
        ],
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=empty_skill_root,
        env={
            **os.environ,
            "PYTHONDONTWRITEBYTECODE": "1",
            "PYTHONIOENCODING": "utf-8",
            **no_token_env,
        },
    )
    if method_api_no_token.returncode == 0:
        raise AssertionError("Expected --method api without token to exit nonzero")
    method_api_payload = json.loads(method_api_no_token.stdout)
    if method_api_payload.get("converted"):
        raise AssertionError(f"--method api without token must not convert locally, got: {method_api_payload}")
    if not method_api_payload.get("errors") or "requires a token" not in method_api_payload["errors"][0].get(
        "error", ""
    ):
        raise AssertionError(f"Expected --method api without token to error, got: {method_api_payload}")

    forced_api_no_token = subprocess.run(
        [
            sys.executable,
            "-B",
            str(STUDENT_OS_SCRIPTS / "materials_convert.py"),
            str(manual_pdf),
            "--output-root",
            str(repo / "references" / "imports" / "forced-api-no-token"),
            "--force-strategy",
            "mineru-api",
            "--overwrite",
        ],
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=empty_skill_root,
        env={
            **os.environ,
            "PYTHONDONTWRITEBYTECODE": "1",
            "PYTHONIOENCODING": "utf-8",
            **no_token_env,
        },
    )
    if forced_api_no_token.returncode == 0:
        raise AssertionError("Expected forced mineru-api without token to exit nonzero")
    forced_api_payload = json.loads(forced_api_no_token.stdout)
    if not forced_api_payload.get("errors") or "requires a token" not in forced_api_payload["errors"][0].get("error", ""):
        raise AssertionError(f"Expected forced API without token to error, got: {forced_api_payload}")

    corrupt_docx = fixture_root / "corrupt.docx"
    corrupt_docx.write_bytes(b"PK\x03\x04not-a-real-docx")
    corrupt_probe_result = subprocess.run(
        [
            sys.executable,
            "-B",
            str(STUDENT_OS_SCRIPTS / "materials_convert.py"),
            str(corrupt_docx),
            "--probe-only",
        ],
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=ROOT,
        env={**os.environ, "PYTHONDONTWRITEBYTECODE": "1", "PYTHONIOENCODING": "utf-8"},
    )
    if corrupt_probe_result.returncode == 0:
        raise AssertionError("Expected corrupt DOCX --probe-only to exit nonzero")
    corrupt_probe_payload = json.loads(corrupt_probe_result.stdout)
    if not corrupt_probe_payload.get("errors") or "Failed to probe DOCX" not in corrupt_probe_payload["errors"][0].get(
        "error", ""
    ):
        raise AssertionError(f"Expected corrupt DOCX probe error, got: {corrupt_probe_payload}")
    corrupt_docx.unlink()

    image_heavy_docx = fixture_root / "image-heavy.docx"
    Document, _, _, _ = load_import_dependencies()
    from docx.shared import Inches
    import struct
    import zlib

    def _png_chunk(tag: bytes, data: bytes) -> bytes:
        return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)

    valid_png = fixture_root / "valid-1x1.png"
    ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    idat = zlib.compress(b"\x00\xff\x00\x00")
    valid_png.write_bytes(b"\x89PNG\r\n\x1a\n" + _png_chunk(b"IHDR", ihdr) + _png_chunk(b"IDAT", idat) + _png_chunk(b"IEND", b""))
    heavy_doc = Document()
    heavy_doc.add_paragraph("pic")
    heavy_doc.add_picture(str(valid_png), width=Inches(1.5))
    heavy_doc.save(str(image_heavy_docx))
    heavy_probe = json.loads(
        run_script(
            "materials_convert.py",
            str(image_heavy_docx),
            "--probe-only",
            "--api-token",
            "test-token",
        )
    )["probes"][0]
    if (
        heavy_probe["tool"] != "mineru-api"
        or heavy_probe["strategy"] != "image-heavy-docx"
        or not heavy_probe.get("needs_ocr")
    ):
        raise AssertionError(f"Expected image-heavy DOCX to prefer MinerU API+OCR, got: {heavy_probe}")
    valid_png.unlink()

    auto_image_payload = json.loads(
        run_script(
            "materials_convert.py",
            str(png_path),
            "--output-root",
            str(repo / "references" / "imports" / "auto-image-api"),
            "--method",
            "auto",
            "--api-token",
            "test-token",
            "--overwrite",
            env={"PYTHONPATH": str(fake_sdk_root)},
        )
    )
    if auto_image_payload["converted"][0].get("import_method") != "mineru-api:vlm":
        raise AssertionError(f"Expected auto image routing to MinerU API, got: {auto_image_payload}")
    if not auto_image_payload["converted"][0].get("ocr"):
        raise AssertionError("Expected auto image routing to enable OCR")

    legacy_doc = fixture_root / "legacy-notes.doc"
    legacy_doc.write_bytes(b"\xd0\xcf\x11\xe0legacy-doc")
    legacy_probe = json.loads(
        run_script(
            "materials_convert.py",
            str(legacy_doc),
            "--probe-only",
            "--api-token",
            "test-token",
        )
    )["probes"][0]
    if legacy_probe["tool"] != "mineru-api":
        raise AssertionError(f"Expected legacy .doc to prefer MinerU API, got: {legacy_probe}")
    legacy_no_token = json.loads(
        run_script(
            "materials_convert.py",
            str(legacy_doc),
            "--probe-only",
            cwd=empty_skill_root,
            env=no_token_env,
        )
    )["probes"][0]
    if legacy_no_token["tool"] != "mineru-api" or not legacy_no_token.get("agent_api"):
        raise AssertionError(f"Expected legacy .doc without token to use MinerU v1 Agent, got: {legacy_no_token}")

    with FakeMineruAgentServer() as agent_server:
        agent_env = {
            **no_token_env,
            "STUDENT_OS_MINERU_AGENT_BASE_URL": agent_server.base_url,
        }
        agent_output_root = repo / "references" / "imports" / "agent-output"
        agent_payload = json.loads(
            run_script(
                "materials_convert.py",
                str(scanned_pdf),
                "--output-root",
                str(agent_output_root),
                "--method",
                "auto",
                "--overwrite",
                cwd=empty_skill_root,
                env=agent_env,
            )
        )
        if agent_payload["applied_method"] != "api":
            raise AssertionError(f"Expected v1 Agent to count as API mode, got: {agent_payload}")
        agent_converted = agent_payload["converted"][0]
        if agent_converted.get("import_method") != "mineru-agent-v1":
            raise AssertionError(f"Expected --method auto without token to use v1 Agent, got: {agent_converted}")
        ensure_contains(agent_output_root / "scanned-blank.pdf.md", "Import method: mineru-agent-v1")
        ensure_contains(agent_output_root / "scanned-blank.pdf.md", "# Agent Parsed - scanned-blank.pdf")
        ensure_not_contains(agent_output_root / "scanned-blank.pdf.md", "- uploaded_bytes: 0")

        agent_split_pdf = fixture_root / "agent-split.pdf"
        write_multipage_pdf_fixture(agent_split_pdf, 21)
        agent_split_root = repo / "references" / "imports" / "agent-split-output"
        agent_split_payload = json.loads(
            run_script(
                "materials_convert.py",
                str(agent_split_pdf),
                "--output-root",
                str(agent_split_root),
                "--method",
                "auto",
                "--overwrite",
                cwd=empty_skill_root,
                env=agent_env,
            )
        )
        split_info = agent_split_payload["converted"][0].get("split")
        if not split_info or split_info.get("part_count") != 2 or split_info.get("chunk_size") != 20:
            raise AssertionError(f"Expected v1 Agent PDF to split at 20 pages, got: {agent_split_payload}")
        ensure_contains(agent_split_root / "agent-split.pdf.md", "<!-- MERGED from 2 parts: pages 1-20, pages 21-21 -->")

        agent_too_large = subprocess.run(
            [
                sys.executable,
                "-B",
                str(STUDENT_OS_SCRIPTS / "materials_convert.py"),
                str(scanned_pdf),
                "--output-root",
                str(repo / "references" / "imports" / "agent-too-large"),
                "--method",
                "auto",
                "--overwrite",
            ],
            check=False,
            capture_output=True,
            text=True,
            encoding="utf-8",
            cwd=empty_skill_root,
            env={
                **os.environ,
                "PYTHONDONTWRITEBYTECODE": "1",
                "PYTHONIOENCODING": "utf-8",
                **agent_env,
                "STUDENT_OS_MINERU_AGENT_MAX_FILE_BYTES": "1",
            },
        )
        if agent_too_large.returncode == 0:
            raise AssertionError("Expected v1 Agent over size limit without token to exit nonzero")
        agent_too_large_payload = json.loads(agent_too_large.stdout)
        if not agent_too_large_payload.get("errors") or "provide MINERU_TOKEN" not in agent_too_large_payload["errors"][
            0
        ].get("error", ""):
            raise AssertionError(f"Expected v1 Agent size-limit token hint, got: {agent_too_large_payload}")
        agent_split_pdf.unlink()

    scanned_pdf.unlink()
    manual_pdf.unlink()
    image_heavy_docx.unlink()
    legacy_doc.unlink()

    docx_path.unlink()
    xlsx_path.unlink()
    pptx_path.unlink()
    pdf_path.unlink()
    large_pdf_path.unlink()
    png_path.unlink()
    binary_path.unlink()


def exercise_exam_census(repo: Path) -> None:
    exams_dir = repo / "courses" / "linear-algebra" / "references" / "exams"
    exams_dir.mkdir(parents=True, exist_ok=True)

    papers = {
        "2019-期中-A.pdf.md": "\n".join(
            [
                "---",
                "type: pdf-import-note",
                'course: "Linear Algebra"',
                "status: active",
                "tags: [import, pdf, exam]",
                "---",
                "",
                "# 2019 Midterm A",
                "",
                "1. Compute the rank of a matrix.",
                "2. Diagonalize a symmetric matrix via eigenvalues.",
                "",
            ]
        ),
        "2020-期中-B.pdf.md": "\n".join(
            [
                "---",
                "type: pdf-import-note",
                'course: "Linear Algebra"',
                "status: active",
                "tags: [import, pdf, exam]",
                "---",
                "",
                "# 2020 Midterm B",
                "",
                "1. Find the rank and nullity.",
                "2. Solve a linear system with Gaussian elimination.",
                "",
            ]
        ),
        "2021-期中-C.pdf.md": "\n".join(
            [
                "---",
                "type: pdf-import-note",
                'course: "Linear Algebra"',
                "status: active",
                "tags: [import, pdf, exam]",
                "---",
                "",
                "# 2021 Midterm C",
                "",
                "1. Eigenvalues and diagonalization.",
                "2. Rank of an augmented matrix.",
                "",
            ]
        ),
    }
    for name, body in papers.items():
        (exams_dir / name).write_text(body + "\n", encoding="utf-8", newline="\n")

    init_payload = json.loads(
        run_script(
            "init_exam_census.py",
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "期中",
            "--papers-dir",
            "courses/linear-algebra/references/exams",
            "--pattern",
            "**/*.pdf.md",
            "--batch-size",
            "2",
        )
    )
    if init_payload.get("paper_count") != 3 or init_payload.get("batch_count") != 2:
        raise AssertionError(f"Unexpected exam-census init payload: {init_payload}")

    state_dir = repo / ".student-os" / "state" / "exam-census" / "linear-algebra" / "期中"
    taxonomy_path = state_dir / "taxonomy.yaml"
    taxonomy_path.write_text(
        "\n".join(
            [
                "version: 1",
                'course: "Linear Algebra"',
                "exam_scope: 期中",
                "types:",
                "  - id: matrix-rank",
                "    name: 矩阵的秩",
                "    aliases: [秩, rank]",
                "    keywords: [秩, 行阶梯]",
                "  - id: eigen-decomp",
                "    name: 特征值与对角化",
                "    aliases: [特征值, diagonalize]",
                "    keywords: [特征值, 对角化]",
                "  - id: gaussian-elim",
                "    name: 高斯消元",
                "    aliases: [消元]",
                "    keywords: [高斯, 行变换]",
                "",
            ]
        ),
        encoding="utf-8",
        newline="\n",
    )

    annotations = {
        "2019-期中-A": {
            "source": "courses/linear-algebra/references/exams/2019-期中-A.pdf.md",
            "exam_label": "2019 期中 A",
            "types_present": ["matrix-rank", "eigen-decomp"],
            "type_counts": {"matrix-rank": 1, "eigen-decomp": 1},
            "confidence": "high",
            "notes": "classic pair",
        },
        "2020-期中-B": {
            "source": "courses/linear-algebra/references/exams/2020-期中-B.pdf.md",
            "exam_label": "2020 期中 B",
            "types_present": ["matrix-rank", "gaussian-elim"],
            "type_counts": {"matrix-rank": 1, "gaussian-elim": 1},
            "confidence": "high",
            "notes": "",
        },
        "2021-期中-C": {
            "source": "courses/linear-algebra/references/exams/2021-期中-C.pdf.md",
            "exam_label": "2021 期中 C",
            "types_present": ["matrix-rank", "eigen-decomp"],
            "type_counts": {"matrix-rank": 1, "eigen-decomp": 1},
            "confidence": "low",
            "notes": "OCR a bit noisy",
        },
    }
    annotations_dir = state_dir / "annotations"
    for stem, payload in annotations.items():
        (annotations_dir / f"{stem}.json").write_text(
            json.dumps(payload, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
            newline="\n",
        )

    stats_payload = json.loads(
        run_script(
            "build_exam_type_stats.py",
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "期中",
            "--overwrite",
        )
    )
    ranked_ids = [item["id"] for item in stats_payload.get("ranked_types", [])]
    if ranked_ids[:2] != ["matrix-rank", "eigen-decomp"]:
        raise AssertionError(f"Expected matrix-rank then eigen-decomp by appearance, got: {ranked_ids}")
    if stats_payload["ranked_types"][0]["paper_count"] != 3:
        raise AssertionError(f"matrix-rank should appear in all 3 papers: {stats_payload}")
    if not stats_payload["ranked_types"][0]["must_know"]:
        raise AssertionError("matrix-rank should be must-know at 100% appearance")

    report_path = repo / "courses" / "linear-algebra" / "reviews" / "期中" / "题型频率统计.md"
    ensure_exists(report_path)
    ensure_contains(report_path, "type: exam-type-frequency-report")
    ensure_contains(report_path, "matrix-rank")
    ensure_contains(report_path, "Low-confidence annotations")
    ensure_contains(report_path, "`2021-期中-C`")
    ensure_contains(report_path, "](../../references/exams/2019-期中-A.pdf.md)")

    first_skeleton = repo / "courses" / "linear-algebra" / "reviews" / "期中" / "题型解析" / "01-matrix-rank.md"
    ensure_exists(first_skeleton)
    ensure_contains(first_skeleton, "type: exam-type-analysis")
    ensure_contains(first_skeleton, "exam_type_id: \"matrix-rank\"")
    ensure_contains(first_skeleton, "exam_type_name:")
    ensure_contains(first_skeleton, "quality: draft")
    ensure_contains(first_skeleton, "必掌握：是")
    ensure_contains(first_skeleton, "](../../../references/exams/2019-期中-A.pdf.md)")
    first_fm = first_skeleton.read_text(encoding="utf-8").split("---", 2)[1]
    if "source_artifacts: [" in first_fm:
        raise AssertionError("题型解析 frontmatter must not contain bulky source_artifacts arrays")
    if "generated_fingerprint" in first_fm:
        raise AssertionError("题型解析 frontmatter must not contain generated_fingerprint")
    for required_field in (
        "exam_type_name:",
        "rank:",
        "paper_count:",
        "must_know:",
        "quality:",
        "status:",
        "source_summary:",
    ):
        if required_field not in first_fm:
            raise AssertionError(f"Missing unified frontmatter field {required_field!r} in {first_skeleton}")
    ensure_contains(first_skeleton, "共 3 份试卷；详见 题型频率统计.md")

    second_skeleton = repo / "courses" / "linear-algebra" / "reviews" / "期中" / "题型解析" / "02-eigen-decomp.md"
    ensure_exists(second_skeleton)
    third_skeleton = repo / "courses" / "linear-algebra" / "reviews" / "期中" / "题型解析" / "03-gaussian-elim.md"
    ensure_exists(third_skeleton)

    # Pipe characters in type names must not break the Markdown table.
    taxonomy_path.write_text(
        taxonomy_path.read_text(encoding="utf-8").replace(
            "    name: 矩阵的秩\n",
            "    name: \"矩阵的秩 P(A|B)\"\n",
        ),
        encoding="utf-8",
        newline="\n",
    )
    run_script(
        "build_exam_type_stats.py",
        str(repo),
        "--course",
        "linear-algebra",
        "--exam-scope",
        "期中",
        "--overwrite",
    )
    ensure_contains(report_path, "矩阵的秩 P(A\\|B)")

    # Determinant / absolute-value pipes must stay in one Markdown table cell.
    utils_spec = importlib.util.spec_from_file_location(
        "exam_census_utils_issue51", STUDENT_OS_SCRIPTS / "exam_census_utils.py"
    )
    if utils_spec is None or utils_spec.loader is None:
        raise RuntimeError("Unable to load exam_census_utils for pipe-escape check")
    exam_census_utils = importlib.util.module_from_spec(utils_spec)
    sys.path.insert(0, str(STUDENT_OS_SCRIPTS))
    try:
        utils_spec.loader.exec_module(exam_census_utils)
    finally:
        if sys.path and sys.path[0] == str(STUDENT_OS_SCRIPTS):
            sys.path.pop(0)
    for raw, expected_fragment in (
        ("|A|", r"\lvert A\rvert"),
        ("$|A|$", r"\lvert A\rvert"),
        ("|λE-A|", r"\lvert λE-A\rvert"),
    ):
        escaped = exam_census_utils.md_table_cell(raw)
        if expected_fragment not in escaped:
            raise AssertionError(f"md_table_cell({raw!r}) missing {expected_fragment!r}: {escaped!r}")
        residual = escaped.replace("\\|", "").replace(r"\lvert", "").replace(r"\rvert", "")
        if "|" in residual:
            raise AssertionError(f"md_table_cell left bare pipe in {escaped!r}")
    sample_table = "\n".join(
        [
            "| 公式 | 说明 |",
            "| --- | --- |",
            f"| {exam_census_utils.md_table_cell('|A|')} | det |",
            f"| {exam_census_utils.md_table_cell('$|A|$')} | abs |",
            f"| {exam_census_utils.md_table_cell('|λE-A|')} | char |",
        ]
    )
    if exam_census_utils.markdown_table_pipe_issues(sample_table):
        raise AssertionError(f"escaped determinant table should be clean: {sample_table}")
    broken_table = "\n".join(
        [
            "| 公式 | 说明 |",
            "| --- | --- |",
            "| $|A|$ | det |",
        ]
    )
    if not exam_census_utils.markdown_table_pipe_issues(broken_table):
        raise AssertionError("expected bare $|A|$ table cell to be flagged as column mismatch")

    # Unknown types_present / type_counts keys must appear in the durable report and fail --validate.
    bad_annotation = dict(annotations["2021-期中-C"])
    bad_annotation["types_present"] = ["matrix-rank", "not-a-real-type"]
    bad_annotation["type_counts"] = {"matrix-rank": 1, "ghost-count": 2}
    (annotations_dir / "2021-期中-C.json").write_text(
        json.dumps(bad_annotation, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
        newline="\n",
    )
    bad_validate = subprocess.run(
        [
            sys.executable,
            "-B",
            str(STUDENT_OS_SCRIPTS / "build_exam_type_stats.py"),
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "期中",
            "--validate",
            "--overwrite",
        ],
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=ROOT,
        env={**os.environ, "PYTHONDONTWRITEBYTECODE": "1", "PYTHONIOENCODING": "utf-8"},
    )
    if bad_validate.returncode == 0:
        raise AssertionError("Expected --validate to fail on unknown type ids / count keys")
    ensure_contains(report_path, "Unknown type ids in types_present")
    ensure_contains(report_path, "not-a-real-type")
    ensure_contains(report_path, "ghost-count")
    bad_payload = json.loads(bad_validate.stdout)
    if not bad_payload.get("skeletons_skipped_due_to_validate"):
        raise AssertionError("Expected skeleton reconcile to be skipped when --validate fails")

    # Invalid type_counts must not silently become 1.
    invalid_annotation = dict(annotations["2021-期中-C"])
    invalid_annotation["types_present"] = ["matrix-rank", "eigen-decomp"]
    invalid_annotation["type_counts"] = {"matrix-rank": "many", "eigen-decomp": -2}
    (annotations_dir / "2021-期中-C.json").write_text(
        json.dumps(invalid_annotation, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
        newline="\n",
    )
    skeleton_before = first_skeleton.read_text(encoding="utf-8")
    invalid_validate = subprocess.run(
        [
            sys.executable,
            "-B",
            str(STUDENT_OS_SCRIPTS / "build_exam_type_stats.py"),
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "期中",
            "--validate",
            "--overwrite",
        ],
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=ROOT,
        env={**os.environ, "PYTHONDONTWRITEBYTECODE": "1", "PYTHONIOENCODING": "utf-8"},
    )
    if invalid_validate.returncode == 0:
        raise AssertionError("Expected --validate to fail on invalid type_counts values")
    ensure_contains(report_path, "Invalid type_counts values")
    if first_skeleton.read_text(encoding="utf-8") != skeleton_before:
        raise AssertionError("Skeleton must not change when --validate fails")

    # Restore clean annotations after validation failure cases.
    for stem, payload in annotations.items():
        (annotations_dir / f"{stem}.json").write_text(
            json.dumps(payload, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
            newline="\n",
        )

    # Nested papers with the same basename must get distinct annotation ids.
    nested_root = exams_dir / "nested"
    (nested_root / "2019").mkdir(parents=True, exist_ok=True)
    (nested_root / "2020").mkdir(parents=True, exist_ok=True)
    (nested_root / "2019" / "paper.pdf.md").write_text("# nested 2019\n", encoding="utf-8", newline="\n")
    (nested_root / "2020" / "paper.pdf.md").write_text("# nested 2020\n", encoding="utf-8", newline="\n")
    nested_init = json.loads(
        run_script(
            "init_exam_census.py",
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "nested-midterm",
            "--papers-dir",
            "courses/linear-algebra/references/exams/nested",
            "--pattern",
            "**/*.pdf.md",
            "--overwrite",
        )
    )
    nested_manifest = Path(nested_init["manifest"])
    nested_stems = {item["stem"] for item in json.loads(nested_manifest.read_text(encoding="utf-8"))["papers"]}
    if nested_stems != {"2019__paper", "2020__paper"}:
        raise AssertionError(f"Expected path-namespaced annotation ids, got: {nested_stems}")

    # Semester-aware courses must not share census state.
    run_script("scaffold_course.py", str(repo), "Calculus II", "--semester", "2026 Fall")
    run_script("scaffold_course.py", str(repo), "Calculus II", "--semester", "2027 Spring")
    fall_exams = repo / "courses" / "2026-fall" / "calculus-ii" / "references" / "exams"
    spring_exams = repo / "courses" / "2027-spring" / "calculus-ii" / "references" / "exams"
    fall_exams.mkdir(parents=True, exist_ok=True)
    spring_exams.mkdir(parents=True, exist_ok=True)
    (fall_exams / "paper.pdf.md").write_text("# fall\n", encoding="utf-8", newline="\n")
    (spring_exams / "paper.pdf.md").write_text("# spring\n", encoding="utf-8", newline="\n")
    fall_init = json.loads(
        run_script(
            "init_exam_census.py",
            str(repo),
            "--course",
            "calculus-ii",
            "--semester",
            "2026 Fall",
            "--exam-scope",
            "期中",
            "--papers-dir",
            "courses/2026-fall/calculus-ii/references/exams",
        )
    )
    spring_init = json.loads(
        run_script(
            "init_exam_census.py",
            str(repo),
            "--course",
            "calculus-ii",
            "--semester",
            "2027 Spring",
            "--exam-scope",
            "期中",
            "--papers-dir",
            "courses/2027-spring/calculus-ii/references/exams",
        )
    )
    if fall_init["course"] != "2026-fall/calculus-ii" or spring_init["course"] != "2027-spring/calculus-ii":
        raise AssertionError(f"Expected semester-qualified course keys, got {fall_init['course']} / {spring_init['course']}")
    if Path(fall_init["state_dir"]) == Path(spring_init["state_dir"]):
        raise AssertionError("Semester-qualified courses must not share exam-census state directories")

    # Rank changes should retire obsolete generated skeletons.
    taxonomy_path.write_text(
        "\n".join(
            [
                "version: 1",
                'course: "Linear Algebra"',
                "exam_scope: 期中",
                "types:",
                "  - id: matrix-rank",
                "    name: 矩阵的秩",
                "    aliases: [秩, rank]",
                "    keywords: [秩, 行阶梯]",
                "  - id: eigen-decomp",
                "    name: 特征值与对角化",
                "    aliases: [特征值, diagonalize]",
                "    keywords: [特征值, 对角化]",
                "  - id: gaussian-elim",
                "    name: 高斯消元",
                "    aliases: [消元]",
                "    keywords: [高斯, 行变换]",
                "",
            ]
        ),
        encoding="utf-8",
        newline="\n",
    )
    for stem, payload in annotations.items():
        (annotations_dir / f"{stem}.json").write_text(
            json.dumps(payload, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
            newline="\n",
        )
    # Boost gaussian-elim above eigen-decomp by swapping one paper's second type.
    boosted = dict(annotations["2021-期中-C"])
    boosted["types_present"] = ["matrix-rank", "gaussian-elim"]
    boosted["type_counts"] = {"matrix-rank": 1, "gaussian-elim": 1}
    (annotations_dir / "2021-期中-C.json").write_text(
        json.dumps(boosted, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
        newline="\n",
    )
    rerank = json.loads(
        run_script(
            "build_exam_type_stats.py",
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "期中",
            "--overwrite",
        )
    )
    if [item["id"] for item in rerank["ranked_types"][:3]] != ["matrix-rank", "gaussian-elim", "eigen-decomp"]:
        raise AssertionError(f"Unexpected reranked order: {rerank['ranked_types']}")
    analysis_dir = repo / "courses" / "linear-algebra" / "reviews" / "期中" / "题型解析"
    ensure_exists(analysis_dir / "02-gaussian-elim.md")
    ensure_exists(analysis_dir / "03-eigen-decomp.md")
    if (analysis_dir / "02-eigen-decomp.md").exists():
        raise AssertionError("Obsolete ranked skeleton 02-eigen-decomp.md should be retired after rerank")
    ensure_contains(analysis_dir / "02-gaussian-elim.md", 'exam_type_id: "gaussian-elim"')
    ensure_contains(analysis_dir / "03-eigen-decomp.md", 'exam_type_id: "eigen-decomp"')

    # User-edited page must follow its type_id across rank swaps (no content cross-wiring).
    for stem, payload in annotations.items():
        (annotations_dir / f"{stem}.json").write_text(
            json.dumps(payload, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
            newline="\n",
        )
    run_script(
        "build_exam_type_stats.py",
        str(repo),
        "--course",
        "linear-algebra",
        "--exam-scope",
        "期中",
        "--overwrite",
    )
    eigen_at_02 = analysis_dir / "02-eigen-decomp.md"
    ensure_exists(eigen_at_02)
    eigen_at_02.write_text(
        eigen_at_02.read_text(encoding="utf-8") + "\n## Unique Eigen Marker\n\nSWAP-TEST-EIGEN\n",
        encoding="utf-8",
        newline="\n",
    )
    boosted = dict(annotations["2021-期中-C"])
    boosted["types_present"] = ["matrix-rank", "gaussian-elim"]
    boosted["type_counts"] = {"matrix-rank": 1, "gaussian-elim": 1}
    (annotations_dir / "2021-期中-C.json").write_text(
        json.dumps(boosted, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
        newline="\n",
    )
    swap = json.loads(
        run_script(
            "build_exam_type_stats.py",
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "期中",
            "--overwrite",
        )
    )
    if [item["id"] for item in swap["ranked_types"][:3]] != ["matrix-rank", "gaussian-elim", "eigen-decomp"]:
        raise AssertionError(f"Unexpected swap order: {swap['ranked_types']}")
    if (analysis_dir / "02-eigen-decomp.md").exists():
        raise AssertionError("Old 02-eigen-decomp.md must disappear after rank swap (moved, not copied)")
    gaussian_page = (analysis_dir / "02-gaussian-elim.md").read_text(encoding="utf-8")
    eigen_page = (analysis_dir / "03-eigen-decomp.md").read_text(encoding="utf-8")
    if 'exam_type_id: "gaussian-elim"' not in gaussian_page or "SWAP-TEST-EIGEN" in gaussian_page:
        raise AssertionError("Rank swap cross-wired eigen user content into gaussian page")
    if 'exam_type_id: "eigen-decomp"' not in eigen_page or "SWAP-TEST-EIGEN" not in eigen_page:
        raise AssertionError("Rank swap lost user-edited eigen content or wrong exam_type_id")
    first_h1 = next((line for line in eigen_page.splitlines() if line.startswith("# ")), "")
    if not first_h1.startswith("# 03 ·"):
        raise AssertionError(f"Migrated eigen page first H1 rank was not refreshed: {first_h1!r}")

    # Interrupted _reconcile_staging must be recovered, not wiped.
    for stem, payload in annotations.items():
        (annotations_dir / f"{stem}.json").write_text(
            json.dumps(payload, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
            newline="\n",
        )
    run_script(
        "build_exam_type_stats.py",
        str(repo),
        "--course",
        "linear-algebra",
        "--exam-scope",
        "期中",
        "--overwrite",
    )
    staging_dir = analysis_dir / "_reconcile_staging"
    staging_dir.mkdir(parents=True, exist_ok=True)
    staged_user = "\n".join(
        [
            "---",
            "type: exam-type-analysis",
            'course: "Linear Algebra"',
            'exam_scope: "期中"',
            'exam_type_id: "matrix-rank"',
            'exam_type_name: "矩阵的秩"',
            "rank: 1",
            "paper_count: 3",
            "must_know: true",
            "quality: draft",
            "status: active",
            'source_summary: "共 3 份试卷；详见 题型频率统计.md"',
            "---",
            "",
            "# 01 · 矩阵的秩",
            "",
            "## Recovered Staging Marker",
            "",
            "STAGING-RECOVERY-MARKER",
            "",
        ]
    )
    (staging_dir / "interrupted-user.md").write_text(staged_user + "\n", encoding="utf-8", newline="\n")
    run_script(
        "build_exam_type_stats.py",
        str(repo),
        "--course",
        "linear-algebra",
        "--exam-scope",
        "期中",
        "--overwrite",
    )
    recovered = (analysis_dir / "01-matrix-rank.md").read_text(encoding="utf-8")
    if "STAGING-RECOVERY-MARKER" not in recovered:
        raise AssertionError("Interrupted staging user content was not recovered on retry")
    if staging_dir.exists() and any(staging_dir.iterdir()):
        raise AssertionError("Staging directory should be empty after successful overwrite")

    # --validate should fail when an annotation is missing.
    (annotations_dir / "2021-期中-C.json").unlink()
    validate_result = subprocess.run(
        [
            sys.executable,
            "-B",
            str(STUDENT_OS_SCRIPTS / "build_exam_type_stats.py"),
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "期中",
            "--validate",
            "--overwrite",
        ],
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=ROOT,
        env={**os.environ, "PYTHONDONTWRITEBYTECODE": "1", "PYTHONIOENCODING": "utf-8"},
    )
    if validate_result.returncode == 0:
        raise AssertionError("build_exam_type_stats.py --validate should fail when annotations are missing")
    # Restore annotation and rebuild final report for the example snapshot.
    for stem, payload in annotations.items():
        (annotations_dir / f"{stem}.json").write_text(
            json.dumps(payload, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
            newline="\n",
        )
    run_script(
        "build_exam_type_stats.py",
        str(repo),
        "--course",
        "linear-algebra",
        "--exam-scope",
        "期中",
        "--overwrite",
    )

    # Path-like exam_scope values must be rejected.
    for bad_scope in ("../escape", "", "C:\\windows"):
        bad_init = subprocess.run(
            [
                sys.executable,
                "-B",
                str(STUDENT_OS_SCRIPTS / "init_exam_census.py"),
                str(repo),
                "--course",
                "linear-algebra",
                "--exam-scope",
                bad_scope,
                "--papers-dir",
                "courses/linear-algebra/references/exams",
            ],
            check=False,
            capture_output=True,
            text=True,
            encoding="utf-8",
            cwd=ROOT,
            env={**os.environ, "PYTHONDONTWRITEBYTECODE": "1", "PYTHONIOENCODING": "utf-8"},
        )
        if bad_init.returncode == 0:
            raise AssertionError(f"Expected init to reject exam_scope={bad_scope!r}")

    # Without --overwrite, existing skeletons must be skipped (not migrated/deleted).
    analysis_dir = repo / "courses" / "linear-algebra" / "reviews" / "期中" / "题型解析"
    existing_skeletons = sorted(path.name for path in analysis_dir.glob("*.md"))
    if not existing_skeletons:
        raise AssertionError("Expected skeletons before no-overwrite rebuild")
    existing_skeleton_text = {
        name: (analysis_dir / name).read_text(encoding="utf-8") for name in existing_skeletons
    }
    archive_dir = analysis_dir / "_archive"
    archive_before = sorted(
        path.relative_to(archive_dir).as_posix() for path in archive_dir.rglob("*") if path.is_file()
    ) if archive_dir.exists() else []
    no_overwrite = json.loads(
        run_script(
            "build_exam_type_stats.py",
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "期中",
        )
    )
    if not no_overwrite.get("skeletons_skipped"):
        raise AssertionError("Expected existing skeletons to be skipped without --overwrite")
    if no_overwrite.get("skeletons_migrated"):
        raise AssertionError(f"No-overwrite rebuild must not report migrated skeletons: {no_overwrite}")
    after_skeletons = sorted(path.name for path in analysis_dir.glob("*.md"))
    if after_skeletons != existing_skeletons:
        raise AssertionError(f"Skeletons changed without --overwrite: {existing_skeletons} -> {after_skeletons}")
    for name, before_text in existing_skeleton_text.items():
        after_text = (analysis_dir / name).read_text(encoding="utf-8")
        if after_text != before_text:
            raise AssertionError(f"Skeleton content changed without --overwrite: {name}")
    archive_after = sorted(
        path.relative_to(archive_dir).as_posix() for path in archive_dir.rglob("*") if path.is_file()
    ) if archive_dir.exists() else []
    if archive_after != archive_before:
        raise AssertionError(f"No-overwrite rebuild should not add archive entries: {archive_before} -> {archive_after}")

    # User-edited skeletons (fingerprint mismatch) are archived, not deleted.
    user_skeleton = analysis_dir / existing_skeletons[-1]
    user_text = user_skeleton.read_text(encoding="utf-8") + "\n## User Notes\n\nKeep me.\n"
    user_skeleton.write_text(user_text, encoding="utf-8", newline="\n")
    user_type = None
    import re as _re

    match = _re.search(r'(?m)^exam_type_id:\s*"?([^"\n]+)"?\s*$', user_text)
    if match:
        user_type = match.group(1)
    if not user_type:
        raise AssertionError(f"Could not read exam_type_id from {user_skeleton}")
    # Drop that type from all annotations so it becomes obsolete and should archive.
    for stem, payload in annotations.items():
        cleaned = dict(payload)
        cleaned["types_present"] = [tid for tid in payload["types_present"] if tid != user_type]
        cleaned["type_counts"] = {
            key: value for key, value in payload["type_counts"].items() if key != user_type
        }
        if "matrix-rank" not in cleaned["types_present"]:
            cleaned["types_present"].append("matrix-rank")
            cleaned["type_counts"]["matrix-rank"] = 1
        (annotations_dir / f"{stem}.json").write_text(
            json.dumps(cleaned, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
            newline="\n",
        )
    run_script(
        "build_exam_type_stats.py",
        str(repo),
        "--course",
        "linear-algebra",
        "--exam-scope",
        "期中",
        "--overwrite",
    )
    if user_skeleton.exists():
        raise AssertionError("Edited obsolete skeleton should leave 题型解析/")
    archived = analysis_dir / "_archive" / user_skeleton.name
    ensure_exists(archived)
    ensure_contains(archived, "Keep me.")

    # Taxonomy JSON round-trip preserves string ids like "true" and comma-bearing names.
    utils_spec = importlib.util.spec_from_file_location(
        "exam_census_utils_smoke", STUDENT_OS_SCRIPTS / "exam_census_utils.py"
    )
    if utils_spec is None or utils_spec.loader is None:
        raise RuntimeError("Unable to load exam_census_utils for round-trip check")
    exam_census_utils = importlib.util.module_from_spec(utils_spec)
    sys.path.insert(0, str(STUDENT_OS_SCRIPTS))
    try:
        utils_spec.loader.exec_module(exam_census_utils)
    finally:
        if sys.path and sys.path[0] == str(STUDENT_OS_SCRIPTS):
            sys.path.pop(0)
    roundtrip_path = state_dir / "taxonomy-roundtrip.yaml"
    exam_census_utils.write_taxonomy(
        roundtrip_path,
        {
            "version": 1,
            "course": "Linear Algebra",
            "exam_scope": "期中",
            "types": [
                {"id": "true", "name": "A, B family", "aliases": ["x,y"], "keywords": ["a,b"]},
            ],
        },
    )
    loaded = exam_census_utils.load_taxonomy_yaml(roundtrip_path)
    if loaded["types"][0]["id"] != "true":
        raise AssertionError(f"Expected id 'true' string round-trip, got {loaded['types'][0]['id']!r}")
    if loaded["types"][0]["name"] != "A, B family":
        raise AssertionError(f"Expected comma-bearing name preserved, got {loaded['types'][0]['name']!r}")
    dumped = exam_census_utils.dump_taxonomy_yaml(loaded)
    if '"true"' not in dumped:
        raise AssertionError("Writer should JSON-quote the id true")

    # Paths with spaces must produce quoted Markdown hrefs.
    spaced = exams_dir / "mid term paper.pdf.md"
    spaced.write_text("# spaced\n", encoding="utf-8", newline="\n")
    spaced_init = json.loads(
        run_script(
            "init_exam_census.py",
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "spaced-scope",
            "--papers-dir",
            "courses/linear-algebra/references/exams",
            "--pattern",
            "mid term paper.pdf.md",
            "--overwrite",
        )
    )
    spaced_state = Path(spaced_init["state_dir"])
    (spaced_state / "taxonomy.yaml").write_text(
        "\n".join(
            [
                "version: 1",
                'course: "Linear Algebra"',
                'exam_scope: "spaced-scope"',
                "types:",
                '  - id: "matrix-rank"',
                '    name: "矩阵的秩"',
                "    aliases: []",
                "    keywords: []",
                "",
            ]
        ),
        encoding="utf-8",
        newline="\n",
    )
    stem = "mid term paper"
    # annotation id from basename without .pdf.md
    (spaced_state / "annotations" / f"{stem}.json").write_text(
        json.dumps(
            {
                "source": "courses/linear-algebra/references/exams/mid term paper.pdf.md",
                "exam_label": "spaced",
                "types_present": ["matrix-rank"],
                "type_counts": {"matrix-rank": 1},
                "confidence": "high",
            },
            ensure_ascii=False,
            indent=2,
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    # Fix annotation filename: annotation_id uses relative path
    spaced_manifest = json.loads((spaced_state / "manifest.json").read_text(encoding="utf-8"))
    spaced_stem = spaced_manifest["papers"][0]["stem"]
    (spaced_state / "annotations" / f"{stem}.json").rename(spaced_state / "annotations" / f"{spaced_stem}.json")
    run_script(
        "build_exam_type_stats.py",
        str(repo),
        "--course",
        "linear-algebra",
        "--exam-scope",
        "spaced-scope",
        "--overwrite",
    )
    spaced_report = repo / "courses" / "linear-algebra" / "reviews" / "spaced-scope" / "题型频率统计.md"
    ensure_contains(spaced_report, "mid%20term%20paper.pdf.md")

    # Restore linear-algebra midterm annotations for any later consumers of this fixture.
    for stem, payload in annotations.items():
        (annotations_dir / f"{stem}.json").write_text(
            json.dumps(payload, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
            newline="\n",
        )

    # Phase A–E mechanical scripts (fill queue, quality gate, multi-dim, deep-dive, cross-val).
    run_script(
        "build_exam_type_stats.py",
        str(repo),
        "--course",
        "linear-algebra",
        "--exam-scope",
        "期中",
        "--validate",
        "--overwrite",
    )
    # Textbook sidecar under course references/ (not exams/) should appear in concept_sources.
    textbook_sidecar = (
        repo / "courses" / "linear-algebra" / "references" / "线性代数教材.pdf.md"
    )
    textbook_sidecar.parent.mkdir(parents=True, exist_ok=True)
    textbook_sidecar.write_text(
        "\n".join(
            [
                "---",
                'course: "Linear Algebra"',
                "---",
                "",
                "# 线性代数教材",
                "",
                "## 第1章 行列式与秩",
                "",
                "秩的定义……",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    fill_payload = json.loads(
        run_script(
            "fill_type_analysis.py",
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "期中",
        )
    )
    if fill_payload.get("item_count", 0) < 1:
        raise AssertionError(f"Expected fill queue items, got: {fill_payload}")
    if fill_payload.get("concept_source_count", 0) < 1:
        raise AssertionError(f"Expected concept_source_count >= 1, got: {fill_payload}")
    ensure_exists(state_dir / "fill-queue.json")
    fill_queue = json.loads((state_dir / "fill-queue.json").read_text(encoding="utf-8"))
    if fill_queue.get("quality_reference") != "references/exam-census-quality.md":
        raise AssertionError(f"Expected skill-relative quality_reference, got: {fill_queue.get('quality_reference')}")
    concept_paths = [str(item.get("path") or "") for item in (fill_queue.get("concept_sources") or [])]
    if not any("线性代数教材.pdf.md" in path.replace("\\", "/") for path in concept_paths):
        raise AssertionError(f"Expected textbook in concept_sources, got: {concept_paths}")
    if not fill_queue["items"][0].get("concept_sources"):
        raise AssertionError("Expected per-item concept_sources on fill-queue items")
    if not fill_queue["items"][0].get("source_papers"):
        raise AssertionError("Expected source_papers on fill-queue items")
    if not fill_queue["items"][0].get("source_instances"):
        raise AssertionError("Expected source_instances on fill-queue items")
    fill_instructions = "\n".join(fill_queue["items"][0].get("instructions") or [])
    for needle in (
        "禁止自行编造",
        "至少 5 道例题",
        "4 道自测题",
        "必须标注来源",
        "不要在 <details> 中使用 $$",
        "不要使用引用块内表格",
        "content-standard v3",
        "表格为主",
        "难度星级",
        "自测答案",
        "快速得分技巧",
        "[表达式]",
        "concept_sources",
        "未参考指定教材",
        "参考：",
    ):
        if needle not in fill_instructions:
            raise AssertionError(f"Expected fill instructions to include {needle!r}, got: {fill_instructions}")

    # Without textbook candidates, concept_sources is empty but disclaimer guidance remains.
    textbook_sidecar.unlink()
    empty_book_payload = json.loads(
        run_script(
            "fill_type_analysis.py",
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "期中",
        )
    )
    empty_book_file = json.loads((state_dir / "fill-queue.json").read_text(encoding="utf-8"))
    if empty_book_payload.get("concept_source_count", 0) != 0:
        raise AssertionError(
            f"Expected concept_source_count 0 after removing textbook, got: {empty_book_payload}"
        )
    if empty_book_file.get("concept_sources"):
        raise AssertionError(
            f"Expected empty concept_sources after removing textbook, got: {empty_book_file.get('concept_sources')}"
        )
    if "未参考指定教材" not in str(empty_book_file.get("concept_sources_note") or ""):
        raise AssertionError(
            f"Expected disclaimer note when no textbooks, got: {empty_book_file.get('concept_sources_note')}"
        )
    # Restore textbook for later consumers of this fixture.
    textbook_sidecar.write_text(
        "\n".join(
            [
                "---",
                'course: "Linear Algebra"',
                "---",
                "",
                "# 线性代数教材",
                "",
                "## 第1章 行列式与秩",
                "",
                "秩的定义……",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    # Exam-like / repair artifacts / unrelated textbooks must not enter concept_sources.
    (repo / "courses" / "linear-algebra" / "references" / "linear-algebra-2019-midterm.pdf.md").write_text(
        "# exam\n", encoding="utf-8", newline="\n"
    )
    (repo / "courses" / "linear-algebra" / "references" / "线性代数教材.pdf.raw.md").write_text(
        "# raw\n", encoding="utf-8", newline="\n"
    )
    vault_books = repo / "references" / "textbooks"
    vault_books.mkdir(parents=True, exist_ok=True)
    (vault_books / "calculus-textbook.pdf.md").write_text("# calc\n", encoding="utf-8", newline="\n")
    (vault_books / "linear-algebra-notes.pdf.md").write_text("# la\n", encoding="utf-8", newline="\n")
    run_script(
        "fill_type_analysis.py",
        str(repo),
        "--course",
        "linear-algebra",
        "--exam-scope",
        "期中",
    )
    filtered_queue = json.loads((state_dir / "fill-queue.json").read_text(encoding="utf-8"))
    filtered_paths = [str(item.get("path") or "").replace("\\", "/") for item in (filtered_queue.get("concept_sources") or [])]
    if any("midterm" in path for path in filtered_paths):
        raise AssertionError(f"Exam-like filename must not enter concept_sources: {filtered_paths}")
    if any(".raw.md" in path for path in filtered_paths):
        raise AssertionError(f"Repair raw sidecar must not enter concept_sources: {filtered_paths}")
    if any("calculus-textbook" in path for path in filtered_paths):
        raise AssertionError(f"Unrelated vault textbook must not enter concept_sources: {filtered_paths}")
    if not any("线性代数教材.pdf.md" in path for path in filtered_paths):
        raise AssertionError(f"Expected course textbook retained in concept_sources: {filtered_paths}")
    if not any("linear-algebra-notes.pdf.md" in path for path in filtered_paths):
        raise AssertionError(f"Expected course-matching vault textbook in concept_sources: {filtered_paths}")

    # Issue #69: course-local lecture-material directories feed concept_sources.
    lecture_material_dir = repo / "courses" / "linear-algebra" / "教材课件"
    lecture_material_dir.mkdir(parents=True, exist_ok=True)
    (lecture_material_dir / "1-2 n阶行列式.pdf.md").write_text(
        "\n".join(
            [
                "---",
                'course: "Linear Algebra"',
                "---",
                "",
                "# 1-2 n阶行列式",
                "",
                "## 行列式的定义",
                "",
                "行列式是一个数……",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    # Exam-like / answer files inside lecture-material dirs must still be excluded.
    (lecture_material_dir / "期中试卷.pdf.md").write_text("# 期中试卷\n", encoding="utf-8", newline="\n")
    (lecture_material_dir / "答案.pdf.md").write_text("# 答案\n", encoding="utf-8", newline="\n")
    # Files under reviews/<scope>/文本/ must not leak into concept_sources.
    review_text_dir = repo / "courses" / "linear-algebra" / "reviews" / "期中" / "文本"
    review_text_dir.mkdir(parents=True, exist_ok=True)
    (review_text_dir / "review-handout.pdf.md").write_text("# review\n", encoding="utf-8", newline="\n")
    run_script(
        "fill_type_analysis.py",
        str(repo),
        "--course",
        "linear-algebra",
        "--exam-scope",
        "期中",
    )
    lecture_queue = json.loads((state_dir / "fill-queue.json").read_text(encoding="utf-8"))
    lecture_paths = [str(item.get("path") or "").replace("\\", "/") for item in (lecture_queue.get("concept_sources") or [])]
    if not any("教材课件/1-2 n阶行列式.pdf.md" in path for path in lecture_paths):
        raise AssertionError(f"Expected lecture-material sidecar in concept_sources: {lecture_paths}")
    if any("期中试卷.pdf.md" in path for path in lecture_paths):
        raise AssertionError(f"Exam-paper sidecar must not enter concept_sources: {lecture_paths}")
    if any("答案.pdf.md" in path for path in lecture_paths):
        raise AssertionError(f"Answer sidecar must not enter concept_sources: {lecture_paths}")
    if any("review-handout.pdf.md" in path for path in lecture_paths):
        raise AssertionError(f"Review-text sidecar must not enter concept_sources: {lecture_paths}")

    # Without any textbook candidates, concept_sources is empty and the disclaimer remains.
    saved_course_textbook = repo / "courses" / "linear-algebra" / "references" / "线性代数教材.pdf.md"
    saved_course_textbook_text = (
        saved_course_textbook.read_text(encoding="utf-8") if saved_course_textbook.exists() else ""
    )
    saved_vault_textbook = repo / "references" / "textbooks" / "linear-algebra-notes.pdf.md"
    saved_vault_textbook_text = (
        saved_vault_textbook.read_text(encoding="utf-8") if saved_vault_textbook.exists() else ""
    )
    shutil.rmtree(lecture_material_dir)
    if saved_course_textbook.exists():
        saved_course_textbook.unlink()
    if saved_vault_textbook.exists():
        saved_vault_textbook.unlink()
    empty_lecture_payload = json.loads(
        run_script(
            "fill_type_analysis.py",
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "期中",
        )
    )
    empty_lecture_file = json.loads((state_dir / "fill-queue.json").read_text(encoding="utf-8"))
    if empty_lecture_payload.get("concept_source_count", 0) != 0:
        raise AssertionError(
            f"Expected concept_source_count 0 after removing textbooks, got: {empty_lecture_payload}"
        )
    if empty_lecture_file.get("concept_sources"):
        raise AssertionError(
            f"Expected empty concept_sources after removing textbooks, got: {empty_lecture_file.get('concept_sources')}"
        )
    if "未参考指定教材" not in str(empty_lecture_file.get("concept_sources_note") or ""):
        raise AssertionError(
            f"Expected disclaimer note when no concept sources, got: {empty_lecture_file.get('concept_sources_note')}"
        )
    # Restore for any later consumers of this fixture.
    lecture_material_dir.mkdir(parents=True, exist_ok=True)
    (lecture_material_dir / "1-2 n阶行列式.pdf.md").write_text(
        "\n".join(
            [
                "---",
                'course: "Linear Algebra"',
                "---",
                "",
                "# 1-2 n阶行列式",
                "",
                "## 行列式的定义",
                "",
                "行列式是一个数……",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    if saved_course_textbook_text:
        saved_course_textbook.write_text(saved_course_textbook_text, encoding="utf-8", newline="\n")
    if saved_vault_textbook_text:
        saved_vault_textbook.write_text(saved_vault_textbook_text, encoding="utf-8", newline="\n")

    review_result = subprocess.run(
        [
            sys.executable,
            "-B",
            str(STUDENT_OS_SCRIPTS / "review_type_analysis.py"),
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "期中",
        ],
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=ROOT,
        env={**os.environ, "PYTHONDONTWRITEBYTECODE": "1", "PYTHONIOENCODING": "utf-8"},
    )
    if review_result.returncode != 1:
        raise AssertionError(f"Expected Phase B exit code 1 for bare skeletons, got {review_result.returncode}")
    ensure_exists(state_dir / "quality-reviews.json")
    quality_report = json.loads((state_dir / "quality-reviews.json").read_text(encoding="utf-8"))
    if quality_report.get("needs_revision_count", 0) <= 0:
        raise AssertionError(f"Expected needs_revision_count > 0, got: {quality_report}")
    ensure_exists(repo / "courses" / "linear-algebra" / "reviews" / "期中" / "analysis" / "质量门禁.md")

    multi = json.loads(
        run_script(
            "build_multi_dim_stats.py",
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "期中",
            "--overwrite",
        )
    )
    if multi.get("pair_count", 0) < 1:
        raise AssertionError(f"Expected co-occurrence pairs from annotations, got: {multi}")
    ensure_exists(repo / "courses" / "linear-algebra" / "reviews" / "期中" / "analysis" / "题型关联分析.md")
    if multi.get("format_labels", {}).get("unspecified", 0) < 1:
        raise AssertionError(f"Expected unspecified format labels when annotations omit format: {multi}")
    reliability_report = repo / "courses" / "linear-algebra" / "reviews" / "期中" / "analysis" / "卷源可靠性分级.md"
    ensure_exists(reliability_report)
    reliability_text = reliability_report.read_text(encoding="utf-8")
    for banned in ("Seeded from", "Paper | Reliability", "unspecified"):
        if banned in reliability_text:
            raise AssertionError(f"analysis/卷源可靠性分级.md must not contain {banned!r}")
    ensure_contains(reliability_report, "| 试卷 | 可靠性 |")
    ensure_contains(reliability_report, "未标注")

    # Phase B CLI must catch Issue #51 analysis defects and pass after fix.
    analysis_dir_reports = repo / "courses" / "linear-algebra" / "reviews" / "期中" / "analysis"
    bad_report_path = analysis_dir_reports / "题型关联分析.md"
    good_report_backup = bad_report_path.read_text(encoding="utf-8")
    bad_report_path.write_text(
        "\n".join(
            [
                "---",
                "type: exam-type-cooccurrence",
                'course: "Linear Algebra"',
                "status: draft",
                "---",
                "",
                "# bad seed",
                "",
                "Seeded from optional annotation.",
                "",
                "| Paper | Reliability |",
                "| --- | --- |",
                "| demo | unspecified |",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    bad_cli = subprocess.run(
        [
            sys.executable,
            "-B",
            str(STUDENT_OS_SCRIPTS / "review_type_analysis.py"),
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "期中",
        ],
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=ROOT,
        env={**os.environ, "PYTHONDONTWRITEBYTECODE": "1", "PYTHONIOENCODING": "utf-8"},
    )
    if bad_cli.returncode == 0:
        raise AssertionError("Expected Phase B CLI to fail when analysis report has English residue")
    bad_quality = json.loads((state_dir / "quality-reviews.json").read_text(encoding="utf-8"))
    bad_analysis_hits = [
        item
        for item in bad_quality.get("analysis_needs_revision", [])
        if item.get("path", "").endswith("题型关联分析.md")
    ]
    if not bad_analysis_hits:
        raise AssertionError(f"Expected analysis_needs_revision for 题型关联分析.md, got: {bad_quality}")
    bad_report_path.write_text(good_report_backup, encoding="utf-8", newline="\n")
    fixed_cli = subprocess.run(
        [
            sys.executable,
            "-B",
            str(STUDENT_OS_SCRIPTS / "review_type_analysis.py"),
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "期中",
        ],
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=ROOT,
        env={**os.environ, "PYTHONDONTWRITEBYTECODE": "1", "PYTHONIOENCODING": "utf-8"},
    )
    if fixed_cli.returncode == 0:
        # Type skeletons are still incomplete, so CLI should still fail overall.
        raise AssertionError("Expected Phase B CLI to keep failing on incomplete type pages")
    fixed_quality = json.loads((state_dir / "quality-reviews.json").read_text(encoding="utf-8"))
    if fixed_quality.get("analysis_needs_revision"):
        raise AssertionError(
            f"Expected analysis reports to pass after restore, got: {fixed_quality.get('analysis_needs_revision')}"
        )
    gate_md = analysis_dir_reports / "质量门禁.md"
    ensure_contains(gate_md, "助手最多修订轮数")
    ensure_contains(gate_md, "必含区块齐全")

    # Direct reviewer helpers still cover table pipe fail/pass fixtures.
    quality_spec = importlib.util.spec_from_file_location(
        "exam_census_quality_issue51", STUDENT_OS_SCRIPTS / "exam_census_quality.py"
    )
    if quality_spec is None or quality_spec.loader is None:
        raise RuntimeError("Unable to load exam_census_quality for Issue #51 gate check")
    quality_mod = importlib.util.module_from_spec(quality_spec)
    sys.path.insert(0, str(STUDENT_OS_SCRIPTS))
    try:
        quality_spec.loader.exec_module(quality_mod)
    finally:
        if sys.path and sys.path[0] == str(STUDENT_OS_SCRIPTS):
            sys.path.pop(0)
    bad_table_doc = "\n".join(
        [
            "---",
            "type: exam-type-analysis",
            'course: "Linear Algebra"',
            'exam_scope: "期中"',
            'exam_type_id: "matrix-rank"',
            'exam_type_name: "矩阵的秩"',
            "rank: 1",
            "paper_count: 1",
            "must_know: true",
            "quality: draft",
            "status: active",
            "---",
            "",
            "| 公式 | 说明 |",
            "| --- | --- |",
            "| $|A|$ | det |",
            "",
        ]
    )
    bad_table_review = quality_mod.structural_review("题型解析/bad-table.md", bad_table_doc)
    if "markdown_tables" not in bad_table_review["failed_checks"]:
        raise AssertionError(f"Expected Phase B to flag unescaped table pipe, got: {bad_table_review}")
    fixed_table_doc = bad_table_doc.replace(
        "| $|A|$ |",
        f"| {exam_census_utils.md_table_cell('$|A|$')} |",
    )
    fixed_table_review = quality_mod.structural_review("题型解析/fixed-table.md", fixed_table_doc)
    if not fixed_table_review["checks"]["markdown_tables"]["pass"]:
        raise AssertionError(f"Expected escaped table to pass markdown_tables check: {fixed_table_review}")

    # --- Issue #63 / #65: fill template / prompt / quality-gate contracts ---
    template_path = ROOT / "student-os" / "templates" / "exam-type-analysis.md"
    template_text = template_path.read_text(encoding="utf-8")
    if re.search(r"(?m)^>\s*\|", template_text):
        raise AssertionError("exam-type-analysis.md must not use blockquote tables (> | ... |)")
    if re.search(r"(?is)<details\b[^>]*>.*?\$\$.*?\$\$.*?</details>", template_text):
        raise AssertionError("exam-type-analysis.md must not place $$ inside <details>")
    for required_heading in (
        "## 核心概念",
        "## 核心方法",
        "## 自测答案",
        "## 快速得分技巧",
        "## 易错点与检查清单",
        "### 意义 / 直观理解",
        "### 教材 / 考纲来源",
    ):
        if required_heading not in template_text:
            raise AssertionError(f"v3 template missing {required_heading}")
    if "{{concept_source_citation_or_disclaimer}}" not in template_text:
        raise AssertionError("exam-type-analysis.md must use concept_source placeholder")
    if "未参考指定教材" in template_text:
        raise AssertionError(
            "exam-type-analysis.md must not embed disclaimer phrase that bypasses concept gate"
        )
    if quality_mod.MIN_WORKED_EXAMPLES != 5 or quality_mod.MIN_SELF_TESTS != 4:
        raise AssertionError(
            f"Expected MIN_WORKED_EXAMPLES=5 / MIN_SELF_TESTS=4, got "
            f"{quality_mod.MIN_WORKED_EXAMPLES}/{quality_mod.MIN_SELF_TESTS}"
        )

    def _fm_header() -> list[str]:
        return [
            "---",
            "type: exam-type-analysis",
            'course: "Linear Algebra"',
            'exam_scope: "期中"',
            'exam_type_id: "matrix-rank"',
            'exam_type_name: "矩阵的秩"',
            "rank: 1",
            "paper_count: 9",
            "must_know: true",
            "quality: draft",
            "status: active",
            'source_summary: "共 9 份试卷"',
            "---",
            "",
            "# 1 · 矩阵的秩",
            "",
            "## 元信息",
            "",
            "**频率**：90%",
            "**分值**：12",
            "**难度**：⭐⭐",
            "**来源**：2018-2024 期中真题",
            "",
            "## 真卷对应题号",
            "",
            "| 题号 | 考点 | 本页先看 | 类型(背/手推) | 高危 |",
            "| --- | --- | --- | --- | --- |",
            "| 3 | 秩 | 是 | 手推 | 否 |",
            "",
            "## 考前速记",
            "",
            "### 方法选择树 / 决策流程",
            "",
            "```text",
            "看题目特征 → 选方法：",
            "├─ 求秩 → 行阶梯",
            "└─ 判断满秩 → 看行列式或主元",
            "```",
            "",
            "### 一眼先记住",
            "",
            "| 看到什么 | 先写什么 |",
            "| --- | --- |",
            "| 求秩 | 行阶梯 |",
            "",
            "### 关键公式表",
            "",
            "| 公式 | 什么时候用 | 先算什么 | 最容易错 |",
            "| --- | --- | --- | --- |",
            "| r(A)=非零行数 | 行阶梯后 | 主元位置 | 漏掉零行 |",
            "",
            "### 本页符号先认清",
            "",
            "| 符号 | 意思 |",
            "| --- | --- |",
            "| r(A) | 矩阵的秩 |",
            "",
            "## 核心概念",
            "",
            "### 定义",
            "",
            "- 秩是行阶梯后非零行数",
            "",
            "### 易混淆概念对比",
            "",
            "| 概念 A | 概念 B | 怎么区分 |",
            "| --- | --- | --- |",
            "| 秩 | 行列式 | 秩看主元，行列式看可逆 |",
            "",
            "### 教材 / 考纲来源",
            "",
            "- 基于考纲整理，未参考指定教材",
            "",
            "## 核心方法",
            "",
            "### 方法 1",
            "",
            "| 项目 | 内容 |",
            "| --- | --- |",
            "| 适用场景 | 求秩 |",
            "| 步骤 | 行变换后数非零行 |",
            "| 关键技巧 | 保留主元位置 |",
            "",
            "### 填空式答题模板",
            "",
            "1. 结论：`[答案]`",
            "2. 关键式：`[表达式]`",
            "",
            "## 零基础先看这里",
            "",
            "### 这类题考试到底在考什么",
            "",
            "- 考查把矩阵化到行阶梯后读出非零行数",
            "",
            "### 30秒认题",
            "",
            "| 看到什么特征 | 先按什么想 |",
            "| --- | --- |",
            "| 求 r(A) | 行变换 |",
            "",
            "### 2分钟下笔模板",
            "",
            "1. 先写：`[步骤1]`",
            "2. 再写：`[步骤2]`",
            "3. 最后：`[答案]`",
            "",
            "### 不会做时先写什么拿步骤分",
            "",
            "- 先写初等行变换记号并保留中间矩阵",
            "",
            "## 快速得分技巧",
            "",
            "| 时间情况 | 策略 | 大约可省 |",
            "| --- | --- | --- |",
            "| 时间充裕 | 完整行阶梯 | 0 |",
            "| 时间紧张 | 只追主元 | 40s |",
            "| 几乎不够 | 写变换记号 | 60s |",
            "| 完全不会 | 写定义拿步骤分 | 90s |",
            "",
            "## 易错点与检查清单",
            "",
            "| 易错点 | 错误做法 | 正确做法 | 原因 |",
            "| --- | --- | --- | --- |",
            "| 漏零行 | 直接数原矩阵行数 | 化阶梯后再数 | 行变换改变外观不改变秩 |",
            "",
        ]

    def _example_block(index: int, *, with_source: bool = True, with_teaching: bool = True, with_stars: bool = True) -> list[str]:
        source = f"来源：201{index}-201{index + 1} 第1学期 第{index}题" if with_source else "来源："
        heading_source = source if with_source else "来源：真题"
        stars = "⭐" * min(index + 1, 5) if with_stars else ""
        heading_tail = f" · 难度：{stars}" if stars else ""
        why = "因为要把未知元消掉；先做什么，为什么要消元；最后验算非零行。" if with_teaching else "按模板计算。"
        return [
            f"### 例题 {index}（{heading_source}{heading_tail}）",
            "",
            f"**题目来源**：{source}",
            "",
            f"**难度**：{stars}" if stars else "**难度**：",
            "",
            "**题目原文/摘录**：",
            "",
            f"求矩阵 A_{index} 的秩。",
            "",
            "**看到题目先判断什么**：",
            "",
            "- 先判断能否直接行变换",
            "",
            f"**方法引用**：【方法引用】2分钟下笔模板 步骤 {index}",
            "",
            "**完整解析**：",
            "",
            f"1. 先做什么，为什么：先写增广结构，{why}",
            "2. 再做什么，为什么：继续消元得到阶梯",
            "3. 如何验算：回代检查主元个数",
            "",
            "**易错点对比**：",
            "",
            "- 错法：直接套行列式",
            "- 正法：先化行阶梯",
            "",
            "**技巧总结**：",
            "",
            "- 主元位置决定秩",
            "",
        ]

    def _self_test_block(index: int, *, with_source: bool = True) -> list[str]:
        source = f"来源：202{index}-202{index + 1} 第2学期 第{index}题" if with_source else ""
        heading = f"来源：202{index}-202{index + 1} 第2学期 第{index}题" if with_source else "年份 题号"
        return [
            f"### 自测 {index}（{heading}）",
            "",
            f"**题目来源**：{source}" if source else "**题目来源**：",
            "",
            "**题目**：",
            "",
            f"判断 B_{index} 是否满秩。",
            "",
            "**提示**：",
            "",
            "- 看主元",
            "",
        ]

    def _self_test_answer_block(index: int) -> list[str]:
        return [
            f"### 自测 {index} 答案",
            "",
            "**答案与解析**：",
            "",
            f"答案：满秩。解析：行阶梯后有 {index} 个主元，因此可以验算 r={index}。",
            "",
        ]

    def _closing() -> list[str]:
        return ["## 来源校对说明", "", "- 已校对", ""]

    thin_doc = "\n".join(
        _fm_header()
        + ["## 例题精讲", ""]
        + _example_block(1)
        + _example_block(2)
        + ["## 自测题", ""]
        + _self_test_block(1)
        + _self_test_block(2)
        + ["## 自测答案", ""]
        + _self_test_answer_block(1)
        + _self_test_answer_block(2)
        + _closing()
    )
    thin_review = quality_mod.structural_review("题型解析/thin.md", thin_doc)
    if thin_review["verdict"] != "needs-revision":
        raise AssertionError(f"Expected 2+2 page to need revision, got: {thin_review}")
    if thin_review["checks"]["worked_examples"]["pass"] or thin_review["checks"]["self_tests"]["pass"]:
        raise AssertionError(f"Expected quantity checks to fail for 2+2 page: {thin_review['checks']}")

    rich_doc = "\n".join(
        _fm_header()
        + ["## 例题精讲", ""]
        + [line for i in range(1, 6) for line in _example_block(i)]
        + ["## 自测题", ""]
        + [line for i in range(1, 5) for line in _self_test_block(i)]
        + ["## 自测答案", ""]
        + [line for i in range(1, 5) for line in _self_test_answer_block(i)]
        + _closing()
    )
    rich_review = quality_mod.structural_review("题型解析/rich.md", rich_doc)
    if not rich_review["checks"]["worked_examples"]["pass"]:
        raise AssertionError(f"Expected 5 filled examples to pass worked_examples: {rich_review['checks']['worked_examples']}")
    if not rich_review["checks"]["self_tests"]["pass"]:
        raise AssertionError(f"Expected 4 filled self-tests to pass self_tests: {rich_review['checks']['self_tests']}")
    if not rich_review["checks"]["source_grounding"]["pass"]:
        raise AssertionError(f"Expected grounded sources to pass: {rich_review['checks']['source_grounding']}")
    if not rich_review["checks"]["teaching_scaffolding"]["pass"]:
        raise AssertionError(f"Expected teaching scaffolding to pass: {rich_review['checks']['teaching_scaffolding']}")
    for v3_check in (
        "concept_explanation",
        "core_methods",
        "scoring_strategy",
        "error_comparison",
        "difficulty_stars",
        "self_test_answer_separation",
        "fill_in_answer_template",
        "decision_tree",
    ):
        if not rich_review["checks"][v3_check]["pass"]:
            raise AssertionError(f"Expected rich v3 page to pass {v3_check}: {rich_review['checks'][v3_check]}")

    # Short concept without textbook cite / disclaimer must fail concept_explanation.
    bare_concept_doc = rich_doc.replace(
        "- 基于考纲整理，未参考指定教材",
        "- （待补）",
    )
    bare_concept_review = quality_mod.structural_review("题型解析/bare-concept.md", bare_concept_doc)
    if "concept_explanation" not in bare_concept_review["failed_checks"]:
        raise AssertionError(
            f"Expected concept_explanation failure without textbook cite/disclaimer, got: {bare_concept_review}"
        )
    # Citation without literal 教材/textbook words must still pass.
    cited_concept_doc = rich_doc.replace(
        "- 基于考纲整理，未参考指定教材",
        "- 参考：linear-algebra-done-right.pdf.md 第1章",
    )
    cited_concept_review = quality_mod.structural_review("题型解析/cited-concept.md", cited_concept_doc)
    if not cited_concept_review["checks"]["concept_explanation"]["pass"]:
        raise AssertionError(
            f"Expected textbook citation to pass concept_explanation: "
            f"{cited_concept_review['checks']['concept_explanation']}"
        )
    # When concept_sources exist, disclaimer alone is not enough.
    disclaimer_with_sources = quality_mod.structural_review(
        "题型解析/disclaimer-with-sources.md",
        rich_doc,
        concept_sources=[{"path": "courses/linear-algebra/references/线性代数教材.pdf.md", "label": "线性代数教材.pdf.md"}],
    )
    if "concept_explanation" not in disclaimer_with_sources["failed_checks"]:
        raise AssertionError(
            f"Expected concept_explanation failure for disclaimer while sources exist, got: {disclaimer_with_sources}"
        )
    # Unfilled template placeholder / instructional default must not pass as a disclaimer.
    template_default_doc = rich_doc.replace(
        "- 基于考纲整理，未参考指定教材",
        "- {{concept_source_citation_or_disclaimer}}",
    )
    template_default_review = quality_mod.structural_review(
        "题型解析/template-default-concept.md",
        template_default_doc,
    )
    if "concept_explanation" not in template_default_review["failed_checks"]:
        raise AssertionError(
            f"Expected concept_explanation failure for unfilled template placeholder, got: {template_default_review}"
        )
    if "no_placeholders" not in template_default_review["failed_checks"]:
        raise AssertionError(
            f"Expected no_placeholders failure for unfilled concept placeholder, got: {template_default_review}"
        )

    no_source_doc = "\n".join(
        _fm_header()
        + ["## 例题精讲", ""]
        + [line for i in range(1, 6) for line in _example_block(i, with_source=False)]
        + ["## 自测题", ""]
        + [line for i in range(1, 5) for line in _self_test_block(i, with_source=False)]
        + ["## 自测答案", ""]
        + [line for i in range(1, 5) for line in _self_test_answer_block(i)]
        + _closing()
    )
    no_source_review = quality_mod.structural_review("题型解析/no-source.md", no_source_doc)
    if "source_grounding" not in no_source_review["failed_checks"]:
        raise AssertionError(f"Expected source_grounding failure without 来源, got: {no_source_review}")

    details_math_doc = rich_doc + "\n<details>\n<summary>答案</summary>\n\n$$\nA=1\n$$\n\n</details>\n"
    details_review = quality_mod.structural_review("题型解析/details-math.md", details_math_doc)
    if "render_safe_markdown" not in details_review["failed_checks"]:
        raise AssertionError(f"Expected render_safe_markdown for details+$$, got: {details_review}")

    quote_table_doc = rich_doc.replace(
        "## 元信息\n\n**频率**：90%",
        "## 元信息\n\n> **元信息**\n> | 项目 | 内容 |\n> | --- | --- |\n> | 频率 | 96% |\n\n**频率**：90%",
    )
    quote_review = quality_mod.structural_review("题型解析/quote-table.md", quote_table_doc)
    if "render_safe_markdown" not in quote_review["failed_checks"]:
        raise AssertionError(f"Expected render_safe_markdown for blockquote table, got: {quote_review}")

    weak_teach_doc = rich_doc
    weak_teach_doc = weak_teach_doc.replace("### 方法选择树 / 决策流程", "### 方法罗列")
    weak_teach_doc = weak_teach_doc.replace("├─ 求秩 → 行阶梯", "- 求秩")
    weak_teach_doc = weak_teach_doc.replace("└─ 判断满秩 → 看行列式或主元", "- 判断满秩")
    weak_teach_doc = weak_teach_doc.replace("### 不会做时先写什么拿步骤分", "### 临时提示")
    weak_teach_doc = weak_teach_doc.replace("**易错点对比**：", "**备注**：")
    weak_teach_doc = weak_teach_doc.replace("- 错法：直接套行列式", "- 注意书写")
    weak_teach_doc = weak_teach_doc.replace("- 正法：先化行阶梯", "- 保持整洁")
    weak_teach_doc = weak_teach_doc.replace("3. 如何验算：回代检查主元个数", "3. 写出结果")
    weak_teach_doc = weak_teach_doc.replace(
        "因为要把未知元消掉；先做什么，为什么要消元；最后验算非零行。",
        "按模板计算。",
    )
    weak_teach_doc = weak_teach_doc.replace("| 易错点 | 错误做法 | 正确做法 | 原因 |", "| 提醒 | 写法A | 写法B | 说明 |")
    weak_teach_doc = weak_teach_doc.replace(
        "| 漏零行 | 直接数原矩阵行数 | 化阶梯后再数 | 行变换改变外观不改变秩 |",
        "| 漏零行 | 粗心 | 细心 | 习惯 |",
    )
    weak_teach_review = quality_mod.structural_review("题型解析/weak-teach.md", weak_teach_doc)
    if "teaching_scaffolding" not in weak_teach_review["failed_checks"]:
        raise AssertionError(f"Expected teaching_scaffolding failure, got: {weak_teach_review}")

    # Blank 【方法引用】 must not count via cross-line \\S+ match onto the next heading.
    blank_method_doc = rich_doc.replace(
        "**方法引用**：【方法引用】2分钟下笔模板 步骤 1",
        "**方法引用**：【方法引用】",
    )
    for i in range(2, 6):
        blank_method_doc = blank_method_doc.replace(
            f"**方法引用**：【方法引用】2分钟下笔模板 步骤 {i}",
            "**方法引用**：【方法引用】",
        )
    blank_method_review = quality_mod.structural_review("题型解析/blank-method.md", blank_method_doc)
    if blank_method_review["checks"]["worked_examples"]["pass"]:
        raise AssertionError("Expected blank same-line 【方法引用】 to fail worked_examples")
    if blank_method_review["checks"]["method_reference"]["pass"]:
        raise AssertionError("Expected blank same-line 【方法引用】 to fail method_reference")

    # Evidence-short + needs-review relaxes quantity gates; same counts without the marker still fail.
    short_base = "\n".join(
        _fm_header()
        + ["## 例题精讲", ""]
        + _example_block(1)
        + _example_block(2)
        + ["## 自测题", ""]
        + _self_test_block(1)
        + ["## 自测答案", ""]
        + _self_test_answer_block(1)
        + ["## 来源校对说明", "", f"- {quality_mod.EVIDENCE_SHORT_MARKER}", ""]
    )
    short_relaxed = short_base.replace("quality: draft", "quality: needs-review")
    short_relaxed_review = quality_mod.structural_review("题型解析/short-relaxed.md", short_relaxed)
    if not short_relaxed_review["checks"]["worked_examples"]["pass"]:
        raise AssertionError(
            f"Expected evidence-short needs-review path to pass worked_examples: "
            f"{short_relaxed_review['checks']['worked_examples']}"
        )
    if not short_relaxed_review["checks"]["self_tests"]["pass"]:
        raise AssertionError(
            f"Expected evidence-short needs-review path to pass self_tests: "
            f"{short_relaxed_review['checks']['self_tests']}"
        )
    if not short_relaxed_review["checks"]["method_reference"]["pass"]:
        raise AssertionError(
            f"Expected evidence-short needs-review path to pass method_reference: "
            f"{short_relaxed_review['checks']['method_reference']}"
        )
    short_strict_review = quality_mod.structural_review("题型解析/short-strict.md", short_base)
    if short_strict_review["checks"]["worked_examples"]["pass"] or short_strict_review["checks"]["self_tests"]["pass"]:
        raise AssertionError(
            f"Expected 2+1 without needs-review to fail quantity checks: {short_strict_review['checks']}"
        )

    # Issue #65: inline answers in 自测题 should fail separation; missing stars should fail difficulty_stars.
    inline_answers_doc = "\n".join(
        _fm_header()
        + ["## 例题精讲", ""]
        + [line for i in range(1, 6) for line in _example_block(i)]
        + ["## 自测题", ""]
        + [
            line
            for i in range(1, 5)
            for line in (
                _self_test_block(i)
                + ["**答案与解析**：", "", f"答案：满秩。解析：有 {i} 个主元。", ""]
            )
        ]
        + ["## 自测答案", "", "- 见上", ""]
        + _closing()
    )
    inline_review = quality_mod.structural_review("题型解析/inline-answers.md", inline_answers_doc)
    if "self_test_answer_separation" not in inline_review["failed_checks"]:
        raise AssertionError(f"Expected self_test_answer_separation failure, got: {inline_review}")

    no_stars_doc = "\n".join(
        _fm_header()
        + ["## 例题精讲", ""]
        + [line for i in range(1, 6) for line in _example_block(i, with_stars=False)]
        + ["## 自测题", ""]
        + [line for i in range(1, 5) for line in _self_test_block(i)]
        + ["## 自测答案", ""]
        + [line for i in range(1, 5) for line in _self_test_answer_block(i)]
        + _closing()
    )
    no_stars_review = quality_mod.structural_review("题型解析/no-stars.md", no_stars_doc)
    if "difficulty_stars" not in no_stars_review["failed_checks"]:
        raise AssertionError(f"Expected difficulty_stars failure, got: {no_stars_review}")

    # Pair answers by question number, not list position.
    mismatched_answers = "\n".join(
        _fm_header()
        + ["## 例题精讲", ""]
        + [line for i in range(1, 6) for line in _example_block(i)]
        + ["## 自测题", ""]
        + [line for i in range(1, 5) for line in _self_test_block(i)]
        + ["## 自测答案", ""]
        + _self_test_answer_block(2)
        + _self_test_answer_block(3)
        + _self_test_answer_block(4)
        + _self_test_answer_block(5)
        + _closing()
    )
    mismatched_review = quality_mod.structural_review("题型解析/mismatch-answers.md", mismatched_answers)
    if mismatched_review["checks"]["self_tests"]["pass"]:
        raise AssertionError("Expected self_tests to fail when answer numbers do not match questions")

    # Inline **答案**： should fail separation even if ## 自测答案 exists.
    answer_colon_doc = "\n".join(
        _fm_header()
        + ["## 例题精讲", ""]
        + [line for i in range(1, 6) for line in _example_block(i)]
        + ["## 自测题", ""]
        + [
            line
            for i in range(1, 5)
            for line in (_self_test_block(i) + ["**答案**：", "", f"满秩，r={i}", ""])
        ]
        + ["## 自测答案", ""]
        + [line for i in range(1, 5) for line in _self_test_answer_block(i)]
        + _closing()
    )
    answer_colon_review = quality_mod.structural_review("题型解析/answer-colon.md", answer_colon_doc)
    if "self_test_answer_separation" not in answer_colon_review["failed_checks"]:
        raise AssertionError(f"Expected separation failure for inline **答案**：, got: {answer_colon_review}")

    # Checklist boxes alone must not satisfy fill-in template gate.
    checklist_only = rich_doc
    for token in ("[答案]", "[表达式]", "[步骤1]", "[步骤2]"):
        checklist_only = checklist_only.replace(token, "（待填）")
    checklist_review = quality_mod.structural_review("题型解析/checklist-only.md", checklist_only)
    if "fill_in_answer_template" not in checklist_review["failed_checks"]:
        raise AssertionError(f"Expected fill_in_answer_template failure for checklist-only, got: {checklist_review}")

    # Empty scoring / error tables should fail even when headers/labels exist.
    blank_scoring = rich_doc
    blank_scoring = blank_scoring.replace("| 时间充裕 | 完整行阶梯 | 0 |", "| 时间充裕 |  |  |")
    blank_scoring = blank_scoring.replace("| 时间紧张 | 只追主元 | 40s |", "| 时间紧张 |  |  |")
    blank_scoring = blank_scoring.replace("| 几乎不够 | 写变换记号 | 60s |", "| 几乎不够 |  |  |")
    blank_scoring = blank_scoring.replace("| 完全不会 | 写定义拿步骤分 | 90s |", "| 完全不会 |  |  |")
    blank_scoring_review = quality_mod.structural_review("题型解析/blank-scoring.md", blank_scoring)
    if "scoring_strategy" not in blank_scoring_review["failed_checks"]:
        raise AssertionError(f"Expected scoring_strategy failure for blank strategies, got: {blank_scoring_review}")

    blank_errors = rich_doc.replace(
        "| 漏零行 | 直接数原矩阵行数 | 化阶梯后再数 | 行变换改变外观不改变秩 |",
        "|  |  |  |  |",
    )
    blank_errors_review = quality_mod.structural_review("题型解析/blank-errors.md", blank_errors)
    if "error_comparison" not in blank_errors_review["failed_checks"]:
        raise AssertionError(f"Expected error_comparison failure for empty rows, got: {blank_errors_review}")

    deep = json.loads(
        run_script(
            "init_exam_deep_dive.py",
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "期中",
            "--limit",
            "2",
            "--overwrite",
        )
    )
    if not deep.get("written"):
        raise AssertionError(f"Expected Phase D deep-dive scaffolds, got: {deep}")
    deep_dir = repo / "courses" / "linear-algebra" / "reviews" / "期中" / "真题精析"
    ensure_exists(deep_dir)
    if not any(deep_dir.glob("*-精析.md")):
        raise AssertionError("Expected at least one deep-dive markdown under 真题精析/")

    assert_prep_pack_templates()
    reviews_midterm = repo / "courses" / "linear-algebra" / "reviews" / "期中"
    analysis_dir_for_pack = reviews_midterm / "题型解析"
    skeleton_names = sorted(path.name for path in analysis_dir_for_pack.glob("*.md"))
    if not skeleton_names:
        raise AssertionError("Expected type-analysis skeletons before prep pack validation")

    # Missing prep pack files must fail Phase E and list missing_files.
    for name in ("备考指南.md", "公式总卡.md", "答题模板速查.md", "考前1小时清单.md"):
        path = reviews_midterm / name
        if path.exists():
            path.unlink()
    missing_code, missing_payload = _run_cross_validate(repo)
    if missing_code == 0:
        raise AssertionError("Expected Phase E to fail when prep pack files are missing")
    missing_files = list((missing_payload.get("prep_pack") or {}).get("missing_files") or [])
    for required in ("备考指南.md", "公式总卡.md", "答题模板速查.md", "考前1小时清单.md"):
        if required not in missing_files:
            raise AssertionError(f"Expected missing_files to include {required}, got: {missing_files}")

    write_minimal_prep_pack(reviews_midterm, skeleton_names)
    cross = json.loads(
        run_script(
            "cross_validate_exam_census.py",
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "期中",
        )
    )
    if not cross.get("ok"):
        raise AssertionError(f"Expected clean cross-validation with prep pack, got: {cross}")
    prep_pack = cross.get("prep_pack") or {}
    if prep_pack.get("missing_files"):
        raise AssertionError(f"Expected no prep_pack.missing_files, got: {prep_pack}")
    if prep_pack.get("layer_link_issues") or prep_pack.get("content_issues"):
        raise AssertionError(f"Expected clean prep_pack issues, got: {prep_pack}")
    for key, filename in (
        ("prep_guide", "备考指南.md"),
        ("formula_card", "公式总卡.md"),
        ("answer_templates", "答题模板速查.md"),
        ("one_hour_checklist", "考前1小时清单.md"),
    ):
        recorded = ((prep_pack.get("files") or {}).get(key) or {}).get("path") or ""
        if recorded.startswith("/") or re.match(r"^[A-Za-z]:[\\/]", recorded):
            raise AssertionError(f"prep_pack path must be repo-relative, got {key}={recorded!r}")
        if filename not in recorded.replace("\\", "/"):
            raise AssertionError(f"Expected repo-relative path containing {filename}, got: {recorded}")
    coverage_report = reviews_midterm / "analysis" / "覆盖率检查.md"
    ensure_exists(coverage_report)
    ensure_contains(coverage_report, "Prep pack 四层结构")
    ensure_contains(coverage_report, "备考指南.md")
    ensure_contains(coverage_report, "公式总卡.md")
    ensure_contains(coverage_report, "答题模板速查.md")
    ensure_contains(coverage_report, "考前1小时清单.md")

    # Prep guide without L3/L4 links → layer_link_issues.
    prep_guide = reviews_midterm / "备考指南.md"
    prep_guide.write_text(
        "\n".join(
            [
                "---",
                "type: exam-prep-guide",
                'course: "linear-algebra"',
                "status: active",
                "review_scope: exam-census",
                "---",
                "",
                "# Prep guide",
                "",
                "## 怎么使用这套资料",
                "",
                "See [题型解析/](题型解析/) only.",
                "",
                "## 题型优先级",
                "",
                "| 优先级 | 题型 | 出现率 | 建议投入时间 | 入口 |",
                "| --- | --- | ---: | --- | --- |",
                "\n".join(
                    f"| P0 | {name} | — | — | [题型解析/{name}](题型解析/{name}) |"
                    for name in skeleton_names
                ),
                "",
                "## 复习时间分配",
                "",
                "| 时间总量 | 先做什么 | 目标 |",
                "| --- | --- | --- |",
                "| 1 小时 | 题型解析 | 定向 |",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    link_code, link_payload = _run_cross_validate(repo)
    if link_code == 0:
        raise AssertionError("Expected Phase E to fail when prep guide lacks L3/L4 links")
    link_issues = list((link_payload.get("prep_pack") or {}).get("layer_link_issues") or [])
    if not any("公式总卡.md" in item for item in link_issues):
        raise AssertionError(f"Expected layer_link_issues about 公式总卡.md, got: {link_issues}")
    if not any("答题模板速查.md" in item for item in link_issues):
        raise AssertionError(f"Expected layer_link_issues about 答题模板速查.md, got: {link_issues}")
    if not any("考前1小时清单.md" in item for item in link_issues):
        raise AssertionError(f"Expected layer_link_issues about 考前1小时清单.md, got: {link_issues}")

    # Formula card with only header/empty rows (no filled data) → content_issues.
    write_minimal_prep_pack(reviews_midterm, skeleton_names)
    (reviews_midterm / "公式总卡.md").write_text(
        "\n".join(
            [
                "---",
                "type: formula-cheat-sheet",
                'course: "linear-algebra"',
                "status: active",
                "review_scope: exam-census",
                "---",
                "",
                "# Formula card",
                "",
                "See [题型解析/](题型解析/).",
                "",
                "## 高频公式速查",
                "",
                "| 题型 | 看到什么 | 公式 / 结论 | 先算什么 | 最容易错 | 来源 |",
                "| --- | --- | --- | --- | --- | --- |",
                "|  |  |  |  |  |  |",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    empty_formula_code, empty_formula_payload = _run_cross_validate(repo)
    if empty_formula_code == 0:
        raise AssertionError("Expected Phase E to fail when formula table has no filled data rows")
    empty_formula_issues = list((empty_formula_payload.get("prep_pack") or {}).get("content_issues") or [])
    if not any("filled formula" in item or "table data" in item for item in empty_formula_issues):
        raise AssertionError(f"Expected empty-formula content_issues, got: {empty_formula_issues}")

    # Formula card without 题型解析/ links → layer_link_issues or content_issues.
    write_minimal_prep_pack(reviews_midterm, skeleton_names)
    (reviews_midterm / "公式总卡.md").write_text(
        "\n".join(
            [
                "---",
                "type: formula-cheat-sheet",
                'course: "linear-algebra"',
                "status: active",
                "review_scope: exam-census",
                "---",
                "",
                "# Formula card",
                "",
                "## 高频公式速查",
                "",
                "| 题型 | 看到什么 | 公式 / 结论 | 先算什么 | 最容易错 | 来源 |",
                "| --- | --- | --- | --- | --- | --- |",
                "| matrix-rank | 秩 | rank | 化阶梯 | 漏零行 | 待补 |",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    formula_code, formula_payload = _run_cross_validate(repo)
    if formula_code == 0:
        raise AssertionError("Expected Phase E to fail when formula card has no 题型解析 link")
    formula_pack = formula_payload.get("prep_pack") or {}
    formula_issues = list(formula_pack.get("layer_link_issues") or []) + list(
        formula_pack.get("content_issues") or []
    )
    if not any("公式总卡" in item and "题型解析" in item for item in formula_issues):
        raise AssertionError(f"Expected formula-card 题型解析 issue, got: {formula_issues}")

    # Answer templates with placeholders only outside 标准答题模板 → content_issues.
    write_minimal_prep_pack(reviews_midterm, skeleton_names)
    (reviews_midterm / "答题模板速查.md").write_text(
        "\n".join(
            [
                "---",
                "type: answer-template-quickref",
                'course: "linear-algebra"',
                "status: active",
                "review_scope: exam-census",
                "---",
                "",
                "# Answer templates",
                "",
                "Usage note only: placeholders like [条件] / [答案] live here, not in the table.",
                "",
                "## 标准答题模板",
                "",
                "| 题型 | 看到什么 | 第一句写什么 | 填空式模板 | 来源 |",
                "| --- | --- | --- | --- | --- |",
                "| matrix-rank | 求秩 | 先化阶梯 | 直接写结论 | [题型解析/](题型解析/) |",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    tmpl_code, tmpl_payload = _run_cross_validate(repo)
    if tmpl_code == 0:
        raise AssertionError(
            "Expected Phase E to fail when fill-in placeholders are only outside 标准答题模板"
        )
    tmpl_issues = list((tmpl_payload.get("prep_pack") or {}).get("content_issues") or [])
    if not any("填空" in item or "placeholder" in item or "标准答题模板" in item for item in tmpl_issues):
        raise AssertionError(f"Expected fill-in placeholder content_issues, got: {tmpl_issues}")

    # One-hour checklist missing time slots → content_issues.
    write_minimal_prep_pack(reviews_midterm, skeleton_names)
    (reviews_midterm / "考前1小时清单.md").write_text(
        "\n".join(
            [
                "---",
                "type: pre-exam-one-hour-checklist",
                'course: "linear-algebra"',
                "status: active",
                "review_scope: exam-census",
                "---",
                "",
                "# Checklist",
                "",
                "See [备考指南.md](备考指南.md), [公式总卡.md](公式总卡.md), "
                "[答题模板速查.md](答题模板速查.md), [题型解析/](题型解析/).",
                "",
                "- [ ] something",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    hour_code, hour_payload = _run_cross_validate(repo)
    if hour_code == 0:
        raise AssertionError("Expected Phase E to fail when one-hour checklist lacks time slots")
    hour_issues = list((hour_payload.get("prep_pack") or {}).get("content_issues") or [])
    if not any("60-45" in item or "time slots" in item for item in hour_issues):
        raise AssertionError(f"Expected time-slot content_issues, got: {hour_issues}")

    # Empty checklist boxes without text → content_issues.
    write_minimal_prep_pack(reviews_midterm, skeleton_names)
    (reviews_midterm / "考前1小时清单.md").write_text(
        "\n".join(
            [
                "---",
                "type: pre-exam-one-hour-checklist",
                'course: "linear-algebra"',
                "status: active",
                "review_scope: exam-census",
                "---",
                "",
                "# Checklist",
                "",
                "| 时间 | 做什么 | 文件 | 目标 |",
                "| --- | --- | --- | --- |",
                "| 60-45 分钟 | P0 | [备考指南.md](备考指南.md) / [题型解析/](题型解析/) | 方法 |",
                "| 45-30 分钟 | 公式 | [公式总卡.md](公式总卡.md) | 背诵 |",
                "| 30-15 分钟 | 模板 | [答题模板速查.md](答题模板速查.md) | 步骤分 |",
                "| 15-5 分钟 | 易错 | [题型解析/](题型解析/) | 避坑 |",
                "| 5-0 分钟 | checklist | 本文件 | 稳住 |",
                "",
                "## 最后检查",
                "",
                "- [ ]",
                "- [ ]   ",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    empty_box_code, empty_box_payload = _run_cross_validate(repo)
    if empty_box_code == 0:
        raise AssertionError("Expected Phase E to fail when checklist boxes have no text")
    empty_box_issues = list((empty_box_payload.get("prep_pack") or {}).get("content_issues") or [])
    if not any("checklist" in item.lower() or "检查" in item for item in empty_box_issues):
        raise AssertionError(f"Expected empty-checklist content_issues, got: {empty_box_issues}")

    # Restore a clean prep pack for any later consumers of this fixture.
    write_minimal_prep_pack(reviews_midterm, skeleton_names)
    restored_code, restored_payload = _run_cross_validate(repo)
    if restored_code != 0 or not restored_payload.get("ok"):
        raise AssertionError(f"Expected restored prep pack to pass Phase E, got: {restored_payload}")

    # Legacy: guide with no type links must still fail Phase E.
    prep_guide.write_text(
        "\n".join(
            [
                "---",
                "type: exam-prep-guide",
                'course: "Linear Algebra"',
                "status: draft",
                "review_scope: exam-census",
                "---",
                "",
                "# Prep guide",
                "",
                "## 怎么使用这套资料",
                "",
                "See [公式总卡.md](公式总卡.md), [答题模板速查.md](答题模板速查.md), "
                "[考前1小时清单.md](考前1小时清单.md), [题型解析/](题型解析/).",
                "",
                "## 题型优先级",
                "",
                "| 优先级 | 题型 | 出现率 | 建议投入时间 | 入口 |",
                "| --- | --- | ---: | --- | --- |",
                "| P0 | matrix-rank | high | 1h | 题型解析/ |",
                "",
                "## 复习时间分配",
                "",
                "| 时间总量 | 先做什么 | 目标 |",
                "| --- | --- | --- |",
                "| 1 小时 | | |",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    linked_fail_code, linked_fail_payload = _run_cross_validate(repo)
    if linked_fail_code == 0:
        raise AssertionError("Expected Phase E to fail when prep guide has no type links")
    if not linked_fail_payload.get("prep_guide_unlinked_types"):
        raise AssertionError(
            f"Expected prep_guide_unlinked_types when guide lacks type links, got: {linked_fail_payload}"
        )
    write_minimal_prep_pack(reviews_midterm, skeleton_names)

    # --- Issue #61 contracts: 文本/ discovery, annotation aliases, taxonomy dump, confidence ---
    chinese_course = repo / "courses" / "线性代数"
    scope_root = chinese_course / "reviews" / "期中"
    text_dir = scope_root / "文本"
    text_dir.mkdir(parents=True, exist_ok=True)
    (text_dir / "2018-期中.pdf.md").write_text(
        "\n".join(
            [
                "---",
                'course: "线性代数"',
                "---",
                "",
                "# 2018 期中",
                "",
                "1. 行列式",
                "",
            ]
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    text_init = json.loads(
        run_script(
            "init_exam_census.py",
            str(repo),
            "--course",
            "线性代数",
            "--exam-scope",
            "期中",
            "--papers-dir",
            "courses/线性代数/reviews/期中",
            "--overwrite",
        )
    )
    if text_init.get("papers_dir_fallback_subdir") != "文本":
        raise AssertionError(f"Expected 文本/ fallback, got: {text_init}")
    if "文本" not in str(text_init.get("papers_dir", "")):
        raise AssertionError(f"Expected papers_dir under 文本/, got: {text_init}")
    text_state = Path(text_init["state_dir"])
    text_manifest = json.loads((text_state / "manifest.json").read_text(encoding="utf-8"))
    text_paper = text_manifest["papers"][0]
    text_stem = text_paper["stem"]
    text_path = text_paper["path"]
    if "文本/" not in text_path.replace("\\", "/"):
        raise AssertionError(f"Manifest path should include 文本/: {text_path}")

    # PyYAML-dump style taxonomy (id not first under list item).
    (text_state / "taxonomy.yaml").write_text(
        "\n".join(
            [
                "version: 1",
                "course: 线性代数",
                "exam_scope: 期中",
                "types:",
                "- aliases: []",
                "  id: determinant",
                "  keywords: []",
                "  name: 行列式",
                "",
            ]
        ),
        encoding="utf-8",
        newline="\n",
    )
    # Wrong filename + wrong source (missing 文本/) + medium confidence.
    wrong_name = text_state / "annotations" / f"{text_stem}.pdf.md.json"
    wrong_name.write_text(
        json.dumps(
            {
                "source": "courses/线性代数/reviews/期中/2018-期中.pdf.md",
                "exam_label": "2018 期中",
                "types_present": ["determinant"],
                "type_counts": {"determinant": 1},
                "confidence": "medium",
                "notes": "alias filename + stale source",
            },
            ensure_ascii=False,
            indent=2,
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    text_stats = json.loads(
        run_script(
            "build_exam_type_stats.py",
            str(repo),
            "--course",
            "线性代数",
            "--exam-scope",
            "期中",
            "--papers-dir",
            "courses/线性代数/reviews/期中/文本",
            "--validate",
            "--overwrite",
        )
    )
    if text_stats.get("unknown_types"):
        raise AssertionError(
            f"PyYAML-style taxonomy must not yield unknown_types: {text_stats['unknown_types']}"
        )
    if not text_stats.get("annotation_aliases_used"):
        raise AssertionError(f"Expected annotation_aliases_used, got: {text_stats}")
    if not text_stats.get("source_mismatches"):
        raise AssertionError(f"Expected source_mismatches, got: {text_stats}")
    if not text_stats.get("medium_confidence"):
        raise AssertionError(f"Expected medium_confidence diagnostics, got: {text_stats}")
    if text_stats.get("low_confidence"):
        raise AssertionError("medium must not be treated as low_confidence")
    if "ignored" not in str(text_stats.get("papers_dir_warning", "")).lower():
        raise AssertionError(f"Expected papers_dir ignored warning, got: {text_stats}")
    text_report = (
        chinese_course / "reviews" / "期中" / "题型频率统计.md"
    ).read_text(encoding="utf-8")
    if "文本/2018-期中.pdf.md" not in text_report.replace("\\", "/"):
        raise AssertionError(
            "Frequency report links must use manifest path including 文本/"
        )
    if "reviews/期中/2018-期中.pdf.md)" in text_report.replace("\\", "/"):
        raise AssertionError("Report must not link the stale annotation.source path")
    if "### Annotation load errors" not in text_report:
        raise AssertionError("Validation report must include Annotation load errors section")

    # Root-relative pattern must keep original papers-dir (not switch into 文本/).
    scope_root = (chinese_course / "reviews" / "期中").resolve()
    resolved_dir, fallback_subdir, effective_pattern = exam_census_utils.resolve_papers_dir(
        scope_root, "文本/*.pdf.md"
    )
    if fallback_subdir is not None:
        raise AssertionError(
            f"Root-relative pattern should not trigger subdir fallback: {fallback_subdir}"
        )
    if effective_pattern != "文本/*.pdf.md":
        raise AssertionError(f"Expected original pattern preserved: {effective_pattern}")
    if resolved_dir != scope_root:
        raise AssertionError(f"papers_dir should stay at scope root: {resolved_dir}")

    # Explicit empty confidence is invalid.
    empty_conf = text_state / "annotations" / f"{text_stem}.json"
    empty_conf.write_text(
        json.dumps(
            {
                "source": text_path,
                "exam_label": "2018 期中",
                "types_present": ["determinant"],
                "type_counts": {"determinant": 1},
                "confidence": "",
            },
            ensure_ascii=False,
            indent=2,
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    wrong_name.unlink(missing_ok=True)
    empty_conf_run = subprocess.run(
        [
            sys.executable,
            "-B",
            str(STUDENT_OS_SCRIPTS / "build_exam_type_stats.py"),
            str(repo),
            "--course",
            "线性代数",
            "--exam-scope",
            "期中",
            "--validate",
            "--overwrite",
        ],
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=ROOT,
        env={**os.environ, "PYTHONDONTWRITEBYTECODE": "1", "PYTHONIOENCODING": "utf-8"},
    )
    if empty_conf_run.returncode == 0:
        raise AssertionError("Expected --validate failure for empty confidence")
    empty_stats = json.loads(empty_conf_run.stdout)
    if not empty_stats.get("invalid_confidence"):
        raise AssertionError(f"Expected invalid_confidence for empty string: {empty_stats}")

    # Invalid confidence fails --validate.
    empty_conf.write_text(
        json.dumps(
            {
                "source": text_path,
                "exam_label": "2018 期中",
                "types_present": ["determinant"],
                "type_counts": {"determinant": 1},
                "confidence": "kinda-sure",
            },
            ensure_ascii=False,
            indent=2,
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    invalid_conf = subprocess.run(
        [
            sys.executable,
            "-B",
            str(STUDENT_OS_SCRIPTS / "build_exam_type_stats.py"),
            str(repo),
            "--course",
            "线性代数",
            "--exam-scope",
            "期中",
            "--validate",
            "--overwrite",
        ],
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
        cwd=ROOT,
        env={**os.environ, "PYTHONDONTWRITEBYTECODE": "1", "PYTHONIOENCODING": "utf-8"},
    )
    if invalid_conf.returncode == 0:
        raise AssertionError("Expected --validate failure for invalid confidence")
    invalid_payload = json.loads(invalid_conf.stdout)
    if not invalid_payload.get("invalid_confidence"):
        raise AssertionError(f"Expected invalid_confidence list, got: {invalid_payload}")

    # Platform adapters install into the vault (not the skill dir).
    adapter_payload = json.loads(
        run_script(
            "install_exam_census_adapters.py",
            str(repo),
            "--platforms",
            "claude,cursor,opencode,github",
            "--json",
        )
    )
    # Default Claude installs skill + command (2 files); plus cursor/opencode/github = 5.
    if adapter_payload.get("installed") != 5:
        raise AssertionError(f"Expected 5 adapter files installed, got: {adapter_payload}")
    claude_skill = repo / ".claude" / "skills" / "exam-census" / "SKILL.md"
    claude_cmd = repo / ".claude" / "commands" / "exam-census.md"
    claude_wf = repo / ".claude" / "workflows" / "exam-census.js"
    cursor_rule = repo / ".cursor" / "rules" / "exam-census.mdc"
    ensure_exists(claude_skill)
    ensure_exists(claude_cmd)
    ensure_exists(cursor_rule)
    ensure_exists(repo / ".opencode" / "exam-census.md")
    ensure_exists(repo / ".github" / "copilot-exam-census.md")
    if claude_wf.exists():
        raise AssertionError("Default Claude install must not copy experimental workflow JS")
    ensure_contains(claude_skill, "name: exam-census")
    ensure_contains(claude_skill, "/exam-census")
    ensure_contains(claude_skill, "init_exam_census.py")
    ensure_contains(claude_skill, "build_exam_type_stats.py")
    ensure_contains(claude_skill, "Runbook")
    ensure_contains(claude_skill, "--papers-dir")
    ensure_contains(claude_skill, "只有本脚本支持")
    ensure_contains(claude_skill, "暂停并询问")
    ensure_contains(claude_skill, "annotations/<stem>.json")
    ensure_contains(claude_skill, "high")
    ensure_contains(claude_skill, "medium")
    skill_text = claude_skill.read_text(encoding="utf-8")
    if "disable-model-invocation: true" in skill_text:
        raise AssertionError(
            "Claude exam-census skill must allow model/natural-language invocation"
        )
    if "Workflow(" in skill_text:
        raise AssertionError(
            "Claude exam-census SKILL.md must not mention Workflow(…) calls"
        )
    ensure_contains(claude_cmd, "/exam-census")
    ensure_contains(claude_cmd, "runbook")
    cmd_text = claude_cmd.read_text(encoding="utf-8")
    if "Workflow(" in cmd_text:
        raise AssertionError(
            "Claude exam-census command must not mention Workflow(…) calls"
        )
    ensure_contains(cursor_rule, "alwaysApply: false")
    ensure_contains(cursor_rule, "exam-census")

    claude_result = next(
        item for item in adapter_payload["results"] if item["platform"] == "claude"
    )
    claude_dests = {Path(d).name for d in claude_result.get("destinations", [])}
    if claude_dests != {"SKILL.md", "exam-census.md"}:
        raise AssertionError(
            f"Expected Claude destinations SKILL.md + exam-census.md, got: {claude_result}"
        )
    if len(claude_result.get("files", [])) != 2:
        raise AssertionError(f"Expected 2 Claude file results, got: {claude_result}")

    student_os = ROOT / "student-os"

    # Claude skill/command adapters must be self-contained runbooks (no Workflow tool).
    for doc_rel in (
        Path("integrations") / "claude" / "skills" / "exam-census" / "SKILL.md",
        Path("integrations") / "claude" / "commands" / "exam-census.md",
    ):
        doc_text = (student_os / doc_rel).read_text(encoding="utf-8")
        if "Workflow(" in doc_text:
            raise AssertionError(
                f"{doc_rel} must not mention Workflow(…) (Issue #59)"
            )
        for line in doc_text.splitlines():
            lowered = line.lower()
            if "invoke:" in lowered and "workflow" in lowered:
                raise AssertionError(
                    f"{doc_rel} must not instruct Invoke Workflow on same line: {line!r}"
                )

    # Broader docs must not recommend Workflow name/scriptPath as the primary entry.
    for doc_rel in (
        Path("commands") / "exam-census.md",
        Path("references") / "exam-census-workflow.md",
    ):
        doc_text = (student_os / doc_rel).read_text(encoding="utf-8")
        recommend_patterns = (
            "推荐使用 Workflow({name",
            'use Workflow({name: "exam-census"})',
            "Invoke via Workflow({name",
            "Invoke: Workflow(",
            "Workflow({name:",
            "Workflow({ name:",
            "Workflow({scriptPath:",
            "Workflow({ scriptPath:",
            'Workflow({scriptPath: ".claude/workflows/exam-census.js"})',
        )
        for pattern in recommend_patterns:
            if pattern in doc_text:
                raise AssertionError(f"{doc_rel} still recommends Workflow entry: {pattern}")

    # Phase 5 prep-pack layer contract must be documented.
    for doc_rel in (
        Path("references") / "exam-census-workflow.md",
        Path("references") / "exam-census-quality.md",
        Path("commands") / "exam-census.md",
        Path("integrations") / "claude" / "skills" / "exam-census" / "SKILL.md",
    ):
        doc_text = (student_os / doc_rel).read_text(encoding="utf-8")
        for needle in ("L1", "L2", "L3", "L4", "备考指南", "公式总卡", "答题模板速查", "考前1小时清单"):
            if needle not in doc_text:
                raise AssertionError(f"{doc_rel} missing prep-pack layer marker {needle!r}")
        if "Phase 5" not in doc_text and "Prep pack" not in doc_text:
            raise AssertionError(f"{doc_rel} must mention Phase 5 / Prep pack")

    # Claude integration templates must not contain dangerous control characters.
    adapter_spec = importlib.util.spec_from_file_location(
        "install_exam_census_adapters_smoke",
        STUDENT_OS_SCRIPTS / "install_exam_census_adapters.py",
    )
    if adapter_spec is None or adapter_spec.loader is None:
        raise RuntimeError("Unable to load install_exam_census_adapters for control-char scan")
    adapter_mod = importlib.util.module_from_spec(adapter_spec)
    sys.path.insert(0, str(STUDENT_OS_SCRIPTS))
    try:
        adapter_spec.loader.exec_module(adapter_mod)
    finally:
        if sys.path and sys.path[0] == str(STUDENT_OS_SCRIPTS):
            sys.path.pop(0)

    claude_templates = [
        student_os / "integrations" / "claude" / "skills" / "exam-census" / "SKILL.md",
        student_os / "integrations" / "claude" / "commands" / "exam-census.md",
        student_os / "integrations" / "claude" / "workflows" / "exam-census.js",
    ]
    for template in claude_templates:
        hits = adapter_mod.scan_integration_template(template)
        if hits:
            raise AssertionError(
                f"Dangerous control characters in {template}: {hits[:5]}"
            )

    # Second install without --force should skip both Claude files.
    skip_payload = json.loads(
        run_script(
            "install_exam_census_adapters.py",
            str(repo),
            "--platforms",
            "claude",
            "--json",
        )
    )
    if skip_payload.get("skipped") != 2 or skip_payload.get("installed") != 0:
        raise AssertionError(f"Expected claude adapter skip on reinstall, got: {skip_payload}")

    # --force backs up existing skill/command files (verify backup content).
    claude_skill.write_text("USER-SKILL-CONTENT\n", encoding="utf-8")
    claude_cmd.write_text("USER-COMMAND-CONTENT\n", encoding="utf-8")
    force_payload = json.loads(
        run_script(
            "install_exam_census_adapters.py",
            str(repo),
            "--platforms",
            "claude",
            "--force",
            "--json",
        )
    )
    if force_payload.get("installed") != 2:
        raise AssertionError(f"Expected force reinstall of 2 Claude files, got: {force_payload}")
    if "git_baseline" not in force_payload:
        raise AssertionError(f"Expected git_baseline in installer JSON, got: {force_payload}")
    skill_bak = claude_skill.with_suffix(claude_skill.suffix + ".bak")
    cmd_bak = claude_cmd.with_suffix(claude_cmd.suffix + ".bak")
    ensure_exists(skill_bak)
    ensure_exists(cmd_bak)
    ensure_contains(skill_bak, "USER-SKILL-CONTENT")
    ensure_contains(cmd_bak, "USER-COMMAND-CONTENT")
    ensure_contains(claude_skill, "name: exam-census")
    ensure_contains(claude_cmd, "/exam-census")
    if "USER-COMMAND-CONTENT" in claude_cmd.read_text(encoding="utf-8"):
        raise AssertionError("Force reinstall must restore command template, not keep sentinel")

    if not adapter_mod.find_dangerous_control_chars("safe\rhidden"):
        raise AssertionError("Bare CR must be reported as dangerous")
    if adapter_mod.find_dangerous_control_chars("safe\r\nok"):
        raise AssertionError("CRLF-normalized text must not report bare CR")

    # Experimental workflow JS only with explicit flag.
    exp_payload = json.loads(
        run_script(
            "install_exam_census_adapters.py",
            str(repo),
            "--platforms",
            "claude",
            "--include-experimental-claude-workflow",
            "--json",
        )
    )
    # skill+command skipped, workflow newly installed → installed=1, skipped=2
    if exp_payload.get("installed") != 1 or exp_payload.get("skipped") != 2:
        raise AssertionError(
            f"Expected experimental workflow install only, got: {exp_payload}"
        )
    ensure_exists(claude_wf)
    ensure_contains(claude_wf, "name: 'exam-census'")
    ensure_contains(claude_wf, "export const meta")
    ensure_contains(claude_wf, "EXPERIMENTAL")

    # Default reinstall without experimental flag must retire legacy workflow JS.
    retire_payload = json.loads(
        run_script(
            "install_exam_census_adapters.py",
            str(repo),
            "--platforms",
            "claude",
            "--json",
        )
    )
    if retire_payload.get("retired") != 1:
        raise AssertionError(
            f"Expected legacy Claude workflow retirement, got: {retire_payload}"
        )
    if claude_wf.exists():
        raise AssertionError("Legacy workflow JS must be retired after default reinstall")
    ensure_exists(claude_wf.with_suffix(claude_wf.suffix + ".bak"))

    # Symlink/junction escape: pre-planted .claude pointing outside the vault must be refused.
    outside = repo.parent / "adapter-escape-target"
    outside.mkdir(parents=True, exist_ok=True)
    escape_vault = repo.parent / "adapter-symlink-vault"
    escape_vault.mkdir(parents=True, exist_ok=True)
    claude_link = escape_vault / ".claude"
    try:
        claude_link.symlink_to(outside, target_is_directory=True)
    except OSError:
        # Windows without SeCreateSymbolicLinkPrivilege: directory junction still escapes.
        junction = subprocess.run(
            ["cmd", "/c", "mklink", "/J", str(claude_link), str(outside)],
            check=False,
            capture_output=True,
            text=True,
            encoding="utf-8",
        )
        if junction.returncode != 0 or not claude_link.exists():
            raise AssertionError(
                f"Unable to create symlink or junction for adapter escape smoke: {junction.stderr or junction.stdout}"
            )
    symlink_out = run_script_failure(
        "install_exam_census_adapters.py",
        str(escape_vault),
        "--platforms",
        "claude",
        "--json",
    )
    symlink_payload = json.loads(symlink_out)
    if symlink_payload.get("errors", 0) < 1:
        raise AssertionError(
            f"Expected symlink destination refusal, got: {symlink_payload}"
        )
    symlink_claude = next(
        item for item in symlink_payload["results"] if item["platform"] == "claude"
    )
    if symlink_claude.get("status") != "error":
        raise AssertionError(f"Expected claude platform error, got: {symlink_claude}")
    err_texts = " ".join(
        str(file_item.get("error", "")).lower() for file_item in symlink_claude["files"]
    )
    if (
        "symlink" not in err_texts
        and "outside vault" not in err_texts
        and "escaping vault" not in err_texts
    ):
        raise AssertionError(
            f"Expected symlink/escape error for claude adapter, got: {symlink_claude}"
        )
    if any(outside.iterdir()):
        raise AssertionError("Installer must not write through a symlinked .claude directory")


def exercise_feedback_lifecycle(repo: Path) -> None:
    feedback_day = date.today()
    expected_feedback_id = f'fb-{feedback_day.strftime("%Y%m%d")}-weekly-plan-omitted-imported-deadline-2'
    raw_path = Path(
        run_script(
            "log_feedback.py",
            str(repo),
            "--title",
            "Weekly plan omitted imported deadline",
            "--feedback-kind",
            "workflow",
            "--severity",
            "high",
            "--reproducibility",
            "always",
            "--source-context",
            "plan-week request after adding a linked homework task",
            "--related-artifacts",
            "tasks/weekly/current.md,tasks/deadlines/sample.md",
            "--related-roles",
            "coordinator,planning-assistant",
            "--what-happened",
            "- The generated weekly plan skipped a deadline that already existed in tasks/deadlines/.",
            "--expected-behavior",
            "- Weekly plans should include near-term deadline tasks alongside inbox and review items.",
            "--why-unsatisfying",
            "- The user still has to manually inspect deadlines after asking for a full weekly plan.",
            "--likely-cause",
            "- The planning workflow did not include imported deadline artifacts in the scan window.",
            "--suggested-improvement",
            "- Expand planning scans to include imported deadline artifacts before sorting upcoming work.",
            "--developer-summary",
            "- Planning scans need a regression check for imported or linked deadline artifacts.",
            "--evidence",
            "- tasks/weekly/current.md and tasks/deadlines/sample.md",
            "--follow-up",
            "Triage into the next planning iteration.",
        )
    )
    triaged_path = Path(
        run_script(
            "triage_feedback.py",
            str(repo),
            str(raw_path),
            "--feedback-kind",
            "workflow",
            "--severity",
            "high",
            "--reproducibility",
            "always",
            "--triage-status",
            "queued-for-planning",
            "--follow-up",
            "Bundle with the next planning workflow fix.",
            "--triage-notes",
            "- Confirmed the issue belongs to planning ingestion rather than task creation.",
        )
    )
    triaged_text = triaged_path.read_text(encoding="utf-8")
    triaged_text = triaged_text.replace('github_issue_url: ""', 'github_issue_url: "https://github.com/Gu-Heping/college-student-workflow/issues/42"')
    triaged_text = triaged_text.replace('github_issue_number: ""', 'github_issue_number: "42"')
    triaged_text = triaged_text.replace('github_issue_status: ""', 'github_issue_status: "open"')
    triaged_text = triaged_text.replace('reported_to_github_at: ""', 'reported_to_github_at: "2026-07-16"')
    triaged_path.write_text(triaged_text, encoding="utf-8", newline="\n")
    resolved_path = Path(
        run_script(
            "resolve_feedback.py",
            str(repo),
            str(triaged_path),
            "--resolution-summary",
            "- Added planning regression coverage so imported deadlines appear in weekly plan summaries.",
            "--fix-version",
            "0.7.0",
            "--changelog-note",
            "- Weekly plans now preserve imported and linked deadline tasks in the upcoming-work scan.",
            "--follow-up",
            "Verify once against a real imported-deadline workflow.",
        )
    )
    second_raw_path = Path(
        run_script(
            "log_feedback.py",
            str(repo),
            "--title",
            "Weekly plan omitted imported deadline",
            "--feedback-kind",
            "workflow",
            "--severity",
            "medium",
            "--reproducibility",
            "sometimes",
            "--source-context",
            "second collision test",
            "--developer-summary",
            "- This duplicate entry should keep a unique filename and feedback_id.",
        )
    )
    second_triaged_path = Path(
        run_script(
            "triage_feedback.py",
            str(repo),
            str(second_raw_path),
            "--triage-status",
            "needs-dedup-review",
            "--triage-notes",
            "- Collision handling preserved the older resolved item.",
        )
    )
    quoted_text = second_triaged_path.read_text(encoding="utf-8")
    quoted_text = quoted_text.replace("status: triaged", 'status: "triaged"')
    quoted_text = quoted_text.replace("severity: medium", 'severity: "medium"')
    second_triaged_path.write_text(quoted_text, encoding="utf-8", newline="\n")
    outside_path = repo.parent / "outside-feedback.md"
    outside_path.write_text("---\nstatus: open\n---\n", encoding="utf-8", newline="\n")
    failure_output = run_script_failure("triage_feedback.py", str(repo), str(outside_path))
    summary_path = Path(
        run_script(
            "summarize_feedback.py",
            str(repo),
            "--title",
            "Developer feedback handoff",
            "--scope",
            "smoke-test",
            "--audience",
            "developer",
        )
    )

    ensure_exists(resolved_path)
    ensure_contains(resolved_path, "feedback_id:")
    ensure_contains(resolved_path, 'github_issue_url: "https://github.com/Gu-Heping/college-student-workflow/issues/42"')
    ensure_contains(resolved_path, 'github_issue_number: "42"')
    ensure_contains(resolved_path, 'github_issue_status: "open"')
    ensure_contains(resolved_path, 'reported_to_github_at: "2026-07-16"')
    ensure_contains(resolved_path, "## Triage Notes")
    ensure_contains(resolved_path, "## Resolution Summary")
    ensure_contains(second_triaged_path, f'feedback_id: "{expected_feedback_id}"')
    ensure_contains(summary_path, "## Developer Handoff")
    ensure_contains(summary_path, "0.7.0")
    ensure_contains(summary_path, "- Triaged items: 1")
    ensure_contains(summary_path, expected_feedback_id)
    if "must stay under" not in failure_output:
        raise AssertionError(f"Expected path-guard failure, got: {failure_output}")

    archive_feedback = Path(
        run_script(
            "log_feedback.py",
            str(repo),
            "--title",
            "Archive should not fake GitHub status",
            "--feedback-kind",
            "docs",
            "--severity",
            "low",
            "--developer-summary",
            "- Archive status should preserve any known GitHub issue state.",
        )
    )
    archived_text = archive_feedback.read_text(encoding="utf-8")
    archived_text = archived_text.replace('github_issue_url: ""', 'github_issue_url: "https://github.com/Gu-Heping/college-student-workflow/issues/77"')
    archived_text = archived_text.replace('github_issue_number: ""', 'github_issue_number: "77"')
    archived_text = archived_text.replace('github_issue_status: ""', 'github_issue_status: "open"')
    archive_feedback.write_text(archived_text, encoding="utf-8", newline="\n")
    archived_path = Path(
        run_script(
            "resolve_feedback.py",
            str(repo),
            str(archive_feedback),
            "--status",
            "archived",
            "--resolution-summary",
            "- Archived locally without changing the public issue state.",
        )
    )
    ensure_contains(archived_path, 'github_issue_status: "open"')

    github_issue_feedback = Path(
        run_script(
            "log_feedback.py",
            str(repo),
            "--title",
            r"Privacy check D:\vault\notes.md sk-proj-1234567890-ABCDEFGHIJKLMNOPQRST",
            "--feedback-kind",
            "install",
            "--severity",
            "high",
            "--source-context",
            r"Codex on Windows with installed version: D:\vault\private-course\version-secret.txt and /Users/alice/My Vault/notes.md",
            "--related-artifacts",
            r"D:\vault\private-course\notes.md,.env,/Users/alice/private-notes.md",
            "--related-roles",
            "feedback-operator,codex",
            "--what-happened",
            "- The installer exposed a private path from D:\\vault\\private-course\\notes.md, C:\\Users\\Alice\\My Vault\\notes.md and /Users/alice/private-notes.md.",
            "--expected-behavior",
            "- Public reports should redact private Windows paths and vault references.",
            "--evidence",
            "- D:\\vault\\private-course\\notes.md\n- C:\\Users\\Alice\\My Vault\\notes.md\n- /Users/alice/My Vault/notes.md\n- .env.local: DATABASE_URL=postgres://secret@example\n- password: \"correct horse battery staple\"\n- sk-proj-1234567890-ABCDEFGHIJKLMNOPQRST\n- github_pat_1234567890ABCDEFGHIJKLMNOP\n- eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJzdWIiOiIxMjM0NTY3ODkwIn0.signature\n- +1 555-123-4567",
        )
    )
    issue_payload = json.loads(
        run_script(
            "prepare_github_issue.py",
            str(repo),
            str(github_issue_feedback),
        )
    )
    if issue_payload["feedback_id"] == "":
        raise AssertionError("prepare_github_issue.py should include feedback_id in its JSON output")
    if issue_payload["labels"] != ["feedback", "feedback:install", "severity:high"]:
        raise AssertionError("prepare_github_issue.py should derive labels from feedback kind and severity")
    if "## Feedback ID" not in issue_payload["body"] or "## Privacy Check" not in issue_payload["body"]:
        raise AssertionError("prepare_github_issue.py should emit the expected issue body sections")
    if "## Completeness Check" not in issue_payload["body"] or "## Environment" not in issue_payload["body"]:
        raise AssertionError("prepare_github_issue.py should emit completeness and environment sections")
    for leaked_text in [
        r"D:\vault\private-course\notes.md",
        r"C:\Users\Alice\My Vault\notes.md",
        "/Users/alice/private-notes.md",
        "/Users/alice/My Vault/notes.md",
        "sk-proj-1234567890-ABCDEFGHIJKLMNOPQRST",
        "github_pat_1234567890ABCDEFGHIJKLMNOP",
        "version-secret.txt",
        "DATABASE_URL=postgres://secret@example",
        "correct horse battery staple",
        "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJzdWIiOiIxMjM0NTY3ODkwIn0.signature",
        "+1 555-123-4567",
        "d-vault-private-course",
        "my-vault-notes",
    ]:
        if leaked_text in issue_payload["body"]:
            raise AssertionError(f"prepare_github_issue.py should redact sensitive text from public issue bodies: {leaked_text}")
    for redacted_marker in [
        "[REDACTED_WINDOWS_PATH]",
        "[REDACTED_UNIX_PATH]",
        "[REDACTED_TOKEN]",
        "[REDACTED_JWT]",
        "[REDACTED_PHONE]",
        "[REDACTED_ENV_FILE]",
    ]:
        if redacted_marker not in issue_payload["body"]:
            raise AssertionError(f"prepare_github_issue.py should include redaction marker {redacted_marker}")
    joined_warnings = "\n".join(issue_payload["privacy_warnings"])
    for expected_warning in ["Windows absolute paths", "Unix-style absolute paths", ".env", "token-like strings", "private vault path"]:
        if expected_warning not in joined_warnings:
            raise AssertionError(f"Expected privacy warning containing {expected_warning!r}, got: {joined_warnings}")
    joined_blockers = "\n".join(issue_payload["privacy_blockers"])
    for expected_blocker in ["secret-like key/value", "JWT-like tokens", "phone-number-like"]:
        if expected_blocker not in joined_blockers:
            raise AssertionError(f"Expected privacy blocker containing {expected_blocker!r}, got: {joined_blockers}")
    if not issue_payload["sanitized"]:
        raise AssertionError("prepare_github_issue.py should mark sanitized drafts")
    if not isinstance(issue_payload["completeness_warnings"], list):
        raise AssertionError("prepare_github_issue.py should return completeness warnings as a list")
    if "fb-public-" not in issue_payload["title"] or "fb-public-" not in issue_payload["body"]:
        raise AssertionError("prepare_github_issue.py should use a neutral public feedback identifier")
    if "- Python: unknown" not in issue_payload["body"]:
        raise AssertionError("prepare_github_issue.py should leave Python version as unknown unless the feedback explicitly captured it")

    stdin_safe = run_script_with_stdin(
        "prepare_github_issue.py",
        "Repro notes for a generic workflow failure without secrets.\n",
        "--stdin",
    )
    if "Repro notes for a generic workflow failure without secrets." not in stdin_safe.stdout:
        raise AssertionError("prepare_github_issue.py --stdin should pass through safe issue drafts")
    if stdin_safe.stderr.strip():
        raise AssertionError("prepare_github_issue.py --stdin should stay quiet for safe drafts")

    stdin_warning_held = run_script_with_stdin(
        "prepare_github_issue.py",
        "Saw a failure while reading C:\\Users\\Alice\\notes.md during import.\n",
        "--stdin",
        check=False,
    )
    if stdin_warning_held.returncode == 0:
        raise AssertionError("prepare_github_issue.py --stdin should hold back warning-bearing drafts without --allow-privacy-warnings")
    if stdin_warning_held.stdout.strip():
        raise AssertionError("prepare_github_issue.py --stdin should not emit stdout when warnings are held back")
    if "WARN:" not in stdin_warning_held.stderr or "Windows absolute paths" not in stdin_warning_held.stderr:
        raise AssertionError("prepare_github_issue.py --stdin should emit WARN lines for privacy warnings")

    stdin_warning = run_script_with_stdin(
        "prepare_github_issue.py",
        "Saw a failure while reading C:\\Users\\Alice\\notes.md during import.\n",
        "--stdin",
        "--allow-privacy-warnings",
    )
    if "[REDACTED_WINDOWS_PATH]" not in stdin_warning.stdout:
        raise AssertionError("prepare_github_issue.py --stdin --allow-privacy-warnings should redact Windows paths in stdout")
    if "C:\\Users\\Alice\\notes.md" in stdin_warning.stdout:
        raise AssertionError("prepare_github_issue.py --stdin should not leave absolute Windows paths in stdout")
    if "WARN:" not in stdin_warning.stderr or "Windows absolute paths" not in stdin_warning.stderr:
        raise AssertionError("prepare_github_issue.py --stdin should emit WARN lines for privacy warnings")

    stdin_blocker = run_script_with_stdin(
        "prepare_github_issue.py",
        "password: \"correct horse battery staple\"\n",
        "--stdin",
        check=False,
    )
    if stdin_blocker.returncode == 0:
        raise AssertionError("prepare_github_issue.py --stdin should abort when privacy blockers are present")
    if "BLOCK:" not in stdin_blocker.stderr:
        raise AssertionError("prepare_github_issue.py --stdin should emit BLOCK lines for privacy blockers")
    if stdin_blocker.stdout.strip():
        raise AssertionError("prepare_github_issue.py --stdin should not emit sanitized stdout when blockers abort publishing")

    stdin_check = run_script_with_stdin(
        "prepare_github_issue.py",
        "password: \"correct horse battery staple\"\n",
        "--stdin",
        "--check-only",
        check=False,
    )
    if stdin_check.returncode == 0:
        raise AssertionError("prepare_github_issue.py --stdin --check-only should exit non-zero for blockers")
    if "BLOCK:" not in stdin_check.stderr:
        raise AssertionError("prepare_github_issue.py --stdin --check-only should report blockers on stderr")
    if stdin_check.stdout.strip():
        raise AssertionError("prepare_github_issue.py --stdin --check-only should not rewrite the draft")

    stdin_check_warning = run_script_with_stdin(
        "prepare_github_issue.py",
        "Saw a failure while reading C:\\Users\\Alice\\notes.md during import.\n",
        "--check-stdin",
        "--check-only",
        check=False,
    )
    if stdin_check_warning.returncode == 0:
        raise AssertionError("prepare_github_issue.py --check-stdin --check-only should exit non-zero for warnings")
    if "WARN:" not in stdin_check_warning.stderr:
        raise AssertionError("prepare_github_issue.py --check-stdin --check-only should report warnings on stderr")
    if stdin_check_warning.stdout.strip():
        raise AssertionError("prepare_github_issue.py --check-stdin --check-only should not rewrite the draft")

    check_stdin_safe = run_script_with_stdin(
        "prepare_github_issue.py",
        "Safe PR review body without secrets or local paths.\n",
        "--check-stdin",
    )
    if "Safe PR review body without secrets or local paths." not in check_stdin_safe.stdout:
        raise AssertionError("prepare_github_issue.py --check-stdin should emit sanitized safe text")

    with tempfile.TemporaryDirectory() as tmp:
        marker = Path(tmp) / "posted.txt"
        safe_post = run_script_with_stdin(
            "sanitize_and_post.py",
            "Safe review comment for the sanitize wrapper.\n",
            "--",
            sys.executable,
            "-c",
            "import sys; from pathlib import Path; Path(sys.argv[1]).write_text(sys.stdin.read(), encoding='utf-8')",
            str(marker),
        )
        if safe_post.returncode != 0:
            raise AssertionError("sanitize_and_post.py should invoke the follow-up command for safe text")
        if not marker.exists() or "Safe review comment for the sanitize wrapper." not in marker.read_text(encoding="utf-8"):
            raise AssertionError("sanitize_and_post.py should feed sanitized stdin to the follow-up command")

        held_marker = Path(tmp) / "should-not-exist.txt"
        held_post = run_script_with_stdin(
            "sanitize_and_post.py",
            "Saw a failure while reading C:\\Users\\Alice\\notes.md during import.\n",
            "--",
            sys.executable,
            "-c",
            "import sys; from pathlib import Path; Path(sys.argv[1]).write_text(sys.stdin.read(), encoding='utf-8')",
            str(held_marker),
            check=False,
        )
        if held_post.returncode == 0:
            raise AssertionError("sanitize_and_post.py should hold back warning-bearing drafts")
        if held_marker.exists():
            raise AssertionError("sanitize_and_post.py must not invoke gh/follow-up when privacy warnings are held back")
        if held_post.stdout.strip():
            raise AssertionError("sanitize_and_post.py should not emit stdout when holding back a draft")

    sanitize_check = run_script_with_stdin(
        "sanitize_and_post.py",
        "Saw a failure while reading C:\\Users\\Alice\\notes.md during import.\n",
        "--check",
        "--allow-privacy-warnings",
    )
    if "[REDACTED_WINDOWS_PATH]" not in sanitize_check.stdout:
        raise AssertionError("sanitize_and_post.py --check --allow-privacy-warnings should redact paths")

    for empty_input, empty_label in (("", "empty"), ("   \n\t  \n", "whitespace-only")):
        with tempfile.TemporaryDirectory() as empty_tmp:
            empty_marker = Path(empty_tmp) / "should-not-post.txt"
            empty_post = run_script_with_stdin(
                "sanitize_and_post.py",
                empty_input,
                "--",
                sys.executable,
                "-c",
                "import sys; from pathlib import Path; Path(sys.argv[1]).write_text(sys.stdin.read(), encoding='utf-8')",
                str(empty_marker),
                check=False,
            )
            if empty_post.returncode == 0:
                raise AssertionError(f"sanitize_and_post.py should reject {empty_label} drafts")
            if empty_marker.exists():
                raise AssertionError(f"sanitize_and_post.py must not invoke follow-up for {empty_label} drafts")
            if "empty after sanitization" not in empty_post.stderr:
                raise AssertionError(f"sanitize_and_post.py should explain rejection of {empty_label} drafts")

    stdin_json = run_script_with_stdin(
        "prepare_github_issue.py",
        json.dumps({"body": "Contact me at alice@example.com about the vault at /Users/alice/My Vault/notes.md\n"}),
        "--stdin",
        "--stdin-format",
        "json",
        "--allow-privacy-warnings",
    )
    if "[REDACTED_EMAIL]" not in stdin_json.stdout or "[REDACTED_UNIX_PATH]" not in stdin_json.stdout:
        raise AssertionError("prepare_github_issue.py --stdin-format json should sanitize the body field")
    if "alice@example.com" in stdin_json.stdout or "/Users/alice/My Vault/notes.md" in stdin_json.stdout:
        raise AssertionError("prepare_github_issue.py --stdin-format json should not leak original sensitive body text")

    stdin_json_title_check = run_script_with_stdin(
        "prepare_github_issue.py",
        json.dumps({"title": "password: hunter2 leaked", "body": "Harmless body describing a workflow bug.\n"}),
        "--stdin",
        "--stdin-format",
        "json",
        "--check-only",
        check=False,
    )
    if stdin_json_title_check.returncode == 0:
        raise AssertionError("prepare_github_issue.py --stdin-format json --check-only should scan the title field for blockers")
    if "BLOCK:" not in stdin_json_title_check.stderr:
        raise AssertionError("prepare_github_issue.py --stdin-format json --check-only should report title blockers on stderr")

    stdin_json_title = run_script_with_stdin(
        "prepare_github_issue.py",
        json.dumps({"title": "Bug hit while reading C:\\Users\\Bob\\notes.md", "body": "Body without sensitive data.\n"}),
        "--stdin",
        "--stdin-format",
        "json",
        "--allow-privacy-warnings",
    )
    title_payload = json.loads(stdin_json_title.stdout)
    if "[REDACTED_WINDOWS_PATH]" not in title_payload["title"]:
        raise AssertionError("prepare_github_issue.py --stdin-format json should sanitize the title field")
    if "C:\\Users\\Bob\\notes.md" in stdin_json_title.stdout:
        raise AssertionError("prepare_github_issue.py --stdin-format json should not leak the original title path")
    if title_payload["body"].strip() != "Body without sensitive data.":
        raise AssertionError("prepare_github_issue.py --stdin-format json should preserve the sanitized body alongside the title")

    privacy_blocked_payload = json.loads(
        run_path_script(
            STUDENT_OS_SCRIPTS / "publish_github_issue.py",
            str(repo),
            str(github_issue_feedback),
            "--github-repo",
            "Gu-Heping/college-student-workflow",
            "--json",
            cwd=repo,
            env={"PATH": ""},
        )
    )
    if privacy_blocked_payload["published"]:
        raise AssertionError("publish_github_issue.py should not publish feedback with privacy warnings by default")
    if privacy_blocked_payload["blocked_reason"] != "privacy-blockers":
        raise AssertionError("publish_github_issue.py should report privacy-blockers when blocking sensitive data is present")
    if "--allow-privacy-warnings" in privacy_blocked_payload["next_step"]:
        raise AssertionError("publish_github_issue.py should not suggest overriding blocking sensitive data")
    if "gh_command" in privacy_blocked_payload:
        raise AssertionError("publish_github_issue.py should not emit a runnable gh command for blocker-only drafts")
    privacy_blocked_stdout = run_script(
        "publish_github_issue.py",
        str(repo),
        str(github_issue_feedback),
        "--github-repo",
        "Gu-Heping/college-student-workflow",
    )
    if "Publishing blocked due to privacy warnings." not in privacy_blocked_stdout:
        raise AssertionError("Non-JSON blocking output should still explain that publishing was blocked")
    if "gh issue create" in privacy_blocked_stdout:
        raise AssertionError("Non-JSON privacy blocking output should not print a ready-to-run publish command")

    publish_failure = json.loads(
        run_script(
            "publish_github_issue.py",
            str(repo),
            str(resolved_path),
            "--github-repo",
            "Gu-Heping/college-student-workflow",
            "--json",
        )
    )
    if publish_failure["blocked_reason"] != "already-linked":
        raise AssertionError("publish_github_issue.py should refuse duplicate publication for already-linked feedback")
    if publish_failure["existing_issue_number"] != "42":
        raise AssertionError("publish_github_issue.py should surface the existing linked issue metadata")

    manual_feedback = repo / "feedback" / "triaged" / "manual-path-feedback.md"
    manual_feedback.write_text(
        "\n".join(
            [
                "---",
                'type: "feedback"',
                'status: "triaged"',
                f'created: "{feedback_day.isoformat()}"',
                f'updated: "{feedback_day.isoformat()}"',
                'tags: ["feedback"]',
                'feedback_id: "../../notes/leak"',
                'feedback_kind: "other"',
                'severity: "medium"',
                'reproducibility: "sometimes"',
                'source_context: "manual import"',
                'related_course: ""',
                'related_artifacts: ""',
                'related_roles: "feedback-operator"',
                'github_issue_url: ""',
                'github_issue_number: ""',
                'github_issue_status: ""',
                'reported_to_github_at: ""',
                "---",
                "",
                "# Feedback - Manual path $(danger) `tick` test",
                "",
                "## What Happened",
                "",
                "- Fallback issue body generation should stay inside feedback/summaries.",
                "",
                "## Expected Behavior",
                "",
                "- Draft issue body should use a safe filename.",
                "",
                "## Why This Was Unsatisfying",
                "",
                "- Path traversal in feedback_id should not escape the summaries directory.",
                "",
                "## Likely Cause",
                "",
                "- The fallback body path trusted frontmatter directly.",
                "",
                "## Suggested Improvement",
                "",
                "- Normalize feedback IDs before using them in filenames.",
                "",
                "## Evidence",
                "",
                "- Manual path traversal test.",
                "",
                "## Follow-up",
                "",
                "- Confirm the generated draft body lives under feedback/summaries.",
                "",
            ]
        ) + "\n",
        encoding="utf-8",
        newline="\n",
    )
    publish_payload = json.loads(
        run_path_script(
            STUDENT_OS_SCRIPTS / "publish_github_issue.py",
            str(repo),
            str(manual_feedback),
            "--github-repo",
            "Gu-Heping/college-student-workflow",
            "--json",
            cwd=repo,
            env={"PATH": ""},
        )
    )
    if publish_payload["published"]:
        raise AssertionError("publish_github_issue.py fallback path should not publish when gh is unavailable")
    if publish_payload["blocked_reason"] != "gh-unavailable":
        raise AssertionError("publish_github_issue.py should distinguish gh-unavailable fallback from privacy blocking")
    body_path = repo / publish_payload["body_path"]
    ensure_exists(body_path)
    try:
        body_path.relative_to(repo / "feedback" / "summaries")
    except ValueError as exc:
        raise AssertionError("Fallback GitHub issue body should stay inside feedback/summaries") from exc
    if body_path.name != "manual-path-feedback-github-issue-body.md":
        raise AssertionError("Fallback body file should use a unique source-based name when feedback_id is unsafe or empty")
    gh_command = publish_payload["gh_command"]
    if "../../notes/leak" in gh_command:
        raise AssertionError("Fallback gh command should not expose the raw feedback_id once a neutral public ID is used")
    if "fb-public-" not in gh_command:
        raise AssertionError("Fallback gh command should use the neutral public feedback identifier in issue titles")
    if "--label" in gh_command:
        raise AssertionError("Fallback gh command should omit labels when the repo label set cannot be verified")
    if publish_payload["omitted_labels"] != ["feedback", "feedback:other", "severity:medium"]:
        raise AssertionError("publish_github_issue.py should report labels omitted from fallback publication guidance")

    single_quote_feedback = repo / "feedback" / "triaged" / "single-quote-empty-issue-fields.md"
    single_quote_feedback.write_text(
        "\n".join(
            [
                "---",
                "type: 'feedback'",
                "status: 'triaged'",
                f"created: '{feedback_day.isoformat()}'",
                f"updated: '{feedback_day.isoformat()}'",
                "tags: ['feedback']",
                "feedback_id: 'single-quote-empty'",
                "feedback_kind: 'docs'",
                "severity: 'low'",
                "reproducibility: 'sometimes'",
                "source_context: 'single quote empty metadata'",
                "related_course: ''",
                "related_artifacts: ''",
                "related_roles: 'feedback-operator'",
                "github_issue_url: ''",
                "github_issue_number: ''",
                "github_issue_status: ''",
                "reported_to_github_at: ''",
                "---",
                "",
                "# Feedback - Single quote empty fields",
                "",
                "## What Happened",
                "",
                "- Empty quoted GitHub fields should not block first-time publish preparation.",
                "",
                "## Expected Behavior",
                "",
                "- The publish helper should treat '' as empty metadata.",
                "",
                "## Why This Was Unsatisfying",
                "",
                "- Migrated feedback can otherwise be misclassified as already linked.",
                "",
                "## Likely Cause",
                "",
                "- Duplicate detection only stripped double quotes.",
                "",
                "## Suggested Improvement",
                "",
                "- Normalize scalar issue metadata before duplicate checks.",
                "",
                "## Evidence",
                "",
                "- Manual migrated feedback fixture.",
                "",
                "## Follow-up",
                "",
                "- Confirm fallback preparation succeeds without gh.",
                "",
            ]
        ) + "\n",
        encoding="utf-8",
        newline="\n",
    )
    single_quote_payload = json.loads(
        run_path_script(
            STUDENT_OS_SCRIPTS / "publish_github_issue.py",
            str(repo),
            str(single_quote_feedback),
            "--github-repo",
            "Gu-Heping/college-student-workflow",
            "--json",
            cwd=repo,
            env={"PATH": ""},
        )
    )
    if single_quote_payload["published"]:
        raise AssertionError("Single-quoted empty GitHub issue fields should not force a publish path")
    if single_quote_payload["blocked_reason"] != "gh-unavailable":
        raise AssertionError("Single-quoted empty GitHub metadata should still allow draft preparation")


def verify_organize_reviews(repo: Path) -> None:
    """Smoke coverage for Issue #52: reviews/<scope> auto-archiving."""
    scope_dir = repo / "courses" / "linear-algebra" / "reviews" / "期中"
    scope_dir.mkdir(parents=True, exist_ok=True)
    (scope_dir / "2019-期中-A.pdf").write_text("pdf", encoding="utf-8", newline="\n")
    (scope_dir / "2019-期中-A.pdf.md").write_text("# 2019 期中 A\n", encoding="utf-8", newline="\n")
    (scope_dir / "2019-期中-A.raw.md").write_text("# raw\n", encoding="utf-8", newline="\n")
    (scope_dir / "2019-期中-A-repair-summary.md").write_text("# repair\n", encoding="utf-8", newline="\n")

    payload = json.loads(
        run_script(
            "organize_reviews.py",
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "期中",
        )
    )
    if payload.get("dry_run"):
        raise AssertionError(f"Expected non-dry-run, got: {payload}")
    moved_targets = [str(item.get("target") or "").replace("\\", "/") for item in (payload.get("moved") or [])]
    if not any("试卷/2019-期中-A.pdf" in target for target in moved_targets):
        raise AssertionError(f"Expected PDF in 试卷/, got: {moved_targets}")
    if not any("文本/2019-期中-A.pdf.md" in target for target in moved_targets):
        raise AssertionError(f"Expected pdf.md in 文本/, got: {moved_targets}")
    if not any("文本/2019-期中-A.raw.md" in target for target in moved_targets):
        raise AssertionError(f"Expected raw.md in 文本/, got: {moved_targets}")
    if not any("归档/2019-期中-A-repair-summary.md" in target for target in moved_targets):
        raise AssertionError(f"Expected repair summary in 归档/, got: {moved_targets}")

    ensure_exists(scope_dir / "试卷" / "2019-期中-A.pdf")
    ensure_exists(scope_dir / "文本" / "2019-期中-A.pdf.md")
    ensure_exists(scope_dir / "文本" / "2019-期中-A.raw.md")
    ensure_exists(scope_dir / "归档" / "2019-期中-A-repair-summary.md")
    ensure_exists(scope_dir / "README.md")

    readme_text = (scope_dir / "README.md").read_text(encoding="utf-8")
    if "2019-期中-A.pdf" not in readme_text:
        raise AssertionError("README should index archived files")
    if "2019-期中-A.pdf.md" not in readme_text:
        raise AssertionError("README should index archived sidecars")

    # Idempotency: re-running should report nothing to do.
    idempotent_payload = json.loads(
        run_script(
            "organize_reviews.py",
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "期中",
        )
    )
    if idempotent_payload.get("moved"):
        raise AssertionError(f"Expected idempotent run to move nothing, got: {idempotent_payload}")
    if (idempotent_payload.get("message") or "") != "already organized, nothing to do":
        raise AssertionError(f"Expected idempotent message, got: {idempotent_payload}")

    # Dry-run on a separate fixture should not modify filesystem.
    dry_scope = repo / "courses" / "linear-algebra" / "reviews" / "期末"
    dry_scope.mkdir(parents=True, exist_ok=True)
    (dry_scope / "2020-期末.pdf").write_text("pdf", encoding="utf-8", newline="\n")
    (dry_scope / "2020-期末.pdf.md").write_text("# 2020 期末\n", encoding="utf-8", newline="\n")

    dry_payload = json.loads(
        run_script(
            "organize_reviews.py",
            str(repo),
            "--course",
            "linear-algebra",
            "--exam-scope",
            "期末",
            "--dry-run",
        )
    )
    if not dry_payload.get("dry_run"):
        raise AssertionError(f"Expected dry_run flag, got: {dry_payload}")
    if (dry_scope / "试卷").exists() or (dry_scope / "文本").exists():
        raise AssertionError("Dry run should not create archive directories")


def verify_generated_index_filtering(repo: Path) -> None:
    previous_sys_path = list(sys.path)
    sys.path.insert(0, str(STUDENT_OS_SCRIPTS))
    try:
        rebuild_module = load_student_os_script_module("rebuild_indexes.py", "student_os_rebuild_indexes_filter_smoke")
    finally:
        sys.path = previous_sys_path

    for relative in (
        PurePosixPath(".student-os/index/recent-activity.md"),
        PureWindowsPath(r".student-os\index\recent-activity.md"),
    ):
        if not rebuild_module.is_generated_index_relative(relative):
            raise AssertionError(f"Generated index filter should match cross-platform path: {relative}")
    if rebuild_module.is_generated_index_relative(PurePosixPath("notes/.student-os/index.md")):
        raise AssertionError("Generated index filter should only match the repository .student-os/index directory")

    user_note = repo / "notes" / ".student-os" / "index.md"
    user_note.parent.mkdir(parents=True, exist_ok=True)
    user_note.write_text("user note\n", encoding="utf-8")
    run_script("rebuild_indexes.py", str(repo))
    user_note_rel = "notes/.student-os/index.md"

    recent_activity = repo / ".student-os" / "index" / "recent-activity.md"
    recent_text = recent_activity.read_text(encoding="utf-8")
    if ".student-os/index/" in recent_text:
        raise AssertionError("Recent activity index should not include generated .student-os/index markdown files")
    if user_note_rel not in recent_text:
        raise AssertionError("Recent activity index should retain user Markdown outside the root generated index directory")

    summary = run_script("summarize_activity.py", str(repo), "--days", "3650")
    if ".student-os/index/" in summary:
        raise AssertionError("Activity summary should not include generated .student-os/index markdown files")
    if user_note_rel not in summary:
        raise AssertionError("Activity summary should retain user Markdown outside the root generated index directory")


def build_single_semester(repo: Path, today: date) -> None:
    due_date = (today + timedelta(days=8)).isoformat()
    default_week_label = f"{today.isoformat()}-plus-7d"
    run_script("scaffold_repo.py", str(repo))
    run_script("scaffold_course.py", str(repo), "Linear Algebra")
    run_script("scaffold_homework.py", str(repo), "linear-algebra", "Worksheet A", "--due", due_date)
    seed_planning_inputs(repo, today)
    run_script("build_review_indexes.py", str(repo))
    exercise_feedback_lifecycle(repo)
    exercise_import_workflows(repo)
    exercise_exam_census(repo)
    verify_organize_reviews(repo)
    run_script("build_review_indexes.py", str(repo))
    run_script("build_week_plan.py", str(repo))
    ensure_contains(repo / "tasks" / "weekly" / f"{default_week_label}.md", "Linear Algebra Midterm")
    ensure_contains(repo / "tasks" / "weekly" / f"{default_week_label}.md", "courses/linear-algebra/dashboard.md")
    default_plan_text = (repo / "tasks" / "weekly" / f"{default_week_label}.md").read_text(encoding="utf-8")
    if "- `linear-algebra` -> prioritize Worksheet A" in default_plan_text:
        raise AssertionError("Course actions should not prioritize out-of-window homework in the default weekly plan")
    if "- `linear-algebra` -> prioritize Linear Algebra Midterm" in default_plan_text:
        raise AssertionError("Course actions should not prioritize out-of-window exam tasks in the default weekly plan")
    (repo / "tasks" / "weekly" / f"{default_week_label}.md").unlink()
    (repo / "dashboards" / "weekly" / f"{default_week_label}.md").unlink()
    run_script("build_week_plan.py", str(repo), "--days", "14")
    run_script("rebuild_indexes.py", str(repo))

    ensure_exists(repo / "courses" / "linear-algebra" / "index.md")
    ensure_exists(repo / "courses" / "linear-algebra" / "homework" / "worksheet-a.md")
    ensure_contains(repo / ".student-os" / "index" / "courses.md", "courses/linear-algebra")
    ensure_contains(repo / "tasks" / "weekly" / f"{today.isoformat()}-plus-14d.md", "## Overdue Carryover")
    ensure_contains(repo / "tasks" / "weekly" / f"{today.isoformat()}-plus-14d.md", "Linear Algebra Midterm")
    ensure_contains(repo / "tasks" / "weekly" / f"{today.isoformat()}-plus-14d.md", "Imported Materials To Curate")
    ensure_contains(repo / "tasks" / "weekly" / f"{today.isoformat()}-plus-14d.md", "courses/linear-algebra/references/outline-import.md")
    ensure_contains(repo / "tasks" / "weekly" / f"{today.isoformat()}-plus-14d.md", "courses/linear-algebra/dashboard.md")
    weekly_plan_text = (repo / "tasks" / "weekly" / f"{today.isoformat()}-plus-14d.md").read_text(encoding="utf-8")
    ensure_contains(
        repo / "tasks" / "weekly" / f"{today.isoformat()}-plus-14d.md",
        f"- `linear-algebra` -> prioritize Linear Algebra Overdue Reading ({(today - timedelta(days=2)).isoformat()})",
    )
    if "manual-repair-sample-repair-summary.md" in weekly_plan_text:
        raise AssertionError("Repair summaries should not appear in imported-material triage")
    dashboard_exam_line = f"{(today + timedelta(days=9)).isoformat()} :: courses/linear-algebra/dashboard.md"
    task_exam_line = f"{(today + timedelta(days=10)).isoformat()} :: Linear Algebra Midterm"
    exams_section = weekly_plan_text.split("## Exams And Countdowns", 1)[1].split("## Course Actions", 1)[0]
    if exams_section.index(task_exam_line) < exams_section.index(dashboard_exam_line):
        raise AssertionError("Exam countdowns should be sorted chronologically across tasks and dashboard signals")
    if "Archived Quiz" in (repo / "tasks" / "weekly" / f"{today.isoformat()}-plus-14d.md").read_text(encoding="utf-8"):
        raise AssertionError("Archived tasks should not be listed in the weekly plan")
    ensure_contains(repo / "dashboards" / "weekly" / f"{today.isoformat()}-plus-14d.md", "Imported materials to review")
    ensure_contains(repo / ".student-os" / "index" / "dashboards.md", f"dashboards/weekly/{today.isoformat()}-plus-14d.md")
    verify_generated_index_filtering(repo)


def build_repo_inside_weekly_parent(repo: Path, today: date) -> None:
    run_script("scaffold_repo.py", str(repo))
    run_script("scaffold_course.py", str(repo), "Signals")
    write_task_fixture(
        repo / "tasks" / "deadlines" / "signals-quiz.md",
        title="Signals Quiz",
        due=(today + timedelta(days=3)).isoformat(),
        area="exam",
        priority="high",
        course="Signals",
        tags="[task, exam]",
    )
    run_script("build_week_plan.py", str(repo))
    ensure_contains(repo / "tasks" / "weekly" / f"{today.isoformat()}-plus-7d.md", "Signals Quiz")


def build_multi_semester(repo: Path, today: date) -> None:
    due_date = (today + timedelta(days=7)).isoformat()
    run_script("scaffold_repo.py", str(repo))
    run_script("scaffold_course.py", str(repo), "CS 101", "--semester", "2026 Fall")
    run_script("scaffold_course.py", str(repo), "Calculus II", "--semester", "2026 Fall")
    run_script("scaffold_course.py", str(repo), "CS 101", "--semester", "2027 Spring")
    run_script("scaffold_course.py", str(repo), "Data & Models", "--semester", "2026 Fall")
    run_script("scaffold_course.py", str(repo), "Data & Models", "--semester", "2027 Spring")
    run_script(
        "scaffold_homework.py",
        str(repo),
        "2026-fall/cs-101",
        "Problem Set 1",
        "--due",
        due_date,
    )
    write_task_fixture(
        repo / "tasks" / "deadlines" / "problem-set-1.md",
        title="Manual CS 101 checkpoint",
        due=due_date,
        area="homework",
        priority="medium",
        course="CS 101",
        course_link="../../courses/2026-fall/cs-101/index.md",
    )
    write_task_fixture(
        repo / "tasks" / "deadlines" / "2026-fall-data-models-reading.md",
        title="Data & Models reading",
        due=due_date,
        area="review",
        priority="medium",
        course="Data & Models",
    )
    write_task_fixture(
        repo / "tasks" / "deadlines" / "problem-set-2.md",
        title="Ambiguous CS 101 follow-up",
        due=due_date,
        area="homework",
        priority="medium",
        course="CS 101",
    )
    run_script("build_review_indexes.py", str(repo))
    run_script("build_week_plan.py", str(repo), "--days", "14")
    run_script("rebuild_indexes.py", str(repo))

    ensure_exists(repo / "semesters" / "2026-fall" / "overview.md")
    ensure_exists(repo / "tasks" / "deadlines" / "2026-fall-cs-101-problem-set-1.md")
    ensure_contains(repo / ".student-os" / "repo-profile.md", "enabled: true")
    ensure_contains(repo / "semesters" / "2026-fall" / "overview.md", "[CS 101]")
    ensure_contains(repo / ".student-os" / "index" / "courses.md", "courses/2026-fall/cs-101")
    ensure_contains(
        repo / "tasks" / "weekly" / f"{today.isoformat()}-plus-14d.md",
        "- `2026-fall/cs-101` -> prioritize CS 101 - Problem Set 1",
    )
    ensure_contains(
        repo / "tasks" / "weekly" / f"{today.isoformat()}-plus-14d.md",
        "- `2026-fall/data-models` -> prioritize Data & Models reading",
    )
    if "- `2027-spring/cs-101` -> prioritize CS 101 - Problem Set 1" in (
        repo / "tasks" / "weekly" / f"{today.isoformat()}-plus-14d.md"
    ).read_text(encoding="utf-8"):
        raise AssertionError("Duplicate course titles across semesters should not steal another semester's deadline priority")
    if "- `2027-spring/data-models` -> prioritize Data & Models reading" in (
        repo / "tasks" / "weekly" / f"{today.isoformat()}-plus-14d.md"
    ).read_text(encoding="utf-8"):
        raise AssertionError("Duplicate normalized titles should not steal another semester's title-based priority")
    if "Ambiguous CS 101 follow-up" in (
        repo / "tasks" / "weekly" / f"{today.isoformat()}-plus-14d.md"
    ).read_text(encoding="utf-8").split("## Course Actions", 1)[1].split("## Review Targets", 1)[0]:
        raise AssertionError("Ambiguous unscoped duplicate-course tasks should stay out of course-action prioritization")


def verify_chinese_slug_support(repo: Path, today: date) -> None:
    due_date = (today + timedelta(days=6)).isoformat()
    run_script("scaffold_repo.py", str(repo))
    course_path = Path(run_script("scaffold_course.py", str(repo), "模电", "--semester", "2026 春"))
    if course_path.relative_to(repo).as_posix() != "courses/2026-春/模电":
        raise AssertionError("Chinese semester and course names should produce stable unicode-aware course paths")
    run_script(
        "scaffold_homework.py",
        str(repo),
        "2026-春/模电",
        "第 1 次作业",
        "--due",
        due_date,
    )
    ensure_exists(repo / "courses" / "2026-春" / "模电" / "homework" / "第-1-次作业.md")
    ensure_exists(repo / "courses" / "2026-春" / "模电" / "homework" / "第-1-次作业-solution.md")
    ensure_exists(repo / "tasks" / "deadlines" / "2026-春-模电-第-1-次作业.md")
    ensure_contains(repo / ".student-os" / "repo-profile.md", "enabled: true")
    ensure_contains(repo / "semesters" / "2026-春" / "overview.md", "[模电]")
    run_script("build_week_plan.py", str(repo), "--days", "14")
    ensure_contains(
        repo / "tasks" / "weekly" / f"{today.isoformat()}-plus-14d.md",
        "- `2026-春/模电` -> prioritize 模电 - 第 1 次作业",
    )
    hindi_course_path = Path(
        run_script("scaffold_course.py", str(repo), "हिन्दी", "--semester", "2026 वसंत")
    )
    if hindi_course_path.relative_to(repo).as_posix() != "courses/2026-वसंत/हिन्दी":
        raise AssertionError("Unicode slugs should preserve combining marks for non-Latin scripts")


def build_legacy_layout(repo: Path, today: date) -> None:
    due_date = (today + timedelta(days=9)).isoformat()
    weekly_plan = repo / "tasks" / "weekly" / f"{today.isoformat()}-plus-14d.md"
    task_path = repo / "tasks" / "deadlines" / "legacy-course-legacy-sheet.md"
    run_script("scaffold_repo.py", str(repo))
    legacy_course = repo / "courses" / "legacy-course"
    (legacy_course / "homework").mkdir(parents=True, exist_ok=True)
    (legacy_course / "reviews").mkdir(parents=True, exist_ok=True)

    run_script("scaffold_homework.py", str(repo), "legacy-course", "Legacy Sheet", "--due", due_date)
    rewrite_legacy_task_link(task_path)
    run_script("build_review_indexes.py", str(repo))
    run_script("build_week_plan.py", str(repo), "--days", "14")

    ensure_exists(legacy_course / "homework" / "legacy-sheet.md")
    ensure_contains(repo / ".student-os" / "index" / "homework-and-reviews.md", "legacy-sheet.md")
    ensure_contains(weekly_plan, "legacy-course")
    ensure_contains(task_path, "legacy course folder without generated course home")


def verify_inspect_repo(repo: Path) -> None:
    payload = json.loads(run_script("inspect_repo.py", str(repo)))
    if "semesters" not in payload["canonical_dirs_present"]:
        raise AssertionError("inspect_repo.py did not report semesters as a canonical directory")
    if payload["dirty_files"]:
        raise AssertionError(f"Expected no dirty files in smoke-test repo, found: {payload['dirty_files']}")
    subprocess.run(["git", "-C", str(repo), "init"], check=True, capture_output=True, text=True)
    subprocess.run(["git", "-C", str(repo), "config", "user.name", "Smoke Test"], check=True, capture_output=True, text=True)
    subprocess.run(["git", "-C", str(repo), "config", "user.email", "smoke@example.com"], check=True, capture_output=True, text=True)
    subprocess.run(["git", "-C", str(repo), "add", "."], check=True, capture_output=True, text=True)
    subprocess.run(["git", "-C", str(repo), "commit", "-m", "baseline"], check=True, capture_output=True, text=True)
    (repo / "tasks" / "deadlines" / "linear-algebra.sync-conflict-20260710.md").write_text(
        "# conflict copy\n",
        encoding="utf-8",
        newline="\n",
    )
    (repo / ".obsidian").mkdir(parents=True, exist_ok=True)
    (repo / ".obsidian" / "workspace.json").write_text("{}", encoding="utf-8", newline="\n")
    (repo / "__pycache__").mkdir(parents=True, exist_ok=True)
    (repo / "__pycache__" / "inspect_repo.cpython-310.pyc").write_bytes(b"\0PYCCACHE")
    (repo / "tmp").mkdir(parents=True, exist_ok=True)
    (repo / "tmp" / "scratch.log").write_text("temporary scratch\n", encoding="utf-8", newline="\n")
    (repo / "references" / "textbooks").mkdir(parents=True, exist_ok=True)
    (repo / "references" / "textbooks" / "linear-algebra-textbook.pdf").write_bytes(b"%PDF-1.4\n")
    (repo / "references" / "slides").mkdir(parents=True, exist_ok=True)
    (repo / "references" / "slides" / "week-1-capture.mp4").write_bytes(b"fake-binary")
    (repo / "references" / "imports" / "raw").mkdir(parents=True, exist_ok=True)
    (repo / "references" / "imports" / "raw" / "ocr-dump.txt").write_text(
        "x" * (10 * 1024 * 1024 + 1),
        encoding="utf-8",
        newline="\n",
    )
    broken_link_supported = True
    broken_link = repo / "references" / "imports" / "raw" / "broken-link.txt"
    try:
        broken_link.symlink_to(repo / "references" / "imports" / "raw" / "missing-target.txt")
    except OSError:
        broken_link_supported = False
    staged_only_large = repo / "references" / "imports" / "raw" / "staged-large.txt"
    staged_only_large.write_text("x" * (10 * 1024 * 1024 + 2), encoding="utf-8", newline="\n")
    subprocess.run(["git", "-C", str(repo), "add", "references/imports/raw/staged-large.txt"], check=True, capture_output=True, text=True)
    staged_only_large.write_text("shrunk\n", encoding="utf-8", newline="\n")

    payload = json.loads(run_script("inspect_repo.py", str(repo)))
    if broken_link_supported and "references/imports/raw/broken-link.txt" in {item["path"] for item in payload["large_files"]}:
        raise AssertionError("inspect_repo.py should skip unreadable broken symlinks instead of treating them as large files")
    if "tasks/deadlines/linear-algebra.sync-conflict-20260710.md" not in payload["conflict_files"]:
        raise AssertionError("inspect_repo.py should report sync-conflict copies")
    if ".obsidian/workspace.json" not in payload["local_only_files"]:
        raise AssertionError("inspect_repo.py should report local workspace files")
    if "__pycache__/inspect_repo.cpython-310.pyc" not in payload["generated_cache_files"]:
        raise AssertionError("inspect_repo.py should report generated cache files")
    if "tmp/scratch.log" not in payload["temp_files"]:
        raise AssertionError("inspect_repo.py should report files under tmp/temp paths")
    if "references/textbooks/linear-algebra-textbook.pdf" not in payload["binary_files"]:
        raise AssertionError("inspect_repo.py should report binary source documents")
    if "references/slides/week-1-capture.mp4" not in payload["binary_files"]:
        raise AssertionError("inspect_repo.py should report binary media files")
    large_file_paths = {item["path"] for item in payload["large_files"]}
    if "references/imports/raw/ocr-dump.txt" not in large_file_paths:
        raise AssertionError("inspect_repo.py should report files above the large-file threshold")
    if "references/imports/raw/staged-large.txt" not in large_file_paths:
        raise AssertionError("inspect_repo.py should report oversized staged blobs even after the worktree copy shrinks")
    binary_zone_paths = {zone["path"] for zone in payload["binary_zones"]}
    if "references/textbooks" not in binary_zone_paths or "references/slides" not in binary_zone_paths:
        raise AssertionError("inspect_repo.py should summarize binary-heavy repository areas")
    warnings = payload.get("hygiene_warnings", [])
    for snippet in [
        "sync-conflict files detected",
        "generated caches detected",
        "local-only workspace or environment files detected",
        "temporary files detected under tmp/temp paths",
        "large files detected",
        "binary-heavy areas detected",
    ]:
        if not any(snippet in warning for warning in warnings):
            raise AssertionError(f"inspect_repo.py should emit a hygiene warning containing: {snippet}")


def verify_git_grouping(repo: Path, today: date) -> None:
    repo.mkdir(parents=True, exist_ok=True)
    subprocess.run(["git", "-C", str(repo), "init"], check=True, capture_output=True, text=True)
    subprocess.run(["git", "-C", str(repo), "config", "user.name", "Smoke Test"], check=True, capture_output=True, text=True)
    subprocess.run(["git", "-C", str(repo), "config", "user.email", "smoke@example.com"], check=True, capture_output=True, text=True)
    run_script("scaffold_repo.py", str(repo))
    run_script("scaffold_course.py", str(repo), "Linear Algebra")
    (repo / "references" / "slides").mkdir(parents=True, exist_ok=True)
    (repo / "references" / "slides" / "old-capture.mp4").write_bytes(b"tracked-binary")
    (repo / ".env.shared").write_text("TRACKED_SECRET=1\n", encoding="utf-8", newline="\n")
    (repo / ".env.tracked").write_text("TRACKED_SECRET=baseline\n", encoding="utf-8", newline="\n")
    nested_env_note = repo / "courses" / "env" / "notes.md"
    nested_env_note.parent.mkdir(parents=True, exist_ok=True)
    nested_env_note.write_text("# Environment course note\n", encoding="utf-8", newline="\n")
    eol_only_note = repo / "notes" / "eol-only.md"
    eol_only_note.parent.mkdir(parents=True, exist_ok=True)
    eol_only_note.write_text("# EOL only\n\nSame content.\n", encoding="utf-8", newline="\n")
    glob_literal_note = repo / "notes" / "[abc].md"
    glob_literal_note.write_text("# Glob literal\n\nSame content.\n", encoding="utf-8", newline="\n")
    glob_matching_note = repo / "notes" / "a.md"
    glob_matching_note.write_text("# Glob matching\n\nBaseline content.\n", encoding="utf-8", newline="\n")
    subprocess.run(["git", "-C", str(repo), "add", "."], check=True, capture_output=True, text=True)
    subprocess.run(["git", "-C", str(repo), "commit", "-m", "baseline"], check=True, capture_output=True, text=True)

    write_task_fixture(
        repo / "tasks" / "deadlines" / "manual-study-block.md",
        title="Manual Study Block",
        due=(today + timedelta(days=5)).isoformat(),
        area="study",
        priority="medium",
        course="Linear Algebra",
    )
    write_task_fixture(
        repo / "tasks" / "deadlines" / "模电-第1次作业.md",
        title="模电 第1次作业",
        due=(today + timedelta(days=6)).isoformat(),
        area="homework",
        priority="high",
        course="模电",
        tags="[task, homework]",
    )
    (repo / "tmp").mkdir(parents=True, exist_ok=True)
    (repo / "tmp" / "scratch.log").write_text("temporary notes\n", encoding="utf-8", newline="\n")
    (repo / "references" / "slides").mkdir(parents=True, exist_ok=True)
    (repo / "references" / "slides" / "lecture-capture.mp4").write_bytes(b"fake-binary")
    (repo / "references" / "textbooks").mkdir(parents=True, exist_ok=True)
    (repo / "references" / "textbooks" / "linear-algebra-textbook.pdf").write_bytes(b"%PDF-1.4\n")
    (repo / "references" / "imports" / "raw").mkdir(parents=True, exist_ok=True)
    (repo / "references" / "imports" / "raw" / "lecture-slides.pptx").write_bytes(b"PK\x03\x04")
    (repo / "references" / "imports" / "raw" / "ocr-dump.txt").write_text(
        "x" * (10 * 1024 * 1024 + 1),
        encoding="utf-8",
        newline="\n",
    )
    staged_only_large = repo / "references" / "imports" / "raw" / "staged-large.txt"
    staged_only_large.write_text("x" * (10 * 1024 * 1024 + 2), encoding="utf-8", newline="\n")
    conflict_path = repo / "tasks" / "deadlines" / "manual-study-block.sync-conflict-20260704.md"
    conflict_path.write_text("# conflict copy\n", encoding="utf-8", newline="\n")
    (repo / ".env").write_text("API_KEY=local\n", encoding="utf-8", newline="\n")
    (repo / ".env.local").write_text("API_KEY=override\n", encoding="utf-8", newline="\n")
    (repo / "env").mkdir(parents=True, exist_ok=True)
    (repo / "env" / "pyvenv.cfg").write_text("home = C:/Python\n", encoding="utf-8", newline="\n")
    (repo / "venv").mkdir(parents=True, exist_ok=True)
    (repo / "venv" / "pyvenv.cfg").write_text("home = C:/Python\n", encoding="utf-8", newline="\n")
    (repo / ".venv").mkdir(parents=True, exist_ok=True)
    (repo / ".venv" / "pyvenv.cfg").write_text("home = C:/Python\n", encoding="utf-8", newline="\n")
    nested_virtualenv = repo / "courses" / "project" / "env"
    nested_virtualenv.mkdir(parents=True, exist_ok=True)
    (nested_virtualenv / "pyvenv.cfg").write_text("home = C:/Python\n", encoding="utf-8", newline="\n")
    (nested_virtualenv / "lib64").mkdir(parents=True, exist_ok=True)
    (nested_virtualenv / "lib64" / "python3.12.txt").write_text("stdlib marker\n", encoding="utf-8", newline="\n")
    mixed_case_virtualenv = repo / "courses" / "Project" / "env"
    mixed_case_virtualenv.mkdir(parents=True, exist_ok=True)
    (mixed_case_virtualenv / "pyvenv.cfg").write_text("home = C:/Python\n", encoding="utf-8", newline="\n")
    (mixed_case_virtualenv / "lib64").mkdir(parents=True, exist_ok=True)
    (mixed_case_virtualenv / "lib64" / "python3.12.txt").write_text("stdlib marker\n", encoding="utf-8", newline="\n")
    (repo / "references" / "slides" / "old-capture.mp4").unlink()
    subprocess.run(
        ["git", "-C", str(repo), "mv", ".env.shared", "tasks/env-note.md"],
        check=True,
        capture_output=True,
        text=True,
    )
    nested_env_note.write_text("# Environment course note\n\nUpdated for grouping test.\n", encoding="utf-8", newline="\n")
    eol_only_note.write_text("# EOL only\r\n\r\nSame content.\r\n", encoding="utf-8", newline="")
    glob_literal_note.write_text("# Glob literal\r\n\r\nSame content.\r\n", encoding="utf-8", newline="")
    glob_matching_note.write_text("# Glob matching\n\nContent changed.\n", encoding="utf-8", newline="\n")
    nested_venv_course_note = repo / "courses" / "venv" / "notes" / "week1.md"
    nested_venv_course_note.parent.mkdir(parents=True, exist_ok=True)
    nested_venv_course_note.write_text("# Venv Course Note\n\nThis is coursework, not a virtualenv.\n", encoding="utf-8", newline="\n")
    nested_env_scripts_note = repo / "courses" / "env" / "scripts" / "week1.md"
    nested_env_scripts_note.parent.mkdir(parents=True, exist_ok=True)
    nested_env_scripts_note.write_text("# Env Scripts Note\n\nThis is coursework, not a virtualenv.\n", encoding="utf-8", newline="\n")
    tmp_slug_note = repo / "courses" / "notmp" / "notes.md"
    tmp_slug_note.parent.mkdir(parents=True, exist_ok=True)
    tmp_slug_note.write_text("# Notmp Course Note\n\nThis should not be treated as a temp path.\n", encoding="utf-8", newline="\n")
    (repo / ".env.tracked").write_text("TRACKED_SECRET=staged-update\n", encoding="utf-8", newline="\n")
    subprocess.run(["git", "-C", str(repo), "add", ".env.tracked"], check=True, capture_output=True, text=True)
    subprocess.run(["git", "-C", str(repo), "add", "references/imports/raw/staged-large.txt"], check=True, capture_output=True, text=True)
    (repo / ".env.tracked").unlink()
    staged_only_large.write_text("shrunk\n", encoding="utf-8", newline="\n")
    gitignore = repo / ".gitignore"
    gitignore.write_text(gitignore.read_text(encoding="utf-8") + ".venv/\n", encoding="utf-8", newline="\n")

    raw_grouping_output = run_script("group_git_changes.py", str(repo))
    if "tasks/deadlines/模电-第1次作业.md" not in raw_grouping_output:
        raise AssertionError("group_git_changes.py CLI output should preserve readable Unicode paths")
    if "\\u6a21\\u7535" in raw_grouping_output:
        raise AssertionError("group_git_changes.py CLI output should not ASCII-escape Unicode task paths")
    payload = json.loads(raw_grouping_output)
    if not payload.get("is_git_repo"):
        raise AssertionError(f"Expected a git repository payload, got: {payload}")
    hold_back = set(payload["hold_back_files"])
    if "tmp/" not in hold_back:
        raise AssertionError("group_git_changes.py should hold back ignored temporary directories")
    if "references/slides/lecture-capture.mp4" not in hold_back:
        raise AssertionError("group_git_changes.py should hold back binary media assets by default")
    if "tasks/deadlines/manual-study-block.sync-conflict-20260704.md" not in hold_back:
        raise AssertionError("group_git_changes.py should hold back sync-conflict files")
    if ".env" not in hold_back:
        raise AssertionError("group_git_changes.py should hold back .env files")
    if ".env.local" not in hold_back:
        raise AssertionError("group_git_changes.py should hold back .env.* files")
    if "env/pyvenv.cfg" not in hold_back:
        raise AssertionError("group_git_changes.py should hold back env/ virtual environment files")
    if ".venv/" not in hold_back:
        raise AssertionError("group_git_changes.py should hold back ignored local virtual environments")
    if "venv/pyvenv.cfg" not in hold_back:
        raise AssertionError("group_git_changes.py should hold back unignored virtual environment files")
    if "courses/project/env/pyvenv.cfg" not in hold_back:
        raise AssertionError("group_git_changes.py should hold back nested env virtual environment files")
    if "courses/project/env/lib64/python3.12.txt" not in hold_back:
        raise AssertionError("group_git_changes.py should hold back nested env virtual environment lib64 files")
    if "references/textbooks/linear-algebra-textbook.pdf" not in hold_back:
        raise AssertionError("group_git_changes.py should hold back raw imported PDF source documents")
    if "references/imports/raw/lecture-slides.pptx" not in hold_back:
        raise AssertionError("group_git_changes.py should hold back raw imported office source documents")
    if "references/imports/raw/ocr-dump.txt" not in hold_back:
        raise AssertionError("group_git_changes.py should hold back oversized text exports by default")
    if "references/imports/raw/staged-large.txt" not in hold_back:
        raise AssertionError("group_git_changes.py should hold back oversized staged blobs even after the worktree copy shrinks")
    if ".env.shared -> tasks/env-note.md" not in hold_back:
        raise AssertionError("group_git_changes.py should hold back renames from environment files")
    if ".env.tracked" not in hold_back:
        raise AssertionError("group_git_changes.py should keep mixed delete statuses for env files on hold-back")

    grouped_tasks = payload["artifact_grouping"].get("tasks", [])
    if "tasks/deadlines/manual-study-block.md" not in grouped_tasks:
        raise AssertionError("group_git_changes.py should keep normal task artifacts in the tasks group")
    if "tasks/deadlines/模电-第1次作业.md" not in grouped_tasks:
        raise AssertionError("group_git_changes.py should preserve readable Unicode task paths in grouped output")
    grouped_imports = payload["artifact_grouping"].get("imports", [])
    if "references/slides/old-capture.mp4" not in grouped_imports:
        raise AssertionError("group_git_changes.py should keep tracked hold-back deletions in commit guidance")
    split_paths = {
        path
        for split in payload["recommended_commit_split"]
        for path in split["paths"]
    }
    if "notes/eol-only.md" not in payload.get("eol_only_files", []):
        raise AssertionError(f"group_git_changes.py should classify Markdown line-ending-only churn separately: {payload}")
    if "notes/[abc].md" not in payload.get("eol_only_files", []):
        raise AssertionError(f"group_git_changes.py should use literal pathspecs for glob-like filenames: {payload}")
    if "notes/eol-only.md" in split_paths:
        raise AssertionError("EOL-only files should not be included in normal commit split guidance")
    if "notes/[abc].md" in split_paths:
        raise AssertionError("Glob-like EOL-only files should not be included in normal commit split guidance")
    if "notes/a.md" not in split_paths:
        raise AssertionError("Literal EOL-only detection must not hide a similarly matching content change")
    for expected_path, message in [
        ("tasks/deadlines/模电-第1次作业.md", "group_git_changes.py should not escape or drop Unicode task paths"),
        ("courses/env/notes.md", "group_git_changes.py should not treat nested env course paths as virtual environments"),
        ("courses/venv/notes/week1.md", "group_git_changes.py should not treat course slugs named venv as virtual environments without evidence"),
        ("courses/env/scripts/week1.md", "group_git_changes.py should not treat coursework under env/scripts as a virtual environment without pyvenv evidence"),
        ("courses/notmp/notes.md", "group_git_changes.py should not treat path components that only end with tmp as temporary directories"),
    ]:
        if expected_path in hold_back:
            raise AssertionError(message)
        if expected_path not in split_paths:
            raise AssertionError(f"{message}; expected {expected_path} to remain in recommended commit guidance")
    unexpected = hold_back & split_paths
    if unexpected:
        raise AssertionError(f"Hold-back files should not be suggested for commit splits: {sorted(unexpected)}")

    reasons = payload.get("hold_back_reasons", {})
    if reasons.get("tmp/") != "temporary file":
        raise AssertionError("Expected a temporary-file reason for tmp/")
    if reasons.get(".env") != "environment file":
        raise AssertionError("Expected an environment-file reason for .env")
    if reasons.get(".env.local") != "environment file":
        raise AssertionError("Expected an environment-file reason for .env.local")
    if reasons.get("env/pyvenv.cfg") != "local virtual environment":
        raise AssertionError("Expected a local-virtual-environment reason for env/pyvenv.cfg")
    if reasons.get(".venv/") != "local virtual environment":
        raise AssertionError("Expected a local-virtual-environment reason for .venv/")
    if reasons.get("venv/pyvenv.cfg") != "local virtual environment":
        raise AssertionError("Expected a local-virtual-environment reason for venv/pyvenv.cfg")
    if reasons.get("courses/project/env/pyvenv.cfg") != "local virtual environment":
        raise AssertionError("Expected a local-virtual-environment reason for nested env/pyvenv.cfg")
    if reasons.get("courses/project/env/lib64/python3.12.txt") != "local virtual environment":
        raise AssertionError("Expected a local-virtual-environment reason for nested env/lib64 virtualenv files")
    if reasons.get("references/textbooks/linear-algebra-textbook.pdf") != "binary source document":
        raise AssertionError("Expected a binary-source-document reason for imported PDFs")
    if reasons.get("references/imports/raw/lecture-slides.pptx") != "binary source document":
        raise AssertionError("Expected a binary-source-document reason for imported office files")
    if reasons.get("references/imports/raw/ocr-dump.txt") != "large file":
        raise AssertionError("Expected a large-file reason for oversized text exports")
    if reasons.get("references/imports/raw/staged-large.txt") != "large file":
        raise AssertionError("Expected a large-file reason for oversized staged blobs")
    if reasons.get(".env.shared -> tasks/env-note.md") != "environment file":
        raise AssertionError("Expected an environment-file reason for renames from environment files")
    if reasons.get(".env.tracked") != "environment file":
        raise AssertionError("Expected an environment-file reason for mixed delete env states")
    if reasons.get("tasks/deadlines/manual-study-block.sync-conflict-20260704.md") != "sync-conflict file":
        raise AssertionError("Expected a sync-conflict reason for the conflict copy")

    group_git_changes = load_group_git_changes_module()
    if not group_git_changes.is_virtualenv_path(repo, "courses/Project/env/pyvenv.cfg"):
        raise AssertionError("group_git_changes.py should detect mixed-case nested env virtual environment pyvenv markers")
    if not group_git_changes.is_virtualenv_path(repo, "courses/Project/env/lib64/python3.12.txt"):
        raise AssertionError("group_git_changes.py should detect mixed-case nested env virtual environment lib64 files")
    if group_git_changes.hold_back_reason(repo, "courses/notmp/notes.md"):
        raise AssertionError("group_git_changes.py should not mark coursework paths like courses/notmp/notes.md as temporary files")
    if group_git_changes.hold_back_reason(repo, "courses/Project/env/lib64/python3.12.txt") != "local virtual environment":
        raise AssertionError("group_git_changes.py should preserve original path casing when classifying virtualenv artifacts from CLI paths")


def init_git_repo(repo: Path) -> None:
    repo.mkdir(parents=True, exist_ok=True)
    subprocess.run(["git", "-C", str(repo), "init", "-b", "main"], check=True, capture_output=True, text=True)
    subprocess.run(["git", "-C", str(repo), "config", "user.name", "Smoke Test"], check=True, capture_output=True, text=True)
    subprocess.run(["git", "-C", str(repo), "config", "user.email", "smoke@example.com"], check=True, capture_output=True, text=True)


def commit_all(repo: Path, message: str) -> str:
    subprocess.run(["git", "-C", str(repo), "add", "."], check=True, capture_output=True, text=True)
    subprocess.run(["git", "-C", str(repo), "commit", "-m", message], check=True, capture_output=True, text=True)
    result = subprocess.run(
        ["git", "-C", str(repo), "rev-parse", "HEAD"],
        check=True,
        capture_output=True,
        text=True,
        encoding="utf-8",
    )
    return result.stdout.strip()


def seed_fake_student_os_source(repo: Path, *, version_label: str) -> str:
    skill_root = repo / "student-os"
    scripts_root = skill_root / "scripts"
    templates_root = skill_root / "templates"
    references_root = skill_root / "references"
    scripts_root.mkdir(parents=True, exist_ok=True)
    templates_root.mkdir(parents=True, exist_ok=True)
    references_root.mkdir(parents=True, exist_ok=True)
    (skill_root / "SKILL.md").write_text(
        "---\n"
        "name: student-os\n"
        "description: fixture skill\n"
        "---\n\n"
        f"# Student OS Fixture {version_label}\n",
        encoding="utf-8",
        newline="\n",
    )
    (scripts_root / "fixture_script.py").write_text(
        "def main():\n"
        f"    return '{version_label}'\n",
        encoding="utf-8",
        newline="\n",
    )
    (templates_root / "fixture.md").write_text(
        f"# Fixture Template {version_label}\n",
        encoding="utf-8",
        newline="\n",
    )
    (references_root / "fixture.md").write_text(
        f"# Fixture Reference {version_label}\n",
        encoding="utf-8",
        newline="\n",
    )
    shutil.copy2(STUDENT_OS_SCRIPTS / "update_student_os.py", scripts_root / "update_student_os.py")
    shutil.copy2(STUDENT_OS_SCRIPTS / "update_student_os_impl.py", scripts_root / "update_student_os_impl.py")
    return commit_all(repo, f"fixture: {version_label}")


def build_copy_install_fixture(base_dir: Path, source_repo: Path, installed_commit: str) -> Path:
    install_module = load_root_script_module("install_student_os.py", "student_os_install_smoke")
    source_skill = source_repo / "student-os"
    target = base_dir / "installed-student-os"
    shutil.copytree(source_skill, target)
    overrides_dir = target / ".student-os-local-overrides"
    overrides_dir.mkdir(parents=True, exist_ok=True)
    (overrides_dir / "notes.md").write_text("# Local override\n", encoding="utf-8", newline="\n")
    (target / ".student-os-install.local.json").write_text('{"theme":"custom"}\n', encoding="utf-8", newline="\n")
    manifest = install_module.build_install_manifest(
        destination=target,
        agent="codex",
        scope="user",
        install_method="copied",
        used_symlink=False,
        source_repo=str(source_repo.resolve()),
        source_ref="main",
        installed_commit=installed_commit,
        linked_source_path="",
    )
    install_module.write_manifest(target, manifest)
    return target


def assert_no_pycache(root: Path) -> None:
    pycache_paths = sorted(str(path) for path in root.rglob("__pycache__"))
    if pycache_paths:
        raise AssertionError(f"Validation should not leave __pycache__ artifacts behind: {pycache_paths}")


def verify_scaffold_gitattributes(tmp_root: Path) -> None:
    fresh = tmp_root / "fresh"
    run_script("scaffold_repo.py", str(fresh))
    fresh_attrs = fresh / ".gitattributes"
    ensure_contains(fresh_attrs, "*.md text eol=lf")
    if b"\r\n" in fresh_attrs.read_bytes():
        raise AssertionError("scaffold_repo.py should write .gitattributes with LF newlines")
    fresh_profile = fresh / ".student-os" / "repo-profile.md"
    if b"\r\n" in fresh_profile.read_bytes():
        raise AssertionError("scaffold_repo.py should write repo-profile.md with LF newlines")

    existing = tmp_root / "existing"
    existing.mkdir(parents=True, exist_ok=True)
    attrs = existing / ".gitattributes"
    attrs.write_text("*.png binary", encoding="utf-8", newline="\n")
    run_script("scaffold_repo.py", str(existing))
    text = attrs.read_text(encoding="utf-8")
    if "*.png binary" not in text or "*.md text eol=lf" not in text:
        raise AssertionError(f"scaffold_repo.py should preserve existing .gitattributes and add Markdown LF rule:\n{text}")
    run_script("scaffold_repo.py", str(existing))
    if attrs.read_text(encoding="utf-8").count("*.md text eol=lf") != 1:
        raise AssertionError("scaffold_repo.py should not duplicate the Markdown LF rule")

    dirty = tmp_root / "dirty-baseline"
    init_git_repo(dirty)
    preexisting = dirty / "notes" / "preexisting.md"
    preexisting.parent.mkdir(parents=True, exist_ok=True)
    preexisting.write_text("# Preexisting\n", encoding="utf-8", newline="\n")
    dirty_gitignore = dirty / ".gitignore"
    dirty_gitignore.write_text("custom-rule\n", encoding="utf-8", newline="\n")
    commit_all(dirty, "baseline")
    preexisting.write_text("# Preexisting\n\nDirty before scaffold.\n", encoding="utf-8", newline="\n")
    dirty_gitignore.write_text("custom-rule\n*.tmp\n", encoding="utf-8", newline="\n")
    output = run_script("scaffold_repo.py", str(dirty))
    changed_lines = "\n".join(line for line in output.splitlines() if line.startswith("CHANGED "))
    if "preexisting.md" in changed_lines:
        raise AssertionError(f"scaffold_repo.py should not attribute pre-existing dirty files to scaffold changes:\n{output}")
    if "CHANGED ?? .gitattributes" not in output:
        raise AssertionError(f"scaffold_repo.py should report .gitattributes as a scaffold-created change:\n{output}")
    if ".gitignore" not in changed_lines:
        raise AssertionError(f"scaffold_repo.py should report managed files it modified even when already dirty:\n{output}")


def verify_ensure_frontmatter(tmp_root: Path) -> None:
    papers = tmp_root / "courses" / "linear-algebra" / "references" / "exams"
    papers.mkdir(parents=True, exist_ok=True)

    missing_name = "线性代数期中试卷.pdf.md"
    missing_path = papers / missing_name
    chinese_body = "# 线性代数期中试卷\n\n这是中文正文，用于确认 UTF-8 不乱码。\n"
    missing_path.write_text(chinese_body, encoding="utf-8", newline="\n")

    existing_path = papers / "已有frontmatter.pdf.md"
    existing_text = (
        "\n".join(
            [
                "---",
                "type: pdf-import-note",
                "course:",
                "status: active",
                "created:",
                "updated:",
                "tags: [import, pdf]",
                'source_file: "已有frontmatter.pdf"',
                "import_method: manual",
                "repair_status:",
                "derived_from_import:",
                "---",
                "",
                "# Keep me",
                "",
            ]
        )
        + "\n"
    )
    existing_path.write_text(existing_text, encoding="utf-8", newline="\n")

    repaired_path = papers / "repaired-sample.pdf.md"
    repaired_body = "# Repaired sample body\n"
    repaired_path.write_text(repaired_body, encoding="utf-8", newline="\n")
    (papers / "repaired-sample.pdf-repair-summary.md").write_text(
        "# Repair Summary\n\n- Example.\n",
        encoding="utf-8",
        newline="\n",
    )
    (papers / "repaired-sample.pdf.raw.md").write_text("# raw\n", encoding="utf-8", newline="\n")

    raw_only = papers / "orphan.raw.md"
    raw_only.write_text("# orphan raw without frontmatter\n", encoding="utf-8", newline="\n")

    missing_before = missing_path.read_bytes()
    existing_before = existing_path.read_bytes()
    repaired_before = repaired_path.read_bytes()

    dry_payload = json.loads(run_script("ensure_frontmatter.py", str(papers), "--dry-run"))
    for key in [
        "scanned",
        "would_update",
        "updated",
        "skipped_existing_frontmatter",
        "skipped_unsupported",
        "errors",
    ]:
        if key not in dry_payload:
            raise AssertionError(f"ensure_frontmatter JSON missing field {key!r}: {dry_payload}")
    if dry_payload["updated"]:
        raise AssertionError("dry-run must not report updated files")
    if missing_path.read_bytes() != missing_before:
        raise AssertionError("dry-run must not modify files missing frontmatter")
    if existing_path.read_bytes() != existing_before:
        raise AssertionError("dry-run must not modify files with existing frontmatter")
    if repaired_path.read_bytes() != repaired_before:
        raise AssertionError("dry-run must not modify repaired-sample.pdf.md")
    if dry_payload["scanned"] != 3:
        raise AssertionError(f"Expected 3 scanned .pdf.md files by default, got: {dry_payload}")
    if len(dry_payload["would_update"]) != 2:
        raise AssertionError(f"Expected 2 would_update paths, got: {dry_payload}")
    if dry_payload["skipped_existing_frontmatter"] != ["已有frontmatter.pdf.md"]:
        raise AssertionError(f"Expected existing frontmatter skip, got: {dry_payload}")
    if dry_payload["errors"]:
        raise AssertionError(f"dry-run should have no errors: {dry_payload}")

    apply_payload = json.loads(run_script("ensure_frontmatter.py", str(papers), "--apply"))
    if apply_payload["would_update"]:
        raise AssertionError("apply should leave would_update empty")
    if sorted(apply_payload["updated"]) != sorted(
        ["repaired-sample.pdf.md", "线性代数期中试卷.pdf.md"]
    ):
        raise AssertionError(f"Unexpected updated set: {apply_payload}")
    if apply_payload["skipped_existing_frontmatter"] != ["已有frontmatter.pdf.md"]:
        raise AssertionError(f"apply should still skip existing frontmatter: {apply_payload}")

    missing_after = missing_path.read_text(encoding="utf-8")
    if not missing_after.startswith("---\n"):
        raise AssertionError("applied frontmatter must start at file beginning")
    if 'source_file: "线性代数期中试卷.pdf"' not in missing_after:
        raise AssertionError(f"source_file inference failed:\n{missing_after}")
    if "这是中文正文，用于确认 UTF-8 不乱码。" not in missing_after:
        raise AssertionError("UTF-8 Chinese body was corrupted after apply")
    if not missing_after.endswith(chinese_body) and chinese_body not in missing_after:
        raise AssertionError("body content must be preserved after prepend")
    if 'course: "linear-algebra"' not in missing_after:
        raise AssertionError(f"course should be inferred from courses/<course>/ path:\n{missing_after}")
    if "verify_status: unverified" not in missing_after:
        raise AssertionError(f"ensure_frontmatter should mark imports unverified:\n{missing_after}")

    if existing_path.read_text(encoding="utf-8") != existing_text:
        raise AssertionError("existing frontmatter file content must remain unchanged")

    repaired_after = repaired_path.read_text(encoding="utf-8")
    if "repair_status: auto-repaired" not in repaired_after:
        raise AssertionError(f"repair summary should set repair_status: auto-repaired:\n{repaired_after}")
    if "verify_status: unverified" not in repaired_after:
        raise AssertionError(f"ensure_frontmatter should set verify_status: unverified:\n{repaired_after}")
    if "\\u7ebf" in missing_after:
        raise AssertionError(f"frontmatter should keep Chinese readable, not JSON-escaped:\n{missing_after}")
    if b"\r\n" in missing_path.read_bytes() or b"\r\n" in repaired_path.read_bytes():
        raise AssertionError("ensure_frontmatter should write Markdown with LF newlines")
    if "derived_from_import:" not in repaired_after or "repaired-sample.pdf.raw.md" not in repaired_after:
        raise AssertionError(f"derived_from_import should point at sibling .raw.md:\n{repaired_after}")
    if repaired_body not in repaired_after:
        raise AssertionError("repaired sample body must be preserved")

    # Second apply is a no-op: all .pdf.md now have frontmatter.
    noop_payload = json.loads(run_script("ensure_frontmatter.py", str(papers), "--apply"))
    if noop_payload["updated"] or noop_payload["would_update"]:
        raise AssertionError(f"second apply should update nothing: {noop_payload}")
    if sorted(noop_payload["skipped_existing_frontmatter"]) != sorted(
        ["repaired-sample.pdf.md", "已有frontmatter.pdf.md", "线性代数期中试卷.pdf.md"]
    ):
        raise AssertionError(f"second apply skip set unexpected: {noop_payload}")

    raw_payload = json.loads(
        run_script("ensure_frontmatter.py", str(papers), "--include-raw", "--dry-run")
    )
    if "orphan.raw.md" not in raw_payload["would_update"]:
        raise AssertionError(f"--include-raw should plan orphan.raw.md update: {raw_payload}")
    if raw_only.read_text(encoding="utf-8") != "# orphan raw without frontmatter\n":
        raise AssertionError("--include-raw dry-run must not modify .raw.md")

    raw_apply = json.loads(run_script("ensure_frontmatter.py", str(raw_only), "--include-raw", "--apply"))
    if raw_apply["updated"] != ["orphan.raw.md"]:
        raise AssertionError(f"single-file --include-raw apply failed: {raw_apply}")
    raw_text = raw_only.read_text(encoding="utf-8")
    if "type: imported-reference" not in raw_text:
        raise AssertionError(f".raw.md should use imported-reference type:\n{raw_text}")

    wrapper_payload = json.loads(run_root_script("ensure_frontmatter.py", str(papers), "--dry-run"))
    if wrapper_payload["scanned"] != 3:
        raise AssertionError(f"root wrapper should match student-os script: {wrapper_payload}")


def verify_token_loader(tmp_root: Path) -> None:
    module = load_student_os_script_module("token_loader.py", "student_os_token_loader_smoke")
    install_module = load_root_script_module("install_student_os.py", "student_os_install_dotenv_smoke")
    update_module = load_student_os_script_module("update_student_os_impl.py", "student_os_update_dotenv_smoke")

    if ".env" not in install_module.LOCAL_OVERRIDE_NAMES:
        raise AssertionError("install_student_os.py should preserve skill-local .env across installs/updates")
    if ".env" not in update_module.LOCAL_OVERRIDE_NAMES:
        raise AssertionError("update_student_os_impl.py should preserve skill-local .env across updates")
    if not (ROOT / "student-os" / ".env.example").exists():
        raise AssertionError("student-os/.env.example should exist as the documented template")
    repo_gitignore = (ROOT / ".gitignore").read_text(encoding="utf-8")
    for needle in [".env", "student-os/.env", "!.env.example"]:
        if needle not in repo_gitignore:
            raise AssertionError(f"repository .gitignore should ignore skill secrets via {needle!r}")

    skill_root = tmp_root / "skill-root"
    cwd = tmp_root / "cwd"
    skill_root.mkdir(parents=True, exist_ok=True)
    cwd.mkdir(parents=True, exist_ok=True)
    (skill_root / ".env").write_text('MINERU_TOKEN="skill-token"\n# comment\nexport MINERU_API_TOKEN=unused\n', encoding="utf-8", newline="\n")
    (cwd / ".env").write_text("MINERU_TOKEN=cwd-token\n", encoding="utf-8", newline="\n")

    scaffold_repo = tmp_root / "scaffold-gitignore"
    run_script("scaffold_repo.py", str(scaffold_repo))
    gitignore_text = (scaffold_repo / ".gitignore").read_text(encoding="utf-8")
    for needle in [".env", "*.env", "!.env.example"]:
        if needle not in gitignore_text:
            raise AssertionError(f"scaffold_repo.py should ignore secrets via {needle!r}")

    existing_vault = tmp_root / "existing-vault"
    existing_vault.mkdir(parents=True, exist_ok=True)
    (existing_vault / ".gitignore").write_text("*.log\nnode_modules/\n", encoding="utf-8", newline="\n")
    run_script("scaffold_repo.py", str(existing_vault))
    merged_gitignore = (existing_vault / ".gitignore").read_text(encoding="utf-8")
    if "*.log" not in merged_gitignore or "node_modules/" not in merged_gitignore:
        raise AssertionError("scaffold_repo.py should preserve existing .gitignore entries")
    for needle in [".env", "*.env", "!.env.example"]:
        if needle not in merged_gitignore:
            raise AssertionError(f"scaffold_repo.py should append missing secret rule {needle!r} to existing .gitignore")

    env_backup = {name: os.environ.get(name) for name in ["MINERU_TOKEN", "MINERU_API_TOKEN"]}
    try:
        os.environ.pop("MINERU_TOKEN", None)
        os.environ.pop("MINERU_API_TOKEN", None)

        if module.load_token(" cli-token ", skill_root=skill_root, cwd=cwd) != "cli-token":
            raise AssertionError("CLI token should win over env and .env files")

        os.environ["MINERU_TOKEN"] = " process-env-token "
        if module.load_token(None, skill_root=skill_root, cwd=cwd) != "process-env-token":
            raise AssertionError("Process environment should win over .env files")
        os.environ.pop("MINERU_TOKEN", None)

        if module.load_token(None, skill_root=skill_root, cwd=cwd) != "skill-token":
            raise AssertionError("Skill-root .env should win over cwd .env")

        (skill_root / ".env").unlink()
        if module.load_token(None, skill_root=skill_root, cwd=cwd) != "cwd-token":
            raise AssertionError("cwd .env should be used when skill-root .env is absent")

        (cwd / ".env").write_text("MINERU_API_TOKEN=alias-token\n", encoding="utf-8", newline="\n")
        if module.load_token(None, skill_root=skill_root, cwd=cwd) != "alias-token":
            raise AssertionError("MINERU_API_TOKEN in .env should be accepted as an alias")

        (cwd / ".env").write_bytes(b"\xff\xfeMINERU_TOKEN=bad-encoding\n")
        if module.load_token(None, skill_root=skill_root, cwd=cwd) is not None:
            raise AssertionError("Unreadable/non-UTF-8 .env files should be skipped without raising")

        repair_only_md = cwd / "repair-only-local.md"
        repair_only_md.write_text(
            "\n".join(
                [
                    "---",
                    "type: imported-reference",
                    "repair_status:",
                    "derived_from_import:",
                    "---",
                    "",
                    "#Broken",
                    "",
                ]
            )
            + "\n",
            encoding="utf-8",
            newline="\n",
        )
        local_payload = json.loads(
            run_script(
                "materials_convert.py",
                str(repair_only_md),
                "--method",
                "local",
                "--repair-only",
                cwd=cwd,
                env={"MINERU_TOKEN": "", "MINERU_API_TOKEN": ""},
            )
        )
        if not local_payload.get("converted"):
            raise AssertionError(f"Local/repair-only conversion should ignore bad cwd .env, got: {local_payload}")
    finally:
        for name, value in env_backup.items():
            if value is None:
                os.environ.pop(name, None)
            else:
                os.environ[name] = value

    # Legacy entry-symlink layout: scripts/ points at the source checkout, but .env lives
    # in the installed skill root and must still be discovered.
    legacy_install = tmp_root / "legacy-entry-symlink-install"
    legacy_scripts = legacy_install / "scripts"
    legacy_scripts.mkdir(parents=True, exist_ok=True)
    try:
        for child in (ROOT / "student-os").iterdir():
            if child.name == "scripts":
                continue
            target = legacy_install / child.name
            if child.is_dir():
                target.symlink_to(child.resolve(), target_is_directory=True)
            else:
                target.symlink_to(child.resolve())
        # Replace real scripts dir with a symlink to the source scripts.
        shutil.rmtree(legacy_scripts)
        legacy_scripts.symlink_to((ROOT / "student-os" / "scripts").resolve(), target_is_directory=True)
    except OSError:
        return
    (legacy_install / ".env").write_text("MINERU_TOKEN=legacy-install-token\n", encoding="utf-8", newline="\n")
    loader_via_symlink = legacy_scripts / "token_loader.py"
    spec = importlib.util.spec_from_file_location("student_os_token_loader_symlink_smoke", loader_via_symlink)
    if spec is None or spec.loader is None:
        raise RuntimeError("Unable to load token_loader via legacy scripts symlink")
    symlink_module = importlib.util.module_from_spec(spec)
    original_flag = sys.dont_write_bytecode
    sys.dont_write_bytecode = True
    try:
        spec.loader.exec_module(symlink_module)
    finally:
        sys.dont_write_bytecode = original_flag
    lexical_root = symlink_module.skill_root_dir()
    if lexical_root.resolve() != legacy_install.resolve():
        # Compare using absolute() semantics: resolved source checkout should NOT win.
        if Path(lexical_root).absolute() != legacy_install.absolute():
            raise AssertionError(
                f"skill_root_dir() should keep the installed root for entry-symlink layouts, got: {lexical_root}"
            )
    env_backup = {name: os.environ.get(name) for name in ["MINERU_TOKEN", "MINERU_API_TOKEN"]}
    try:
        os.environ.pop("MINERU_TOKEN", None)
        os.environ.pop("MINERU_API_TOKEN", None)
        if symlink_module.load_token(None, cwd=tmp_root / "empty-cwd") != "legacy-install-token":
            raise AssertionError("Entry-symlink installs should read .env from the installed skill root")
    finally:
        for name, value in env_backup.items():
            if value is None:
                os.environ.pop(name, None)
            else:
                os.environ[name] = value


def verify_install_manifest_generation(tmp_root: Path) -> None:
    install_module = load_root_script_module("install_student_os.py", "student_os_install_manifest_smoke")
    codex_home = tmp_root / "codex-home"
    payload = json.loads(
        run_root_script(
            "install_student_os.py",
            "--agent",
            "codex",
            "--scope",
            "user",
            "--mode",
            "copy",
            "--json",
            env={"CODEX_HOME": str(codex_home)},
        )
    )
    result = payload["results"][0]
    destination = Path(result["destination"])
    manifest_path = destination / ".student-os-install.json"
    if not manifest_path.exists():
        raise AssertionError("install_student_os.py should write an install manifest into the installed skill")
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    for key in [
        "skill_name",
        "source_repo",
        "source_ref",
        "installed_commit",
        "installed_at",
        "install_method",
        "agent",
        "scope",
        "target_path",
        "used_symlink",
        "tracked_files",
    ]:
        if key not in manifest:
            raise AssertionError(f"install manifest should include {key}")
    if manifest["skill_name"] != "student-os":
        raise AssertionError("install manifest should record the skill name")
    if manifest["install_method"] != "copied":
        raise AssertionError("copy installs should record copied as the install method")
    if manifest["source_repo"] != install_module.discover_source_repo():
        raise AssertionError("install manifests should record the actual source repo by default")

    (destination / ".env").write_text("MINERU_TOKEN=preserve-me\n", encoding="utf-8", newline="\n")
    forced = json.loads(
        run_root_script(
            "install_student_os.py",
            "--agent",
            "codex",
            "--scope",
            "user",
            "--mode",
            "copy",
            "--force",
            "--json",
            env={"CODEX_HOME": str(codex_home)},
        )
    )
    forced_destination = Path(forced["results"][0]["destination"])
    preserved_env = forced_destination / ".env"
    if not preserved_env.exists():
        raise AssertionError("Forced copy reinstall should preserve skill-local .env")
    if "preserve-me" not in preserved_env.read_text(encoding="utf-8"):
        raise AssertionError("Forced copy reinstall should restore the previous .env contents")
    if ".env" not in forced["results"][0].get("preserved_overrides", ""):
        raise AssertionError("Forced reinstall payload should report preserved .env overrides")


def read_install_manifest(destination: Path) -> dict:
    manifest_path = destination / ".student-os-install.json"
    if not manifest_path.exists():
        raise AssertionError(f"Expected install manifest to exist: {manifest_path}")
    return json.loads(manifest_path.read_text(encoding="utf-8"))


def parse_frontmatter_scalar_fields(frontmatter: str) -> dict[str, object]:
    fields: dict[str, object] = {}
    for line in frontmatter.splitlines():
        if not line.strip() or line.lstrip().startswith("#"):
            continue
        key, separator, value = line.partition(":")
        if not separator:
            raise AssertionError(f"student-os/SKILL.md frontmatter line should be key: value YAML: {line!r}")
        key = key.strip()
        value = value.strip()
        if not re.fullmatch(r"[A-Za-z_][A-Za-z0-9_-]*", key):
            raise AssertionError(f"student-os/SKILL.md frontmatter key should be a plain YAML scalar: {key!r}")
        if value.startswith(("'", '"')):
            quote = value[0]
            if len(value) < 2 or not value.endswith(quote):
                raise AssertionError(f"student-os/SKILL.md frontmatter quoted value should be closed: {value!r}")
        elif value.endswith(("'", '"')):
            raise AssertionError(f"student-os/SKILL.md frontmatter quoted value should be opened: {value!r}")
        if value.startswith(("[", "{")) or value.endswith(("[", "{")):
            raise AssertionError(f"student-os/SKILL.md frontmatter value should be a plain YAML scalar: {value!r}")
        fields[key] = value.strip('"\'')
    return fields


def parse_dsh_skill_frontmatter(text: str) -> dict[str, object]:
    if not text.startswith("---\n"):
        raise AssertionError("student-os/SKILL.md must start with YAML frontmatter for DSH discovery")
    closing = text.find("\n---\n", 4)
    if closing < 0:
        raise AssertionError("student-os/SKILL.md must close YAML frontmatter before the body")
    frontmatter = text[4:closing]
    try:
        import yaml  # type: ignore
    except ImportError:
        yaml = None

    if yaml is not None:
        try:
            loaded = yaml.safe_load(frontmatter)
        except Exception as exc:
            raise AssertionError(f"student-os/SKILL.md frontmatter must be valid YAML: {exc}") from exc
        if not isinstance(loaded, dict):
            raise AssertionError(f"student-os/SKILL.md frontmatter should parse to a mapping, got: {type(loaded).__name__}")
        return {str(key): value for key, value in loaded.items()}
    return parse_frontmatter_scalar_fields(frontmatter)


def verify_dsh_skill_frontmatter_contract() -> None:
    skill_path = ROOT / "student-os" / "SKILL.md"
    fields = parse_dsh_skill_frontmatter(skill_path.read_text(encoding="utf-8"))
    try:
        parse_dsh_skill_frontmatter("---\nname: student-os\ndescription: [unterminated\n---\n")
    except AssertionError:
        pass
    else:
        raise AssertionError("DSH frontmatter validation should reject invalid YAML")
    try:
        parse_dsh_skill_frontmatter('---\nname: student-os\ndescription: "unterminated\n---\n')
    except AssertionError:
        pass
    else:
        raise AssertionError("DSH frontmatter validation should reject unterminated quoted scalars")
    name = fields.get("name")
    description = fields.get("description")
    if name != "student-os":
        raise AssertionError(f"DSH skill name should be student-os, got: {name!r}")
    if not description:
        raise AssertionError("DSH skill description should be non-empty")
    if not re.fullmatch(r"[a-z0-9]+(?:-[a-z0-9]+)*", name):
        raise AssertionError(f"DSH skill name should be kebab-case, got: {name!r}")


def verify_dsh_install_paths(tmp_root: Path) -> None:
    dsh_home = tmp_root / "dsh-home"
    user_payload = json.loads(
        run_root_script(
            "install_student_os.py",
            "--agent",
            "dsh",
            "--scope",
            "user",
            "--mode",
            "copy",
            "--json",
            env={"DSH_HOME": str(dsh_home)},
        )
    )
    user_result = user_payload["results"][0]
    user_destination = dsh_home / "skills" / "student-os"
    if Path(user_result["destination"]) != user_destination.resolve():
        raise AssertionError(f"DSH user install should honor DSH_HOME, got: {user_result}")
    ensure_exists(user_destination / "SKILL.md")
    user_manifest = read_install_manifest(user_destination)
    if user_manifest["agent"] != "dsh" or user_manifest["scope"] != "user":
        raise AssertionError(f"DSH user manifest should record agent/scope, got: {user_manifest}")
    if Path(user_manifest["target_path"]) != user_destination.resolve():
        raise AssertionError("DSH user manifest should record the DSH_HOME destination")

    project_root = tmp_root / "project-root"
    project_root.mkdir(parents=True, exist_ok=True)
    project_payload = json.loads(
        run_root_script(
            "install_student_os.py",
            "--agent",
            "dsh",
            "--scope",
            "project",
            "--project-root",
            str(project_root),
            "--mode",
            "copy",
            "--json",
            env={"DSH_HOME": str(tmp_root / "unused-dsh-home")},
        )
    )
    project_result = project_payload["results"][0]
    project_destination = project_root / ".dsh" / "skills" / "student-os"
    if Path(project_result["destination"]) != project_destination.resolve():
        raise AssertionError(f"DSH project install should use .dsh/skills under --project-root, got: {project_result}")
    ensure_exists(project_destination / "SKILL.md")
    ensure_exists(project_destination / "scripts")
    ensure_exists(project_destination / "references")
    project_manifest = read_install_manifest(project_destination)
    if project_manifest["agent"] != "dsh" or project_manifest["scope"] != "project":
        raise AssertionError(f"DSH project manifest should record agent/scope, got: {project_manifest}")

    both_home = tmp_root / "both-dsh-home"
    both_project = tmp_root / "both-project-root"
    both_project.mkdir(parents=True, exist_ok=True)
    both_payload = json.loads(
        run_root_script(
            "install_student_os.py",
            "--agent",
            "dsh",
            "--scope",
            "both",
            "--project-root",
            str(both_project),
            "--mode",
            "copy",
            "--json",
            env={"DSH_HOME": str(both_home)},
        )
    )
    both_scopes = {result["scope"] for result in both_payload["results"]}
    if both_scopes != {"user", "project"}:
        raise AssertionError(f"DSH --scope both should install user and project scopes, got: {both_payload}")
    ensure_exists(both_home / "skills" / "student-os" / "SKILL.md")
    ensure_exists(both_project / ".dsh" / "skills" / "student-os" / "SKILL.md")

    overlap_project = tmp_root / "overlap-project-root"
    overlap_home = overlap_project / ".dsh"
    overlap_project.mkdir(parents=True, exist_ok=True)
    overlap_payload = json.loads(
        run_root_script(
            "install_student_os.py",
            "--agent",
            "dsh",
            "--scope",
            "both",
            "--project-root",
            str(overlap_project),
            "--mode",
            "copy",
            "--json",
            env={"DSH_HOME": str(overlap_home)},
        )
    )
    overlap_destination = overlap_project / ".dsh" / "skills" / "student-os"
    if len(overlap_payload["results"]) != 1:
        raise AssertionError(f"DSH overlapping user/project destinations should be deduped, got: {overlap_payload}")
    if Path(overlap_payload["results"][0]["destination"]) != overlap_destination.resolve():
        raise AssertionError(f"DSH deduped destination should be the shared .dsh skill path, got: {overlap_payload}")
    ensure_exists(overlap_destination / "SKILL.md")

    env_backup = {name: os.environ.get(name) for name in ["HOME", "USERPROFILE", "DSH_HOME"]}
    try:
        isolated_home = tmp_root / "isolated-home"
        os.environ["HOME"] = str(isolated_home)
        os.environ["USERPROFILE"] = str(isolated_home)
        os.environ.pop("DSH_HOME", None)
        install_module = load_root_script_module("install_student_os.py", "student_os_install_dsh_default_home_smoke")
    finally:
        for name, value in env_backup.items():
            if value is None:
                os.environ.pop(name, None)
            else:
                os.environ[name] = value
    expected_default = (isolated_home / ".dsh" / "skills").resolve()
    if install_module.AGENT_PATHS["dsh"]["user"].resolve() != expected_default:
        raise AssertionError("DSH user install should default to ~/.dsh/skills when DSH_HOME is unset")


def verify_dsh_home_resolution(tmp_root: Path) -> None:
    env_backup = {name: os.environ.get(name) for name in ["HOME", "USERPROFILE", "DSH_HOME"]}
    previous_cwd = Path.cwd()
    try:
        isolated_home = tmp_root / "isolated-home"
        isolated_home.mkdir(parents=True, exist_ok=True)
        os.environ["HOME"] = str(isolated_home)
        os.environ["USERPROFILE"] = str(isolated_home)

        for label, dsh_home_value in [("empty", ""), ("blank", "   ")]:
            os.environ["DSH_HOME"] = dsh_home_value
            os.chdir(tmp_root)
            install_module = load_root_script_module("install_student_os.py", f"student_os_install_dsh_{label}_home")
            update_module = load_student_os_script_module(
                "update_student_os_impl.py",
                f"student_os_update_dsh_{label}_home",
            )
            expected_home = isolated_home / ".dsh"
            if install_module.DSH_HOME.resolve() != expected_home.resolve():
                raise AssertionError(f"Installer should treat blank DSH_HOME as unset, got: {install_module.DSH_HOME}")
            if update_module.resolve_dsh_home().resolve() != expected_home.resolve():
                raise AssertionError("Updater should treat blank DSH_HOME as unset")
            if install_module.DSH_HOME.resolve() == tmp_root.resolve():
                raise AssertionError("Blank DSH_HOME must not resolve to cwd")

        os.environ["DSH_HOME"] = "~/custom-dsh"
        tilde_install_module = load_root_script_module("install_student_os.py", "student_os_install_dsh_tilde_home")
        tilde_update_module = load_student_os_script_module(
            "update_student_os_impl.py",
            "student_os_update_dsh_tilde_home",
        )
        expected_tilde_home = isolated_home / "custom-dsh"
        if tilde_install_module.DSH_HOME.resolve() != expected_tilde_home.resolve():
            raise AssertionError(f"Installer should expand DSH_HOME=~/custom-dsh, got: {tilde_install_module.DSH_HOME}")
        if tilde_update_module.resolve_dsh_home().resolve() != expected_tilde_home.resolve():
            raise AssertionError("Updater should expand DSH_HOME=~/custom-dsh")

        os.environ.pop("DSH_HOME", None)
        os.environ["HOME"] = str(tmp_root / "home-for-current-platform")
        os.environ["USERPROFILE"] = str(tmp_root / "alternate-userprofile")
        platform_install_module = load_root_script_module(
            "install_student_os.py",
            "student_os_install_dsh_platform_home",
        )
        platform_update_module = load_student_os_script_module(
            "update_student_os_impl.py",
            "student_os_update_dsh_platform_home",
        )
        if platform_update_module.resolve_dsh_home().resolve() != platform_install_module.DSH_HOME.resolve():
            raise AssertionError("Installer and updater should resolve default DSH home with the same home semantics")
    finally:
        os.chdir(previous_cwd)
        for name, value in env_backup.items():
            if value is None:
                os.environ.pop(name, None)
            else:
                os.environ[name] = value


def verify_dsh_update_discovery(tmp_root: Path) -> None:
    env_backup = {name: os.environ.get(name) for name in ["HOME", "USERPROFILE", "CODEX_HOME", "DSH_HOME"]}
    previous_cwd = Path.cwd()
    try:
        isolated_home = tmp_root / "isolated-home"
        dsh_home = tmp_root / "dsh-home"
        os.environ["HOME"] = str(isolated_home)
        os.environ["USERPROFILE"] = str(isolated_home)
        os.environ["CODEX_HOME"] = str(isolated_home / ".codex")
        os.environ["DSH_HOME"] = str(dsh_home)
        run_root_script(
            "install_student_os.py",
            "--agent",
            "dsh",
            "--scope",
            "user",
            "--mode",
            "copy",
            "--json",
            env={
                "HOME": str(isolated_home),
                "USERPROFILE": str(isolated_home),
                "CODEX_HOME": str(isolated_home / ".codex"),
                "DSH_HOME": str(dsh_home),
            },
        )
        update_user_module = load_student_os_script_module(
            "update_student_os_impl.py",
            "student_os_update_dsh_user_discovery",
        )
        user_discovery_cwd = tmp_root / "user-discovery-cwd"
        user_discovery_cwd.mkdir(parents=True, exist_ok=True)
        os.chdir(user_discovery_cwd)
        user_target = dsh_home / "skills" / "student-os"
        discovered_user_candidates = [update_user_module.absolute_path(path) for path in update_user_module.user_targets()]
        if user_target.resolve() not in discovered_user_candidates:
            raise AssertionError(f"Updater user targets should include the DSH user install, got: {discovered_user_candidates}")
        update_user_module.user_targets = lambda: [user_target]
        update_user_module.project_candidates = list
        discovered_user = update_user_module.discover_target(None)
        if discovered_user != user_target.resolve():
            raise AssertionError(f"Updater should discover DSH user install, got: {discovered_user}")

        shutil.rmtree(dsh_home)
        project_root = tmp_root / "project"
        project_root.mkdir(parents=True, exist_ok=True)
        (project_root / ".git").mkdir()
        os.environ["DSH_HOME"] = str(tmp_root / "empty-dsh-home")
        run_root_script(
            "install_student_os.py",
            "--agent",
            "dsh",
            "--scope",
            "project",
            "--project-root",
            str(project_root),
            "--mode",
            "copy",
            "--json",
            env={
                "HOME": str(isolated_home),
                "USERPROFILE": str(isolated_home),
                "CODEX_HOME": str(isolated_home / ".codex"),
                "DSH_HOME": str(tmp_root / "empty-dsh-home"),
            },
        )
        nested_workdir = project_root / "nested" / "deeper"
        (project_root / "nested" / ".dsh").mkdir(parents=True, exist_ok=True)
        nested_workdir.mkdir(parents=True, exist_ok=True)
        os.chdir(nested_workdir)
        update_project_module = load_student_os_script_module(
            "update_student_os_impl.py",
            "student_os_update_dsh_project_discovery",
        )
        update_project_module.user_targets = list
        project_target = project_root / ".dsh" / "skills" / "student-os"
        discovered_project = update_project_module.discover_target(None)
        if discovered_project != project_target.resolve():
            raise AssertionError(f"Updater should discover DSH project install, got: {discovered_project}")

        dsh_parent = tmp_root / "dsh-only-parent"
        dsh_child = dsh_parent / "child"
        dsh_child.mkdir(parents=True, exist_ok=True)
        for root in (dsh_parent, dsh_child):
            run_root_script(
                "install_student_os.py",
                "--agent",
                "dsh",
                "--scope",
                "project",
                "--project-root",
                str(root),
                "--mode",
                "copy",
                "--json",
                env={
                    "HOME": str(isolated_home),
                    "USERPROFILE": str(isolated_home),
                    "CODEX_HOME": str(isolated_home / ".codex"),
                    "DSH_HOME": str(tmp_root / "empty-dsh-home"),
                },
            )
        dsh_child_workdir = dsh_child / "deeper"
        dsh_child_workdir.mkdir()
        os.chdir(dsh_child_workdir)
        update_dsh_boundary_module = load_student_os_script_module(
            "update_student_os_impl.py",
            "student_os_update_dsh_project_boundary",
        )
        update_dsh_boundary_module.user_targets = list
        dsh_child_target = dsh_child / ".dsh" / "skills" / "student-os"
        discovered_dsh_child = update_dsh_boundary_module.discover_target(None)
        if discovered_dsh_child != dsh_child_target.resolve():
            raise AssertionError(f"Updater should stop at the nearest DSH project install, got: {discovered_dsh_child}")

        dsh_empty_child = dsh_parent / "empty-child"
        dsh_empty_child_workdir = dsh_empty_child / "deeper"
        (dsh_empty_child / ".dsh").mkdir(parents=True)
        dsh_empty_child_workdir.mkdir()
        os.chdir(dsh_empty_child_workdir)
        update_empty_dsh_boundary_module = load_student_os_script_module(
            "update_student_os_impl.py",
            "student_os_update_empty_dsh_project_boundary",
        )
        update_empty_dsh_boundary_module.user_targets = list
        if update_empty_dsh_boundary_module.discover_targets():
            raise AssertionError("Updater should not cross an empty DSH-only project boundary to a parent install")
    finally:
        os.chdir(previous_cwd)
        for name, value in env_backup.items():
            if value is None:
                os.environ.pop(name, None)
            else:
                os.environ[name] = value


def verify_legacy_link_install_detection(tmp_root: Path) -> None:
    install_module = load_root_script_module("install_student_os.py", "student_os_install_link_smoke")
    legacy_destination = tmp_root / "legacy-linked-student-os"
    try:
        legacy_destination.symlink_to((ROOT / "student-os").resolve(), target_is_directory=True)
    except OSError:
        return
    if not install_module.same_link_install(legacy_destination):
        raise AssertionError("install_student_os.py should recognize legacy whole-directory symlink installs")
    child_link_install = tmp_root / "child-linked-student-os"
    child_link_install.mkdir(parents=True, exist_ok=True)
    (child_link_install / "SKILL.md").symlink_to((ROOT / "student-os" / "SKILL.md").resolve())
    if not install_module.same_link_install(child_link_install):
        raise AssertionError("install_student_os.py should still recognize linked installs that symlink top-level children")
    created_entries = install_module.sync_linked_entries(child_link_install, ROOT / "student-os")
    if "references" not in created_entries:
        raise AssertionError("install_student_os.py should recreate missing top-level symlink entries for linked installs")
    if not (child_link_install / "references").is_symlink():
        raise AssertionError("install_student_os.py should materialize recreated linked entries as symlinks")
    source_manifest = ROOT / "student-os" / ".student-os-install.json"
    if source_manifest.exists():
        source_manifest.unlink()
    legacy_root = tmp_root / "legacy-link-root"
    legacy_root.mkdir(parents=True, exist_ok=True)
    legacy_install = legacy_root / "student-os"
    try:
        legacy_install.symlink_to((ROOT / "student-os").resolve(), target_is_directory=True)
    except OSError:
        return
    result = install_module.install_one(
        install_module.InstallTarget(agent="codex", scope="user", root=legacy_root),
        "auto",
        False,
        source_repo="https://example.com/fork.git",
        source_ref="main",
    )
    if result["manifest"] != "":
        raise AssertionError("Legacy whole-directory symlink installs should not write manifests into the source checkout")
    if source_manifest.exists():
        raise AssertionError("Legacy whole-directory symlink installs should not create a manifest in the source repo")


def verify_dsh_native_plugin() -> bool:
    plugin_root = ROOT / "integrations" / "dsh"
    if not (plugin_root / "package.json").exists():
        raise AssertionError("DSH native plugin package.json is missing")
    if not (plugin_root / "package-lock.json").exists():
        raise AssertionError("DSH native plugin package-lock.json is required for npm ci")
    build_module = load_root_script_module("dsh_plugin_build.py", "student_os_dsh_plugin_build_smoke")
    node_exe = shutil.which("node.exe") or shutil.which("node")
    npm_exe = shutil.which("npm.cmd") or shutil.which("npm")
    if node_exe is None or npm_exe is None:
        return False

    def git_status_snapshot() -> set[str]:
        result = subprocess.run(
            ["git", "status", "--porcelain=v1", "--untracked-files=all"],
            cwd=ROOT,
            check=True,
            capture_output=True,
            text=True,
            encoding="utf-8",
        )
        return set(result.stdout.splitlines())

    env = {
        **os.environ,
        "PYTHONIOENCODING": "utf-8",
        "PYTHONDONTWRITEBYTECODE": "1",
        "STUDENT_OS_REPO_ROOT": str(ROOT),
    }
    before_status = git_status_snapshot()
    try:
        build_module.ensure_dsh_plugin_build(env=env)
        subprocess.run(
            [npm_exe, "run", "test"],
            cwd=plugin_root,
            check=True,
            capture_output=True,
            text=True,
            encoding="utf-8",
            env=env,
        )
    finally:
        after_status = git_status_snapshot()
        if after_status != before_status:
            added = sorted(after_status - before_status)
            removed = sorted(before_status - after_status)
            raise AssertionError(
                "DSH native plugin smoke changed Git status; "
                f"new/changed entries={added}, cleared entries={removed}"
            )
    return True


def write_fake_dsh_plugin_package(plugin_root: Path) -> None:
    (plugin_root / "src").mkdir(parents=True)
    (plugin_root / "package.json").write_text(
        json.dumps({"scripts": {"build": "fake-build"}}, ensure_ascii=False),
        encoding="utf-8",
    )
    (plugin_root / "package-lock.json").write_text(
        json.dumps({"lockfileVersion": 3, "packages": {}}, ensure_ascii=False),
        encoding="utf-8",
    )
    (plugin_root / "tsconfig.json").write_text("{}", encoding="utf-8")
    (plugin_root / "src" / "index.ts").write_text("export const value = 1\n", encoding="utf-8")


def write_fake_node_and_npm(bin_dir: Path) -> None:
    bin_dir.mkdir(parents=True)
    fake_npm = bin_dir / "fake_npm.py"
    fake_npm.write_text(
        "\n".join(
            [
                "import os, sys, time",
                "from pathlib import Path",
                "root = Path.cwd()",
                "args = sys.argv[1:]",
                "if os.environ.get('STUDENT_OS_DSH_SLEEP'):",
                "    time.sleep(float(os.environ['STUDENT_OS_DSH_SLEEP']))",
                "if args[:1] == ['ci']:",
                "    if os.environ.get('STUDENT_OS_DSH_FAIL_CI'): sys.exit(7)",
                "    (root / 'node_modules').mkdir(exist_ok=True)",
                "    sys.exit(0)",
                "if args[:1] == ['ls']:",
                "    sys.exit(0 if (root / 'node_modules').is_dir() else 1)",
                "if args[:2] == ['run', 'build']:",
                "    if os.environ.get('STUDENT_OS_DSH_FAIL_BUILD'): sys.exit(9)",
                "    (root / 'dist').mkdir(exist_ok=True)",
                "    (root / 'dist' / 'index.js').write_text('export default {}\\n', encoding='utf-8')",
                "    (root / 'dist' / 'index.d.ts').write_text('export {}\\n', encoding='utf-8')",
                "    sys.exit(0)",
                "sys.exit(0)",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    if os.name == "nt":
        (bin_dir / "node.cmd").write_text("@echo off\r\nexit /b 0\r\n", encoding="utf-8")
        (bin_dir / "npm.cmd").write_text(f'@echo off\r\n"{sys.executable}" "{fake_npm}" %*\r\nexit /b %ERRORLEVEL%\r\n', encoding="utf-8")
    else:
        node = bin_dir / "node"
        npm = bin_dir / "npm"
        node.write_text("#!/bin/sh\nexit 0\n", encoding="utf-8")
        npm.write_text(f"#!/bin/sh\nexec {sys.executable!r} {str(fake_npm)!r} \"$@\"\n", encoding="utf-8")
        node.chmod(0o755)
        npm.chmod(0o755)


def read_dsh_build_log(log_path: Path) -> list[list[str]]:
    if not log_path.exists():
        return []
    return [json.loads(line)["argv"][1:] for line in log_path.read_text(encoding="utf-8").splitlines() if line.strip()]


def clear_dsh_build_log(log_path: Path) -> None:
    log_path.write_text("", encoding="utf-8")


def verify_dsh_plugin_build_lifecycle(tmp_root: Path) -> None:
    build_module = load_root_script_module("dsh_plugin_build.py", "student_os_dsh_build_lifecycle_smoke")
    fake_bin = tmp_root / "fake-bin"
    write_fake_node_and_npm(fake_bin)
    plugin_root = tmp_root / "plugin"
    write_fake_dsh_plugin_package(plugin_root)
    log_path = tmp_root / "build-commands.jsonl"
    env = {"PATH": str(fake_bin) + os.pathsep + os.environ.get("PATH", ""), "STUDENT_OS_DSH_BUILD_LOG": str(log_path)}

    first = build_module.ensure_dsh_plugin_build(plugin_root=plugin_root, env=env, lock_timeout_seconds=10)
    if first["npm_ci"] is not True or first["build"] is not True or first["reused"] is not False:
        raise AssertionError(f"First DSH plugin build should run npm ci and build, got: {first}")
    first_commands = read_dsh_build_log(log_path)
    if ["ci", "--ignore-scripts", "--no-audit", "--no-fund"] not in first_commands or ["run", "build"] not in first_commands:
        raise AssertionError(f"First DSH plugin build should log npm ci and build, got: {first_commands}")

    clear_dsh_build_log(log_path)
    second = build_module.ensure_dsh_plugin_build(plugin_root=plugin_root, env=env, lock_timeout_seconds=10)
    if second["reused"] is not True or second["npm_ci"] is not False or second["build"] is not False:
        raise AssertionError(f"Second DSH plugin build should reuse existing build, got: {second}")
    if read_dsh_build_log(log_path):
        raise AssertionError("Reused DSH plugin build should not execute npm commands")

    (plugin_root / "src" / "index.ts").write_text("export const value = 2\n", encoding="utf-8")
    source_changed = build_module.ensure_dsh_plugin_build(plugin_root=plugin_root, env=env, lock_timeout_seconds=10)
    if source_changed["npm_ci"] is not False or source_changed["build"] is not True:
        raise AssertionError(f"Source-only DSH plugin change should rebuild without npm ci, got: {source_changed}")
    source_commands = read_dsh_build_log(log_path)
    if ["ci", "--ignore-scripts", "--no-audit", "--no-fund"] in source_commands or ["run", "build"] not in source_commands:
        raise AssertionError(f"Source-only DSH plugin change command set is wrong: {source_commands}")

    shutil.rmtree(plugin_root / "node_modules")
    clear_dsh_build_log(log_path)
    missing_deps = build_module.ensure_dsh_plugin_build(plugin_root=plugin_root, env=env, lock_timeout_seconds=10)
    if missing_deps["npm_ci"] is not True or missing_deps["build"] is not True:
        raise AssertionError(f"Missing DSH plugin dependencies should run npm ci and build, got: {missing_deps}")

    (plugin_root / "src" / "index.ts").write_text("export const value = 3\n", encoding="utf-8")
    clear_dsh_build_log(log_path)
    try:
        build_module.ensure_dsh_plugin_build(
            plugin_root=plugin_root,
            env={**env, "STUDENT_OS_DSH_FAIL_BUILD": "1"},
            lock_timeout_seconds=10,
        )
    except RuntimeError:
        pass
    else:
        raise AssertionError("Failing DSH plugin build should raise")
    if (plugin_root / "dist" / ".student-os-build.lock").exists():
        raise AssertionError("DSH plugin build lock should be released after failures")
    recovered = build_module.ensure_dsh_plugin_build(plugin_root=plugin_root, env=env, lock_timeout_seconds=10)
    if recovered["build"] is not True:
        raise AssertionError(f"DSH plugin build should recover after a failed build, got: {recovered}")

    concurrent_root = tmp_root / "concurrent-plugin"
    write_fake_dsh_plugin_package(concurrent_root)
    concurrent_log = tmp_root / "concurrent-build-commands.jsonl"
    concurrent_env = {
        **os.environ,
        "PATH": str(fake_bin) + os.pathsep + os.environ.get("PATH", ""),
        "STUDENT_OS_DSH_BUILD_LOG": str(concurrent_log),
        "STUDENT_OS_DSH_SLEEP": "0.25",
        "PYTHONIOENCODING": "utf-8",
        "PYTHONDONTWRITEBYTECODE": "1",
    }
    script = (
        "import importlib.util, pathlib, sys;"
        f"p=pathlib.Path({str(ROOT_SCRIPTS / 'dsh_plugin_build.py')!r});"
        "s=importlib.util.spec_from_file_location('dsh_plugin_build_subprocess', p);"
        "m=importlib.util.module_from_spec(s); s.loader.exec_module(m);"
        f"m.ensure_dsh_plugin_build(plugin_root=pathlib.Path({str(concurrent_root)!r}), lock_timeout_seconds=10)"
    )
    processes = [
        subprocess.Popen(
            [sys.executable, "-B", "-c", script],
            cwd=ROOT,
            env=concurrent_env,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
        )
        for _ in range(2)
    ]
    for process in processes:
        stdout, stderr = process.communicate(timeout=30)
        if process.returncode != 0:
            raise AssertionError(f"Concurrent DSH plugin build failed, rc={process.returncode}, stdout={stdout}, stderr={stderr}")
    concurrent_commands = read_dsh_build_log(concurrent_log)
    ci_count = sum(1 for argv in concurrent_commands if argv[:1] == ["ci"])
    build_count = sum(1 for argv in concurrent_commands if argv[:2] == ["run", "build"])
    if ci_count != 1 or build_count != 1:
        raise AssertionError(f"Concurrent DSH plugin build should run one npm ci and one build, got: {concurrent_commands}")


def verify_dsh_bootstrap(tmp_root: Path, native_plugin_available: bool) -> bool:
    git_exe = shutil.which("git.exe") or shutil.which("git")
    if git_exe is None:
        raise AssertionError("git is required for DSH bootstrap smoke tests")

    failure_vault = tmp_root / "missing-node-vault"
    git_only_path = str(Path(git_exe).parent)
    failure_output = run_root_script_failure(
        "bootstrap_dsh.py",
        "--project-root",
        str(failure_vault),
        "--json",
        env={"PATH": git_only_path},
    )
    failure_payload = json.loads(failure_output)
    if failure_payload.get("ok") is not False:
        raise AssertionError("DSH bootstrap failure should return ok:false")
    if failure_payload.get("stage") != "plugin-build":
        raise AssertionError(f"DSH bootstrap missing-node failure should report plugin-build, got: {failure_payload}")
    if "activation" in failure_payload:
        raise AssertionError("DSH bootstrap failure should not claim activation instructions")

    root_failure_output = run_root_script_failure(
        "bootstrap_dsh.py",
        "--project-root",
        str(ROOT),
        "--json",
    )
    root_failure_payload = json.loads(root_failure_output)
    if root_failure_payload.get("ok") is not False or root_failure_payload.get("stage") != "project-root":
        raise AssertionError(f"DSH bootstrap should refuse the source checkout as project root, got: {root_failure_payload}")

    symlink_vault = tmp_root / "external-dsh-symlink-vault"
    external_dsh = tmp_root / "external-dsh-target"
    symlink_vault.mkdir(parents=True)
    external_dsh.mkdir()
    try:
        (symlink_vault / ".dsh").symlink_to(external_dsh, target_is_directory=True)
    except OSError:
        pass
    else:
        symlink_failure_output = run_root_script_failure(
            "bootstrap_dsh.py",
            "--project-root",
            str(symlink_vault),
            "--json",
        )
        symlink_failure_payload = json.loads(symlink_failure_output)
        if symlink_failure_payload.get("ok") is not False or symlink_failure_payload.get("stage") != "project-root":
            raise AssertionError(f"DSH bootstrap should refuse external .dsh symlinks, got: {symlink_failure_payload}")

    overlay_symlink_vault = tmp_root / "overlay-symlink-vault"
    overlay_symlink_vault_dsh = overlay_symlink_vault / ".dsh"
    overlay_symlink_vault_dsh.mkdir(parents=True)
    external_overlay = tmp_root / "external-overlay.yml"
    external_overlay.write_text("# external overlay\n", encoding="utf-8")
    try:
        (overlay_symlink_vault_dsh / "student-os.cordis.yml").symlink_to(external_overlay)
    except OSError:
        pass
    else:
        overlay_symlink_failure_output = run_root_script_failure(
            "bootstrap_dsh.py",
            "--project-root",
            str(overlay_symlink_vault),
            "--force-overlay",
            "--json",
        )
        overlay_symlink_failure_payload = json.loads(overlay_symlink_failure_output)
        if overlay_symlink_failure_payload.get("ok") is not False or overlay_symlink_failure_payload.get("stage") != "project-root":
            raise AssertionError(
                f"DSH bootstrap should refuse symlinked overlay files even with --force-overlay, got: {overlay_symlink_failure_payload}"
            )

    if not native_plugin_available:
        return False

    vault = tmp_root / "vault"
    bootstrap_build_log = tmp_root / "bootstrap-build-commands.jsonl"
    bootstrap_env = {"STUDENT_OS_DSH_BUILD_LOG": str(bootstrap_build_log)}
    init_git_repo(vault)
    first_output = run_root_script(
        "bootstrap_dsh.py",
        "--project-root",
        str(vault),
        "--json",
        env=bootstrap_env,
    )
    first_payload = json.loads(first_output)
    if first_payload.get("ok") is not True:
        raise AssertionError(f"DSH bootstrap should succeed, got: {first_payload}")

    skill_path = vault / ".dsh" / "skills" / "student-os"
    overlay_path = vault / ".dsh" / "student-os.cordis.yml"
    plugin_entry = ROOT / "integrations" / "dsh" / "dist" / "index.js"
    manifest_path = skill_path / ".student-os-install.json"

    ensure_exists(skill_path / "SKILL.md")
    ensure_exists(plugin_entry)
    ensure_exists(overlay_path)
    ensure_exists(manifest_path)

    if Path(first_payload["skill"]["path"]).resolve() != skill_path.resolve():
        raise AssertionError("DSH bootstrap JSON should report the project skill path")
    if Path(first_payload["plugin"]["entry"]).resolve() != plugin_entry.resolve():
        raise AssertionError("DSH bootstrap JSON should report the built plugin entry")
    if Path(first_payload["overlay"]["path"]).resolve() != overlay_path.resolve():
        raise AssertionError("DSH bootstrap JSON should report the project-local overlay path")
    activation = first_payload.get("activation", {})
    if activation.get("active_in_current_process") is not False or activation.get("restart_required") is not True:
        raise AssertionError(f"DSH bootstrap activation flags are wrong: {activation}")
    expected_argv = ["dsh", "web", "--patch", str(overlay_path.resolve())]
    if activation.get("argv") != expected_argv:
        raise AssertionError(f"DSH bootstrap restart argv mismatch: {activation.get('argv')}")
    git_payload = first_payload.get("git", {})
    if git_payload.get("before", {}).get("available") is not True or git_payload.get("after", {}).get("available") is not True:
        raise AssertionError(f"DSH bootstrap should report vault Git status when available: {git_payload}")
    if not any(".dsh/" in entry or ".dsh\\" in entry for entry in git_payload.get("added", [])):
        raise AssertionError(f"DSH bootstrap should report newly created project DSH files: {git_payload}")
    dsh_files = git_payload.get("dsh_files", {})
    if ".dsh/student-os.cordis.yml" not in dsh_files.get("added", []):
        raise AssertionError(f"DSH bootstrap should report created overlay content: {git_payload}")

    overlay = overlay_path.read_text(encoding="utf-8")
    expected_entry = plugin_entry.resolve().as_uri()
    if expected_entry not in overlay:
        raise AssertionError("DSH bootstrap overlay should reference the absolute plugin entry as a file:// URL")
    if str(plugin_entry) in overlay and "\\" in str(plugin_entry):
        raise AssertionError("DSH bootstrap overlay should not use backslash-escaped Windows paths")
    if len(list((vault / ".dsh").glob("student-os*.cordis.yml"))) != 1:
        raise AssertionError("DSH bootstrap should create exactly one project-local Student OS overlay")

    first_manifest = manifest_path.read_text(encoding="utf-8")
    clear_dsh_build_log(bootstrap_build_log)
    second_output = run_root_script(
        "bootstrap_dsh.py",
        "--project-root",
        str(vault),
        "--json",
        env=bootstrap_env,
    )
    second_payload = json.loads(second_output)
    if second_payload.get("ok") is not True:
        raise AssertionError(f"Second DSH bootstrap should succeed, got: {second_payload}")
    if second_payload["skill"]["path"] != first_payload["skill"]["path"]:
        raise AssertionError("DSH bootstrap should keep the same skill path on repeated runs")
    if second_payload["overlay"]["path"] != first_payload["overlay"]["path"]:
        raise AssertionError("DSH bootstrap should keep the same overlay path on repeated runs")
    if manifest_path.read_text(encoding="utf-8") != first_manifest:
        raise AssertionError("DSH bootstrap idempotency should not rewrite the install manifest")
    if len(list((vault / ".dsh").glob("student-os*.cordis.yml"))) != 1:
        raise AssertionError("Repeated DSH bootstrap should not create extra overlays")
    if second_payload.get("plugin", {}).get("reused") is not True:
        raise AssertionError(f"Second DSH bootstrap should reuse the existing plugin build, got: {second_payload.get('plugin')}")
    if read_dsh_build_log(bootstrap_build_log):
        raise AssertionError("Second DSH bootstrap should not execute npm ci or npm run build")
    ensure_exists(skill_path / "SKILL.md")

    overlay_path.write_text("# user overlay\n", encoding="utf-8")
    conflict_output = run_root_script_failure(
        "bootstrap_dsh.py",
        "--project-root",
        str(vault),
        "--json",
    )
    conflict_payload = json.loads(conflict_output)
    if conflict_payload.get("ok") is not False or conflict_payload.get("stage") != "overlay-write":
        raise AssertionError(f"DSH bootstrap should reject a different existing overlay, got: {conflict_payload}")
    dangling_backup = overlay_path.with_name(overlay_path.name + ".bak")
    dangling_target = tmp_root / "missing-overlay-backup-target.yml"
    try:
        dangling_backup.symlink_to(dangling_target)
    except OSError:
        pass
    else:
        dangling_backup_output = run_root_script_failure(
            "bootstrap_dsh.py",
            "--project-root",
            str(vault),
            "--force-overlay",
            "--json",
        )
        dangling_backup_payload = json.loads(dangling_backup_output)
        if dangling_backup_payload.get("ok") is not False or dangling_backup_payload.get("stage") != "overlay-write":
            raise AssertionError(f"DSH bootstrap should reject symlink backup paths, got: {dangling_backup_payload}")
        dangling_backup.unlink()
    force_output = run_root_script(
        "bootstrap_dsh.py",
        "--project-root",
        str(vault),
        "--force-overlay",
        "--json",
    )
    force_payload = json.loads(force_output)
    if force_payload.get("ok") is not True:
        raise AssertionError(f"DSH bootstrap --force-overlay should succeed, got: {force_payload}")
    backup = force_payload.get("overlay", {}).get("backup")
    if not backup:
        raise AssertionError("DSH bootstrap --force-overlay should report an overlay backup")
    ensure_exists(Path(backup))
    force_git = force_payload.get("git", {})
    if ".dsh/student-os.cordis.yml" not in force_git.get("dsh_files", {}).get("modified", []):
        raise AssertionError(f"DSH bootstrap --force-overlay should report modified overlay content: {force_git}")

    overlay_path.write_text("# concurrent overlay\n", encoding="utf-8")
    concurrent_env = {**os.environ, "PYTHONIOENCODING": "utf-8", "PYTHONDONTWRITEBYTECODE": "1"}
    concurrent_args = [
        sys.executable,
        "-B",
        str(ROOT_SCRIPTS / "bootstrap_dsh.py"),
        "--project-root",
        str(vault),
        "--force-overlay",
        "--json",
    ]
    processes = [
        subprocess.Popen(
            concurrent_args,
            cwd=ROOT,
            env=concurrent_env,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
        )
        for _ in range(2)
    ]
    for process in processes:
        stdout, stderr = process.communicate(timeout=120)
        if process.returncode != 0:
            raise AssertionError(f"Concurrent DSH bootstrap should succeed, rc={process.returncode}, stdout={stdout}, stderr={stderr}")
        payload = json.loads(stdout)
        if payload.get("ok") is not True:
            raise AssertionError(f"Concurrent DSH bootstrap returned failure payload: {payload}")
    if "# concurrent overlay" in overlay_path.read_text(encoding="utf-8"):
        raise AssertionError("Concurrent DSH bootstrap should restore the generated overlay content")
    return True


def verify_update_source_override_and_project_copy_detection(tmp_root: Path) -> None:
    install_module = load_root_script_module("install_student_os.py", "student_os_install_override_smoke")
    update_module = load_root_script_module("update_student_os.py", "student_os_update_override_smoke")

    override_repo = tmp_root / "override-source"
    init_git_repo(override_repo)
    override_commit = seed_fake_student_os_source(override_repo, version_label="override")

    target = tmp_root / "override-install"
    shutil.copytree(override_repo / "student-os", target)
    manifest = install_module.build_install_manifest(
        destination=target,
        agent="codex",
        scope="user",
        install_method="copied",
        used_symlink=False,
        source_repo="https://github.com/Gu-Heping/college-student-workflow.git",
        source_ref="main",
        installed_commit=override_commit,
        linked_source_path="",
    )
    install_module.write_manifest(target, manifest)
    info = update_module.build_install_info(target, str(override_repo.resolve()), "main")
    if info.source_repo != str(override_repo.resolve()):
        raise AssertionError("Explicit --repo equivalents should override manifest source_repo values")
    if info.source_ref != "main":
        raise AssertionError("Explicit --ref equivalents should override manifest source_ref values")

    vault_repo = tmp_root / "vault-repo"
    init_git_repo(vault_repo)
    (vault_repo / "README.md").write_text("# Vault\n", encoding="utf-8", newline="\n")
    commit_all(vault_repo, "init vault")
    project_install = vault_repo / ".codex" / "skills" / "student-os"
    project_install.parent.mkdir(parents=True, exist_ok=True)
    shutil.copytree(override_repo / "student-os", project_install)
    project_manifest = install_module.build_install_manifest(
        destination=project_install,
        agent="codex",
        scope="project",
        install_method="copied",
        used_symlink=False,
        source_repo=str(override_repo.resolve()),
        source_ref="main",
        installed_commit=override_commit,
        linked_source_path="",
    )
    install_module.write_manifest(project_install, project_manifest)
    project_info = update_module.build_install_info(project_install, None, None)
    if project_info.install_kind != "copy":
        raise AssertionError("Copied project installs inside a git repository should still be treated as copy installs")

    env_backup = {name: os.environ.get(name) for name in ["HOME", "USERPROFILE", "CODEX_HOME"]}
    previous_cwd = Path.cwd()
    try:
        isolated_home = tmp_root / "isolated-home"
        os.environ["HOME"] = str(isolated_home)
        os.environ["USERPROFILE"] = str(isolated_home)
        os.environ["CODEX_HOME"] = str(isolated_home / ".codex")
        os.chdir(vault_repo)
        sys.modules.pop("update_student_os_impl", None)
        project_discovery_module = load_root_script_module("update_student_os.py", "student_os_update_discovery_smoke")
        discovered_target = project_discovery_module.discover_target(None)
    finally:
        os.chdir(previous_cwd)
        for name, value in env_backup.items():
            if value is None:
                os.environ.pop(name, None)
            else:
                os.environ[name] = value
    if discovered_target != project_install:
        raise AssertionError("update_student_os.py should discover project-scoped installs by default")

    legacy_project_install = vault_repo / ".claude" / "skills" / "student-os"
    legacy_project_install.parent.mkdir(parents=True, exist_ok=True)
    shutil.copytree(override_repo / "student-os", legacy_project_install)
    legacy_project_info = update_module.build_install_info(legacy_project_install, None, None)
    if legacy_project_info.install_kind != "copy":
        raise AssertionError("Manifestless copied installs inside a parent repository should default to copy mode")

    missing_manifest_failure = run_root_script_failure(
        "update_student_os.py",
        "--apply",
        "--target",
        str(legacy_project_install),
        "--repo",
        str(override_repo.resolve()),
        "--ref",
        "main",
    )
    if "--force" not in missing_manifest_failure:
        raise AssertionError("Manifestless copied installs should require --force before replacement")


def verify_self_update_workflow(tmp_root: Path) -> None:
    source_repo = tmp_root / "student-os-source"
    init_git_repo(source_repo)
    commit_v1 = seed_fake_student_os_source(source_repo, version_label="v1")
    install_target = build_copy_install_fixture(tmp_root, source_repo, commit_v1)

    seed_fake_student_os_source(source_repo, version_label="v2")
    (source_repo / "student-os" / "templates" / "fixture.md").write_text(
        "# Fixture Template v2 updated\n",
        encoding="utf-8",
        newline="\n",
    )
    commit_v2 = commit_all(source_repo, "fixture: v2 template tweak")

    check_payload = json.loads(
        run_root_script(
            "update_student_os.py",
            "--check",
            "--target",
            str(install_target),
            "--repo",
            str(source_repo.resolve()),
            "--ref",
            "main",
            "--json",
        )
    )
    if not check_payload["update_available"]:
        raise AssertionError("update_student_os.py --check should report updates when the remote commit is newer")
    if check_payload["current_commit"] != commit_v1:
        raise AssertionError("update_student_os.py --check should report the installed commit from the manifest")
    if check_payload["latest_commit"] != commit_v2:
        raise AssertionError("update_student_os.py --check should resolve the latest remote commit")

    apply_payload = json.loads(
        run_root_script(
            "update_student_os.py",
            "--apply",
            "--target",
            str(install_target),
            "--repo",
            str(source_repo.resolve()),
            "--ref",
            "main",
            "--json",
        )
    )
    if not apply_payload["updated"]:
        raise AssertionError("copy-mode self-update should report an applied update")
    if not apply_payload["backup_path"]:
        raise AssertionError("copy-mode self-update should create and report a backup path")
    if "--restore-backup" not in apply_payload["rollback_command"]:
        raise AssertionError("copy-mode self-update should print a restore command for rollback")
    if ".student-os-local-overrides" not in apply_payload["preserved_override_paths"]:
        raise AssertionError("copy-mode self-update should preserve the documented override directory")
    if ".student-os-install.local.json" not in apply_payload["preserved_override_paths"]:
        raise AssertionError("copy-mode self-update should preserve the documented local override file")
    ensure_contains(install_target / "SKILL.md", "Fixture v2")
    ensure_contains(install_target / ".student-os-local-overrides" / "notes.md", "Local override")
    ensure_contains(install_target / ".student-os-install.local.json", '"theme":"custom"')
    manifest = json.loads((install_target / ".student-os-install.json").read_text(encoding="utf-8"))
    if manifest["installed_commit"] != commit_v2:
        raise AssertionError("copy-mode self-update should refresh the install manifest commit")
    if manifest["source_repo"] != str(source_repo.resolve()):
        raise AssertionError("copy-mode self-update should preserve the configured source repo in the manifest")
    if Path(manifest["target_path"]) != install_target.resolve():
        raise AssertionError("copy-mode self-update should record the final install target in the refreshed manifest")
    installed_updater = install_target / "scripts" / "update_student_os.py"
    ensure_exists(installed_updater)
    installed_check_payload = json.loads(
        run_path_script(
            installed_updater,
            "--check",
            "--target",
            str(install_target),
            "--repo",
            str(source_repo.resolve()),
            "--ref",
            "main",
            "--json",
            cwd=install_target,
        )
    )
    if Path(installed_check_payload["target_path"]) != install_target:
        raise AssertionError("Installed updater entrypoint should work from the installed skill directory")
    if "local_changes" in installed_check_payload and installed_check_payload["local_changes"]:
        raise AssertionError("Installed updater should not report its own __pycache__ artifacts as local changes")
    assert_no_pycache(install_target)
    assert_no_pycache(source_repo)

    (install_target / "SKILL.md").write_text(
        (install_target / "SKILL.md").read_text(encoding="utf-8") + "\nLocal drift\n",
        encoding="utf-8",
        newline="\n",
    )
    (source_repo / "student-os" / "SKILL.md").write_text(
        (source_repo / "student-os" / "SKILL.md").read_text(encoding="utf-8") + "\nFixture v3\n",
        encoding="utf-8",
        newline="\n",
    )
    commit_all(source_repo, "fixture: v3")
    failure = run_root_script_failure(
        "update_student_os.py",
        "--apply",
        "--target",
        str(install_target),
        "--repo",
        str(source_repo.resolve()),
        "--ref",
        "main",
    )
    if "--force" not in failure:
        raise AssertionError("self-update should refuse overwriting local installed-skill drift without --force")


def verify_extract_release_notes(tmp_root: Path) -> None:
    output_plain = tmp_root / "notes-0.6.0.md"
    output_prefixed = tmp_root / "notes-v0.6.0.md"
    run_root_script(
        "extract_release_notes.py",
        "--version",
        "0.6.0",
        "--output",
        str(output_plain),
    )
    run_root_script(
        "extract_release_notes.py",
        "--version",
        "v0.6.0",
        "--output",
        str(output_prefixed),
    )
    plain_text = output_plain.read_text(encoding="utf-8")
    prefixed_text = output_prefixed.read_text(encoding="utf-8")
    if "多学期知识库支持" not in plain_text:
        raise AssertionError("extract_release_notes.py should extract the 0.6.0 section body")
    if "## [0.6.0]" in plain_text or "## [0.5.0]" in plain_text:
        raise AssertionError("extract_release_notes.py should omit version headings from the body")
    if "## [Unreleased]" in plain_text:
        raise AssertionError("extract_release_notes.py must not include the Unreleased section")
    if plain_text != prefixed_text:
        raise AssertionError("0.6.0 and v0.6.0 should extract identical release notes")

    missing = run_root_script_failure(
        "extract_release_notes.py",
        "--version",
        "9.9.9",
        "--output",
        str(tmp_root / "missing.md"),
    )
    if "9.9.9" not in missing:
        raise AssertionError("extract_release_notes.py should fail clearly when the version is missing")

    unreleased = run_root_script_failure(
        "extract_release_notes.py",
        "--version",
        "Unreleased",
        "--output",
        str(tmp_root / "unreleased.md"),
    )
    if "Unreleased" not in unreleased:
        raise AssertionError("extract_release_notes.py should refuse Unreleased as release notes")


def main() -> int:
    parser = argparse.ArgumentParser(description="Run smoke tests for student-os scaffolding workflows.")
    parser.add_argument(
        "--refresh-examples",
        action="store_true",
        help="Replace example repositories in ./examples with fresh generated snapshots.",
    )
    args = parser.parse_args()
    today = date.today()

    with tempfile.TemporaryDirectory(prefix="student-os-smoke-") as tmp:
        tmp_root = Path(tmp)
        single_repo = tmp_root / "single-semester-demo"
        multi_repo = tmp_root / "multi-semester-demo"
        legacy_repo = tmp_root / "legacy-layout-demo"
        weekly_parent_repo = tmp_root / "weekly" / "nested-weekly-parent-demo"
        grouping_repo = tmp_root / "grouping-demo"

        build_single_semester(single_repo, today)
        build_multi_semester(multi_repo, today)
        build_legacy_layout(legacy_repo, today)
        build_repo_inside_weekly_parent(weekly_parent_repo, today)
        inspect_repo_fixture = tmp_root / "inspect-repo-demo"
        copy_repo(multi_repo, inspect_repo_fixture)
        verify_inspect_repo(inspect_repo_fixture)
        verify_git_grouping(grouping_repo, today)
        verify_chinese_slug_support(tmp_root / "unicode-course-demo", today)
        verify_install_manifest_generation(tmp_root / "install-manifest-demo")
        verify_dsh_skill_frontmatter_contract()
        verify_dsh_install_paths(tmp_root / "dsh-install-demo")
        verify_dsh_home_resolution(tmp_root / "dsh-home-resolution-demo")
        verify_dsh_update_discovery(tmp_root / "dsh-update-discovery-demo")
        verify_dsh_plugin_build_lifecycle(tmp_root / "dsh-build-lifecycle-demo")
        dsh_native_plugin_ran = verify_dsh_native_plugin()
        dsh_bootstrap_ran = verify_dsh_bootstrap(tmp_root / "dsh-bootstrap-demo", dsh_native_plugin_ran)
        verify_legacy_link_install_detection(tmp_root / "legacy-link-install-demo")
        verify_update_source_override_and_project_copy_detection(tmp_root / "update-override-demo")
        verify_self_update_workflow(tmp_root / "self-update-demo")
        verify_token_loader(tmp_root / "token-loader-demo")
        verify_scaffold_gitattributes(tmp_root / "scaffold-gitattributes-demo")
        verify_ensure_frontmatter(tmp_root / "ensure-frontmatter-demo")
        verify_extract_release_notes(tmp_root / "extract-release-notes-demo")

        if args.refresh_examples:
            EXAMPLES_ROOT.mkdir(parents=True, exist_ok=True)
            copy_repo(single_repo, EXAMPLES_ROOT / "single-semester-demo")
            copy_repo(multi_repo, EXAMPLES_ROOT / "multi-semester-demo")
            copy_repo(legacy_repo, EXAMPLES_ROOT / "legacy-layout-demo")
            for example_root, source_root in [
                (EXAMPLES_ROOT / "single-semester-demo", single_repo),
                (EXAMPLES_ROOT / "multi-semester-demo", multi_repo),
                (EXAMPLES_ROOT / "legacy-layout-demo", legacy_repo),
            ]:
                materialize_empty_dirs(example_root)
                scrub_example_paths(example_root, source_root)
                normalize_text_files(example_root)

    print("OK single-semester-demo")
    print("OK multi-semester-demo")
    print("OK legacy-layout-demo")
    print("OK unicode-course-demo")
    print("OK install-manifest-demo")
    print("OK dsh-skill-frontmatter")
    print("OK dsh-install-demo")
    print("OK dsh-home-resolution-demo")
    print("OK dsh-update-discovery-demo")
    print("OK dsh-build-lifecycle-demo")
    print("OK dsh-native-plugin" if dsh_native_plugin_ran else "SKIP dsh-native-plugin (node/npm unavailable)")
    print("OK dsh-bootstrap-demo" if dsh_bootstrap_ran else "SKIP dsh-bootstrap-demo (node/npm unavailable)")
    print("OK legacy-link-install-demo")
    print("OK update-override-demo")
    print("OK self-update-demo")
    print("OK token-loader-demo")
    print("OK scaffold-gitattributes-demo")
    print("OK ensure-frontmatter-demo")
    print("OK extract-release-notes-demo")
    if args.refresh_examples:
        print(f"REFRESHED {EXAMPLES_ROOT}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
